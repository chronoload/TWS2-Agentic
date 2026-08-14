"""
TS2 FastAPI 服务端应用
参考思源笔记 API 设计：POST + JSON body，统一返回格式 {code, msg, data}
支持：文件 CRUD、目录浏览、搜索、WebSocket 实时推送、文件同步分发
      文件上传/下载（手机端传输）、局域网访问
"""

import asyncio
import json
import logging
import os
import re
import sys
import shutil
import socket
import sqlite3
import subprocess
import threading
import time
import concurrent.futures
import datetime
import uuid
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

# 高性能 JSON 序列化（orjson 优先）
try:
    import orjson as _orjson
    def _json_dumps(obj, **kwargs):
        return _orjson.dumps(obj, option=_orjson.OPT_INDENT_2).decode('utf-8')
    def _json_loads(s, **kwargs):
        return _orjson.loads(s)
    def _json_dumps_compact(obj):
        """紧凑 JSON（无缩进），用于 SSE 等需要单行输出的场景"""
        return _orjson.dumps(obj).decode('utf-8')
    logger_json = "orjson"
except ImportError:
    _json_dumps = json.dumps
    _json_loads = json.loads
    def _json_dumps_compact(obj):
        return json.dumps(obj, separators=(',', ':'), ensure_ascii=False)
    logger_json = "json"

from fastapi import FastAPI, HTTPException, WebSocket, WebSocketDisconnect, Request, UploadFile, File, Form
from fastapi.responses import JSONResponse, HTMLResponse, FileResponse, StreamingResponse, Response
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from starlette.middleware.gzip import GZipMiddleware

from .ws import WebSocketManager, get_ws_manager
from .sync import FileSyncEngine, SyncEvent, SyncStatus, SyncChange
from .network import NetworkSettings, configure_firewall, check_network_access, set_network_profile_private, get_all_network_interfaces
from .tunnel import TunnelManager, get_tunnel_manager

logger = logging.getLogger(__name__)

# ─── 请求/响应模型 ───────────────────────────────────────────

class APIResponse(BaseModel):
    """统一 API 响应格式（参考思源笔记）"""
    code: int = 0
    msg: str = ""
    data: Any = None


class FileReadRequest(BaseModel):
    path: str


class FileWriteRequest(BaseModel):
    path: str
    content: str


class FileRemoveRequest(BaseModel):
    path: str


class FileRenameRequest(BaseModel):
    old_path: str
    new_path: str


class DirReadRequest(BaseModel):
    path: str = ""


class DirCreateRequest(BaseModel):
    path: str


class SearchRequest(BaseModel):
    query: str
    subdir: str = ""
    sort_by: str = "name"     # name | size | mtime | type
    order: str = "asc"        # asc | desc
    type_filter: str = ""     # "" | dir | file | .ext


class ScanRequest(BaseModel):
    subdir: str = ""


class SyncRequest(BaseModel):
    """同步请求（参考 siyuan-android mobileSwitch 参数）"""
    direction: str = "both"  # upload, download, both
    paths: list = []  # 指定路径，空=全部
    mobile_switch: bool = False  # 是否来自移动端前后台切换（参考 siyuan-android）


class WsCommand(BaseModel):
    """WebSocket 命令"""
    cmd: str
    req_id: float = 0
    param: Dict[str, Any] = {}


class TaskUpdateRequest(BaseModel):
    """任务更新请求"""
    id: str
    title: Optional[str] = None
    description: Optional[str] = None
    due_date: Optional[str] = None
    priority: Optional[str] = None
    status: Optional[str] = None
    start_time: Optional[str] = None
    duration: Optional[int] = None
    recurrence: Optional[str] = None


class TaskCreateRequest(BaseModel):
    """任务创建请求"""
    title: str
    description: str = ""
    due_date: str = ""
    priority: str = "中"
    status: str = "待办"
    start_time: str = ""
    duration: int = 60
    recurrence: str = "不循环"


class TaskDeleteRequest(BaseModel):
    """任务删除请求"""
    id: str


class BookmarkAddRequest(BaseModel):
    """添加书签请求"""
    name: str
    url: str
    category: str = "其他"
    icon: str = "🔖"


class BookmarkDeleteRequest(BaseModel):
    """删除书签请求"""
    id: str


class CourseProgressRequest(BaseModel):
    """课程进度请求"""
    course_id: str


class LessonStatusRequest(BaseModel):
    """课时状态请求（status: 'completed' / 'not_started' 等字符串状态）"""
    course_id: str
    lesson_number: int
    status: str


class LessonReviewRequest(BaseModel):
    """复习调度请求（workload: 工作量/难度，参考 course_tracker.py 的 update_review_schedule）
    用于复习间隔计算：<=5→7天, <=15→4天, <=30→2天, >30→1天"""
    course_id: str
    lesson_number: int
    workload: int = 5


class TimetableSlotCreateRequest(BaseModel):
    """课程表课时创建请求"""
    timetable_id: str = ""
    course_name: str
    day_of_week: int  # 1-7
    start_time: str  # "08:00"
    end_time: str  # "08:45"
    location: str = ""
    teacher: str = ""
    period_idx: int = 0
    color: str = ""


class TimetableSlotDeleteRequest(BaseModel):
    """课程表课时删除请求"""
    timetable_id: str = ""
    slot_id: str


class TimetableCreateRequest(BaseModel):
    """课程表创建请求"""
    name: str
    semester_start: str = ""
    semester_end: str = ""


class TimetableSetActiveRequest(BaseModel):
    """设置激活课程表请求"""
    timetable_id: str


class TimetableDeleteRequest(BaseModel):
    """课程表删除请求"""
    timetable_id: str


class NoteAnalyzeRequest(BaseModel):
    """笔记分析请求"""
    path: str
    course_id: str = ""
    lesson_num: int = 0


class FillBlankRequest(BaseModel):
    """填空题生成请求"""
    content: str


class SpacedReviewRequest(BaseModel):
    """间隔复习评分请求"""
    elem_id: str
    quality: int


class SpacedLoadRequest(BaseModel):
    """间隔复习加载请求"""
    path: str
    course_id: str = ""
    lesson_num: int = 0


class SpacedAnnotateRequest(BaseModel):
    """间隔复习批注请求"""
    elem_id: str
    text: str


class NotePairsRequest(BaseModel):
    """获取笔记中定理-证明/问题-解答配对"""
    path: str
    course_id: str = ""
    lesson_num: int = 0
    pair_type: str = "theorem_proof"  # theorem_proof | problem_solution


class NotePairSaveRequest(BaseModel):
    """重读写入定理/证明或问题/解答到笔记"""
    path: str
    pair_type: str  # theorem_proof | problem_solution
    # 元素在配对列表中的索引
    index: int
    # 用户编辑后的定理/问题内容（raw）
    primary_content: str
    # 用户编辑后的证明/解答内容（raw）
    secondary_content: str
    # 原始配对信息（含 raw_start/raw_end 用于定位替换）
    primary_raw_start: int
    primary_raw_end: int
    secondary_raw_start: int = -1
    secondary_raw_end: int = -1
    primary_elem_type: str = ""  # 定理类型（theorem/corollary/lemma/...）或 "problem"


class UpdateLessonRequest(BaseModel):
    """更新课时信息请求"""
    course_id: str
    lesson_number: int
    updates: Dict[str, Any]


class UpdateCourseRequest(BaseModel):
    """更新课程信息请求"""
    course_id: str
    updates: Dict[str, Any]


class AgentChatRequest(BaseModel):
    """Agent 聊天请求"""
    message: str
    context: Optional[Dict[str, Any]] = None
    session_id: str = ""  # 前端持久化的会话 ID（参考 Crush Session.ID）
    attachments: Optional[List[Dict[str, str]]] = None  # 多模态附件 [{"kind":"image","data_url":"data:...","path":"..."}]


class AgentSessionSwitchRequest(BaseModel):
    """Agent 会话切换/删除请求"""
    session_id: str


class AgentInjectSkillRequest(BaseModel):
    """注入技能请求：把选中的 SKILL.md 全文作为指令注入当前会话"""
    skill_name: str
    session_id: str = ""
    as_system: bool = False  # True=注入为 system 指令；False=作为 user 消息呈现
    direct_text: str = ""    # 非技能类型（tool/mcp/workflow/plugin）直接注入的指令文本


# ─── 工具函数 ────────────────────────────────────────────────

def get_local_ip() -> str:
    """获取本机局域网 IP"""
    try:
        s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
        s.connect(("8.8.8.8", 80))
        ip = s.getsockname()[0]
        s.close()
        return ip
    except Exception:
        return "127.0.0.1"


def find_available_port(start_port: int = 6906, max_tries: int = 100,
                        host: str = "0.0.0.0") -> int:
    """从 start_port 开始查找可用端口，跳过已被占用的端口

    Args:
        start_port: 起始端口号
        max_tries: 最大尝试次数
        host: 检测绑定的地址

    Returns:
        可用的端口号

    Raises:
        OSError: 在 max_tries 范围内找不到可用端口
    """
    for port in range(start_port, start_port + max_tries):
        try:
            with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as sock:
                sock.setsockopt(socket.SOL_SOCKET, socket.SO_REUSEADDR, 1)
                sock.bind((host, port))
                # bind 成功 = 端口可用
                return port
        except OSError:
            continue
    raise OSError(f"在 {start_port}-{start_port + max_tries - 1} 范围内找不到可用端口")


# ─── API 配置加载 ──────────────────────────────────────────────

CONFIG_DIR = Path.home() / ".ts2"
CONFIG_FILE = CONFIG_DIR / "ts2_api_config.json"

def load_api_config() -> dict:
    """读取 ts2_api_config.json，不存在则返回默认"""
    if CONFIG_FILE.exists():
        return json.loads(CONFIG_FILE.read_text(encoding="utf-8"))
    return {
        "api_token": "",
        "source_auth_code": "",
        "workspaces": [
            {
                "name": "默认工作区",
                "path": str(Path.cwd().resolve()),
                "auth_code": "",
                "readable": True,
                "writable": True,
                "relaxed": True,
            }
        ],
    }

def check_path_access(file_path: str, mode: str = "read") -> bool:
    """检查文件路径是否在允许的工作区目录内"""
    config = load_api_config()
    real_path = Path(file_path).resolve()
    for ws in config.get("workspaces", []):
        ws_root = Path(ws["path"]).resolve()
        if not (str(real_path).startswith(str(ws_root) + os.sep) or real_path == ws_root):
            continue
        if ws.get("relaxed"):
            return True
        if mode == "read" and ws.get("readable"):
            return True
        if mode == "write" and ws.get("writable"):
            return True
    return False


def _is_resource_path(file_path: str, workspace_dir: str = "") -> bool:
    """检查路径是否在 resource_index.json 中注册（课程资源）"""
    if not workspace_dir:
        workspace_dir = str(Path.cwd())
    real_path = Path(file_path).resolve()
    for p in [Path(workspace_dir) / "data" / "resource_index.json", Path(workspace_dir) / "resource_index.json"]:
        if p.exists():
            try:
                data = json.loads(p.read_text(encoding="utf-8"))
                for resources in data.values():
                    for r in resources:
                        rp = r.get("path", "")
                        if rp and Path(rp).resolve() == real_path:
                            return True
            except Exception:
                pass
            break
    return False


# ─── 鉴权辅助 ──────────────────────────────────────────────────

# ─── 服务端 Session ──────────────────────────────────────────────

# {session_token: {"code": str, "created_at": float}}
_sessions: Dict[str, Dict] = {}
_SESSION_TTL = 86400 * 7  # 7 天过期
_RATE_LIMIT = 10          # 10 次
_RATE_WINDOW = 300        # 5 分钟内
_attempts: Dict[str, list] = {}

def _create_session(code: str) -> str:
    token = uuid.uuid4().hex
    _sessions[token] = {"code": code, "created_at": time.time(), "authed": []}
    return token

def _get_session_code(token: str) -> str:
    sess = _sessions.get(token)
    if not sess:
        return ""
    if time.time() - sess["created_at"] > _SESSION_TTL:
        del _sessions[token]
        return ""
    return sess["code"]

def _session_add_workspace(token: str, resolved_path: str):
    sess = _sessions.get(token)
    if sess and resolved_path not in sess["authed"]:
        sess["authed"].append(resolved_path)

def _session_has_workspace(token: str, resolved_path: str) -> bool:
    sess = _sessions.get(token)
    return bool(sess and resolved_path in sess["authed"])

def _delete_session(token: str):
    _sessions.pop(token, None)

def _check_rate_limit(key: str) -> bool:
    """每个 IP 5 分钟最多 _RATE_LIMIT 次登录尝试"""
    now = time.time()
    lst = _attempts.setdefault(key, [])
    lst[:] = [t for t in lst if now - t < _RATE_WINDOW]
    if len(lst) >= _RATE_LIMIT:
        return False
    lst.append(now)
    return True

PUBLIC_ENDPOINTS = {
    "/api/system/health",
    "/api/system/authInfo",
    "/api/system/loginAuth",
    "/api/system/sourceAuth",
    "/api/system/stats",
    "/api/system/version",
}

def _get_all_auth_codes(config: dict) -> set:
    """从所有工作区中收集授权码（兼容旧版全局 access_auth_code）"""
    codes = set()
    old_global = config.get("access_auth_code", "")
    if old_global:
        codes.add(old_global)
    for ws in config.get("workspaces", []):
        code = ws.get("auth_code", "")
        if code:
            codes.add(code)
    return codes

async def check_auth(request: Request, config: dict) -> bool:
    client_host = request.client.host if request.client else "127.0.0.1"
    # localhost 免检
    if client_host in ("127.0.0.1", "::1", "localhost"):
        return True

    api_token = config.get("api_token", "")
    auth_codes = _get_all_auth_codes(config)

    # 如果未配置任何凭证，则放行
    if not api_token and not auth_codes:
        return True

    # --- 通用 Token 检查 ---
    token_ok = False
    auth_header = request.headers.get("authorization", "")
    if auth_header.startswith(("Token ", "Bearer ")):
        token_val = auth_header.split(" ", 1)[1]
        if token_val == api_token and api_token:
            token_ok = True
    # 兼容 query param
    if not token_ok and request.query_params.get("token") == api_token and api_token:
        token_ok = True
    # 兼容 session cookie（如果 cookie 中存的是 token）
    if not token_ok:
        session_token = request.cookies.get("ts2_session", "")
        if session_token:
            sess_code = _get_session_code(session_token)
            if sess_code and sess_code == api_token:
                token_ok = True

    # --- Code 检查（仅用于登录端点） ---
    code_ok = False
    if auth_codes:
        # 从 cookie 或 header 或 Basic 中提取 code
        session_token = request.cookies.get("ts2_session", "")
        if session_token:
            sess_code = _get_session_code(session_token)
            if sess_code and sess_code in auth_codes:
                code_ok = True
        if not code_ok and request.headers.get("x-auth-code", "") in auth_codes:
            code_ok = True
        if not code_ok and auth_header.startswith("Basic "):
            import base64
            try:
                decoded = base64.b64decode(auth_header.split(" ", 1)[1]).decode("utf-8")
                if ":" in decoded and decoded.split(":", 1)[1] in auth_codes:
                    code_ok = True
            except Exception:
                pass

    # --- 根据路径决定鉴权策略 ---
    path = request.url.path

    # 登录端点：必须同时满足（与逻辑）
    if path == "/api/system/loginAuth":
        need_token = bool(api_token)
        need_code = bool(auth_codes)
        if need_token and need_code:
            return token_ok and code_ok
        if need_token:
            return token_ok
        if need_code:
            return code_ok
        return True

    # 普通端点：仅验证 token（如果配置了 token）；否则退化为验证 code
    if api_token:
        return token_ok
    if auth_codes:
        # 为了兼容旧逻辑，如果未配置 token，则使用 code
        return code_ok
    return True

def _check_ws_auth(websocket: WebSocket, config: dict) -> bool:
    """WebSocket 鉴权：同 check_auth，但基于 WebSocket 的 query_params 和 cookies"""
    client_host = websocket.client.host if websocket.client else "127.0.0.1"
    if client_host in ("127.0.0.1", "::1", "localhost"):
        return True

    api_token = config.get("api_token", "")
    auth_codes = _get_all_auth_codes(config)

    if not api_token and not auth_codes:
        return True

    token_ok = False
    if api_token:
        qp_token = websocket.query_params.get("token", "")
        if qp_token == api_token:
            token_ok = True
        if not token_ok:
            sess_token = websocket.cookies.get("ts2_session", "")
            if sess_token:
                sess_code = _get_session_code(sess_token)
                if sess_code == api_token:
                    token_ok = True

    code_ok = False
    if auth_codes:
        sess_token = websocket.cookies.get("ts2_session", "")
        if sess_token:
            sess_code = _get_session_code(sess_token)
            if sess_code and sess_code in auth_codes:
                code_ok = True
        if not code_ok:
            from urllib.parse import parse_qs
            raw_qs = websocket.url.query if hasattr(websocket.url, 'query') else ""
            qs_params = parse_qs(raw_qs) if raw_qs else {}
            qp_code = qs_params.get("auth_code", [None])[0]
            if qp_code and qp_code in auth_codes:
                code_ok = True

    need_token = bool(api_token)
    need_code = bool(auth_codes)
    if need_token and need_code:
        return token_ok and code_ok
    if need_token:
        return token_ok
    if need_code:
        return code_ok
    return True


def _check_source_auth(request: Request, require_write: bool = False) -> bool:
    """源码浏览器独立鉴权：读取需 source_auth_code 或工作区 auth_code，写入必须 source_auth_code"""
    config = load_api_config()
    source_code = config.get("source_auth_code", "")

    # 无 source_auth_code 时退化为普通鉴权
    if not source_code:
        return True

    # 本地免检读取，写入仍需验证
    client_host = request.client.host if request.client else "127.0.0.1"
    if not require_write and client_host in ("127.0.0.1", "::1", "localhost"):
        return True

    # Session 中已标记 source_authed
    session_token = request.cookies.get("ts2_session", "")
    if session_token:
        sess = _sessions.get(session_token)
        if sess and sess.get("source_authed"):
            return True

    # 从请求头取 source_code
    src_code = request.headers.get("X-Source-Auth", "")
    if not src_code:
        src_code = request.query_params.get("source_code", "")

    if src_code == source_code:
        return True

    return False

# ─── 应用创建 ────────────────────────────────────────────────

def _get_model_selector():
    from ..model_selector import get_model_selector
    return get_model_selector()

def create_app(workspace_dir: Optional[str] = None, host: str = "0.0.0.0",
               port: int = 6906) -> FastAPI:
    """创建 FastAPI 应用实例"""

    if workspace_dir is None:
        workspace_dir = os.getcwd()

    app = FastAPI(
        title="TS2 Server",
        description="TS2 本地文件同步分发服务（参考思源笔记架构）",
        version="1.0.0",
    )

    # ─── 自定义 CORS 中间件 ──────────────────────────────
    # 反射 Origin + credentials=True，methods/headers 用显式值（w3c 规范要求）
    @app.middleware("http")
    async def cors_middleware(request: Request, call_next):
        # 获取 Origin
        origin = request.headers.get("origin", "")
        # 判断是否为 file:// 或 null 来源
        is_null_or_file = origin in ("null", "file://", "file:", "")
        effective_origin = "*" if is_null_or_file else origin

        # 无 Origin 直接放行（非浏览器请求）
        if not origin:
            return await call_next(request)

        methods = "GET, POST, PUT, DELETE, PATCH, OPTIONS"
        headers = "Content-Type, Authorization, X-Auth-Code, X-Requested-With, Range, If-None-Match"

        # OPTIONS 预检请求
        if request.method == "OPTIONS":
            resp = JSONResponse(content="")
            resp.headers["Access-Control-Allow-Origin"] = effective_origin
            if not is_null_or_file:
                resp.headers["Access-Control-Allow-Credentials"] = "true"
            resp.headers["Access-Control-Allow-Methods"] = methods
            resp.headers["Access-Control-Allow-Headers"] = headers
            resp.headers["Access-Control-Max-Age"] = "86400"
            return resp

        # 实际业务请求
        response = await call_next(request)
        response.headers["Access-Control-Allow-Origin"] = effective_origin
        if not is_null_or_file:
            response.headers["Access-Control-Allow-Credentials"] = "true"
        response.headers["Access-Control-Allow-Methods"] = methods
        response.headers["Access-Control-Allow-Headers"] = headers
        return response
    # GZip 压缩中间件
    app.add_middleware(GZipMiddleware, minimum_size=1000)

    # 初始化引擎
    sync_engine = FileSyncEngine(workspace_dir)
    ws_manager = get_ws_manager()
    network_settings = NetworkSettings()
    network_settings.port = port

    # 文件变更回调 → WebSocket 推送
    async def on_file_change(event: SyncEvent):
        await ws_manager.notify_path_change(
            path=event.path,
            change_type=event.event_type,
            data=event.to_dict(),
        )
        # 同时推送文件树刷新
        if event.event_type in ("created", "deleted", "renamed"):
            await ws_manager.push_reload_filetree()

    sync_engine.set_change_callback(on_file_change)

    # 检查点事件回调 → WebSocket 推送（参考 Crush PubSub）
    def _on_checkpoint_event(event_type: str, payload: dict):
        """中间件事件 → WS 广播（同步→异步桥接）"""
        try:
            loop = asyncio.get_event_loop()
            if loop.is_running():
                loop.call_soon_threadsafe(
                    lambda: asyncio.ensure_future(
                        ws_manager.broadcast(event_type, 0, "", payload)
                    )
                )
            else:
                # 事件循环未运行时（罕见），直接忽略
                pass
        except Exception:
            pass

    # 延迟导入，避免循环依赖
    try:
        from ..middleware.shadow_checkpoint import CheckpointMiddleware
        CheckpointMiddleware.on_event(_on_checkpoint_event)
    except Exception:
        pass

    # ─── Workflow 引擎事件 → WS 广播桥 ──────────────────────
    # 在模块加载时先用 workspace 路径固定进程内 WorkflowEngine 单例，
    # 后续 config_ui 入口 / WorkflowTool 工具 / /api/workflow/* 拿到的是同一实例，
    # agent 创建处的 set_agent 注入也随之全局生效。
    try:
        from ..workflow_engine import get_workflow_engine as _get_wf_engine

        def _on_workflow_event(event):
            """WorkflowEngine 事件（工作流线程回调）→ WS 广播（同步→异步桥接）"""
            try:
                loop = asyncio.get_event_loop()
                if loop.is_running():
                    loop.call_soon_threadsafe(
                        lambda: asyncio.ensure_future(
                            ws_manager.broadcast("workflow_" + event.type.value, 0, "", {
                                "instance_id": event.instance_id,
                                "event": event.type.value,
                                **event.data,
                            })
                        )
                    )
            except Exception:
                pass

        _wf_db_path = os.path.join(str(workspace_dir), "data", "workflow.db")
        _wf_engine = _get_wf_engine(_wf_db_path)
        _wf_engine.on_event(_on_workflow_event)
        app.state.workflow_engine = _wf_engine
        logger.info(f"WorkflowEngine 事件桥已注册: {_wf_db_path}")
    except Exception as e:
        logger.warning(f"WorkflowEngine 事件桥注册失败: {e}")

    # 将引擎挂到 app.state
    app.state.sync_engine = sync_engine
    app.state.ws_manager = ws_manager
    app.state.network_settings = network_settings
    app.state.host = host
    app.state.port = port
    app.state.workspace_dir = workspace_dir
    app.state.start_time = time.time()
    app.state.local_ip = get_local_ip()
    app.state.tunnel_manager = get_tunnel_manager()
    app.state.jupyter_process = None
    app.state.jupyter_port = 8888
    app.state.jupyter_url = None
    app.state._exec_proc = None

    # ─── 鉴权中间件（OPTIONS 已由 cors_middleware 在之前处理）──────
    @app.middleware("http")
    async def auth_middleware(request: Request, call_next):
        if request.method == "OPTIONS":
            return await call_next(request)
        path = request.url.path
        # 静态文件和非 API 路由直接放行（让前端页面能加载）
        if not path.startswith("/api/"):
            return await call_next(request)
        # 公开 API 端点和 tunnel 放行
        if path in PUBLIC_ENDPOINTS or path.startswith("/api/tunnel/"):
            return await call_next(request)
        config = load_api_config()
        if not await check_auth(request, config):
            logger.warning(f"AUTH FAIL {request.method} {path} client={request.client.host if request.client else '?'} auth={request.headers.get('authorization','(none)')[:80]}")
            return JSONResponse(
                status_code=401,
                content={"code": 401, "msg": "未授权访问，请提供有效的 Token 或授权码", "data": None},
            )
        return await call_next(request)

    # ─── 公开端点 ────────────────────────────────────────────

    @app.get("/api/system/authInfo")
    async def auth_info(request: Request):
        """返回鉴权信息（客户端判断是否需要登录 + 工作区授权状态）"""
        client_host = request.client.host if request.client else "127.0.0.1"
        local = client_host in ("127.0.0.1", "::1", "localhost")
        config = load_api_config()
        auth_codes = _get_all_auth_codes(config)
        workspaces = config.get("workspaces", [])
        session_token = request.cookies.get("ts2_session", "")
        session_code = _get_session_code(session_token) if session_token else ""
        api_token = config.get("api_token", "")
        ws_access = []
        for ws in workspaces:
            ws_code = ws.get("auth_code", "")
            ws_resolved = str(Path(ws["path"]).resolve()) if ws.get("path") else ""
            accessible = local
            if not accessible and session_code:
                accessible = (session_code == ws_code) or (session_code == api_token)
            if not accessible and session_token and ws_resolved:
                accessible = _session_has_workspace(session_token, ws_resolved)
            ws_access.append({
                "name": ws.get("name", ""),
                "path": ws.get("path", ""),
                "accessible": accessible,
            })
        return ok({
            "needAuth": not local and bool(auth_codes or config.get("api_token")),
            "local": local,
            "hasAuthCode": bool(auth_codes),
            "hasToken": bool(config.get("api_token")),
            "hasSourceAuth": bool(config.get("source_auth_code")),
            "workspaceAccess": ws_access,
        })

    @app.post("/api/system/loginAuth")
    async def login_auth(request: Request):
        """远程登录：验证授权码和/或 token（与逻辑），创建服务端 session"""
        body = await request.json()
        code = body.get("code", "")
        token = body.get("token", "")
        config = load_api_config()
        api_token = config.get("api_token", "")
        auth_codes = _get_all_auth_codes(config)
        need_token = bool(api_token)
        need_code = bool(auth_codes)

        # 频率限制（本地免检）
        client_host = request.client.host if request.client else "unknown"
        if client_host not in ("127.0.0.1", "::1", "localhost"):
            if not _check_rate_limit(f"login:{client_host}"):
                return err(429, "登录尝试过于频繁，请稍后再试")

        # 与逻辑验证
        token_ok = (not need_token) or (token and token == api_token)
        code_ok = (not need_code) or (code and code in auth_codes)

        if not (token_ok and code_ok):
            if need_token and need_code:
                return err(401, "需要同时提供正确的 Token 和授权码")
            if need_token:
                return err(401, "Token 错误")
            return err(401, "授权码错误")

        # 验证通过，创建 session
        session_code = token if (token and token == api_token) else code
        session_token = _create_session(session_code)
        # 标记已认证的工作区
        for ws in config.get("workspaces", []):
            if ws.get("path"):
                if not need_code or ws.get("auth_code") == code:
                    _session_add_workspace(session_token, str(Path(ws["path"]).resolve()))
        # 如果用 token 登录，放行所有工作区
        if token and token == api_token:
            for ws in config.get("workspaces", []):
                if ws.get("path"):
                    _session_add_workspace(session_token, str(Path(ws["path"]).resolve()))

        response = JSONResponse(content=ok(data={"needToken": need_token, "needCode": need_code}))
        response.set_cookie(key="ts2_session", value=session_token, httponly=True, samesite="lax")
        return response

    @app.post("/api/system/sourceAuth")
    async def source_auth(request: Request):
        """验证源码授权码，成功后在 session 中标记"""
        body = await request.json()
        code = body.get("code", "")
        config = load_api_config()
        source_code = config.get("source_auth_code", "")
        if not source_code:
            return ok(data={"sourceAuthed": True})
        if code == source_code:
            session_token = request.cookies.get("ts2_session", "")
            if session_token:
                sess = _sessions.get(session_token)
                if sess:
                    sess["source_authed"] = True
            else:
                # 无 session 时创建一个
                session_token = _create_session(code)
                _sessions[session_token]["source_authed"] = True
            response = JSONResponse(content=ok(data={"sourceAuthed": True}))
            if not request.cookies.get("ts2_session"):
                response.set_cookie(key="ts2_session", value=session_token, httponly=True, samesite="lax")
            return response
        return err(403, "源码授权码错误")

    @app.post("/api/system/logoutAuth")
    async def logout_auth(request: Request):
        """登出：删除服务端 session"""
        session_token = request.cookies.get("ts2_session", "")
        if session_token:
            _delete_session(session_token)
        response = JSONResponse(content=ok())
        response.delete_cookie("ts2_session")
        return response

    @app.post("/api/system/reloadConfig")
    async def reload_config():
        """重新加载配置（无需重启服务）"""
        config = load_api_config()
        return ok({"token": config.get("api_token", "")}, "配置已重新加载")

    @app.get("/api/system/workspaces")
    async def system_workspaces():
        """返回可用工作区列表（不含敏感字段），标注当前激活工作区与默认工作区"""
        config = load_api_config()
        workspaces = config.get("workspaces", [])
        current = os.path.normcase(str(Path(app.state.workspace_dir).resolve()))
        default = os.path.normcase(str(Path(workspace_dir).resolve()))
        safe = []
        for ws in workspaces:
            item = {k: v for k, v in ws.items() if k != "auth_code"}
            ws_path = os.path.normcase(str(Path(ws["path"]).resolve()))
            item["active"] = ws_path == current
            item["isDefault"] = ws_path == default
            safe.append(item)
        return ok(data=safe)

    @app.post("/api/system/switchWorkspace")
    async def switch_workspace(req: Request):
        """切换当前工作区（热切换 sync_engine.root），同时校验工作区授权"""
        body = await req.json()
        ws_path = body.get("path", "")
        ws_code = body.get("code", "")
        if not ws_path or not Path(ws_path).exists():
            return err(400, f"工作区路径不存在: {ws_path}")
        if not check_path_access(ws_path, "read"):
            return err(403, "无权访问该工作区")

        config = load_api_config()
        api_token = config.get("api_token", "")

        # Authorization header token（API 客户端）
        auth_header = req.headers.get("authorization", "")
        token_valid = False
        if auth_header.startswith(("Token ", "Bearer ")):
            header_token = auth_header.split(" ", 1)[1]
            if header_token == api_token and api_token:
                token_valid = True

        # 查找目标工作区
        target_ws = None
        for ws in config.get("workspaces", []):
            if Path(ws["path"]).resolve() == Path(ws_path).resolve():
                target_ws = ws
                break

        # 鉴权：token 放行 / 无 auth_code 放行 / session 已认证过 / 传入正确 code
        session_token = req.cookies.get("ts2_session", "")
        resolved = str(Path(ws_path).resolve())
        known_ws = _session_has_workspace(session_token, resolved) if session_token else False

        if target_ws:
            need_check = bool(target_ws.get("auth_code", ""))
            if need_check and not token_valid and not known_ws:
                if ws_code != target_ws["auth_code"]:
                    return err(403, f"授权码错误，无权访问工作区「{target_ws.get('name', '')}」")
                if session_token:
                    _session_add_workspace(session_token, resolved)
        elif not token_valid and not known_ws:
            return err(403, "无权访问该工作区")

        app.state.workspace_dir = ws_path
        engine = FileSyncEngine(ws_path)
        # 默认工作区（create_app 时的 workspace_dir）保留 EXPOSED_DIRS 安全限制；
        # 切换到其他工作区时暴露全部文件（不受 EXPOSED_DIRS 限制）。
        if os.path.normcase(str(Path(ws_path).resolve())) != os.path.normcase(str(Path(workspace_dir).resolve())):
            engine.EXPOSED_DIRS = set()
            engine.EXPOSED_ROOT_FILES = set()
        app.state.sync_engine = engine
        app.state.sync_engine.set_change_callback(on_file_change)
        logger.info(f"Switched workspace to {ws_path}")

        # 更新 session cookie
        resp = ok(msg=f"已切换到: {ws_path}")
        response = JSONResponse(content=resp)
        if ws_code and not session_token:
            session_token = _create_session(ws_code)
            _session_add_workspace(session_token, resolved)
            response.set_cookie(key="ts2_session", value=session_token, httponly=True, samesite="lax")
        return response

    # ─── 工具函数 ────────────────────────────────────────────

    def ok(data: Any = None, msg: str = "") -> dict:
        return {"code": 0, "msg": msg, "data": data}

    def err(code: int = -1, msg: str = "") -> dict:
        return {"code": code, "msg": msg, "data": None}

    # ─── note_analyzer / 间隔重复管理器初始化 ──────────────────
    _note_analyzer_ready = False
    _srm = None
    try:
        import sys as _sys
        _project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..'))
        if _project_root not in _sys.path:
            _sys.path.insert(0, _project_root)
        from note_analyzer import NoteAnalyzer, NoteElement, SpacedRepetitionManager, KnowledgeGraph, KnowledgeCategory
        from pathlib import Path as _Path
        _srm_data_path = _Path(os.path.join(workspace_dir, 'data', 'spaced_review_data.json'))
        _srm = SpacedRepetitionManager(_srm_data_path)
        _note_analyzer_ready = True
        logger.info("note_analyzer loaded, SRM data: %s", _srm_data_path)
    except Exception as _e:
        logger.warning("note_analyzer import failed: %s", _e)

    # ─── 数据枢纽初始化（与 course_tracker 共用 data_hub/data_hub.db）───
    _hub_ready = False
    _hub = None
    try:
        from ws2_data_hub import init_data_hub, get_data_hub
        _hub = get_data_hub()
        if _hub is None:
            _hub = init_data_hub(Path(workspace_dir))
        _hub_ready = _hub is not None
        logger.info("data_hub initialized at %s", _hub.db_path if _hub else "N/A")
    except Exception as _he:
        _hub = None
        _hub_ready = False
        logger.warning("data_hub init failed: %s", _he)

    # ─── Agent 核心能力已内建于 Agent 类（见 agent.py __init__ 的 attach_agent_core）
    def _get_agent_core():
        return None

    def _get_hub():
        if not _hub_ready:
            return None
        return _hub

    async def _hub_body(req: Request) -> dict:
        """安全读取 JSON body（空 body 返回 {}）"""
        try:
            return await req.json()
        except Exception:
            return {}

    def _resolve_ws_path(rel_path: str) -> str:
        """将相对路径解析为工作区绝对路径"""
        if os.path.isabs(rel_path):
            return rel_path
        return os.path.join(workspace_dir, rel_path)

    # ─── Vue SPA 挂载（必须在 catch-all 路由之前）─────────────
    static_dir = Path(__file__).parent / "static"
    web_dist = Path(__file__).parent / "web" / "dist"
    if web_dist.exists():
        # 挂载 /app 的静态资源（js/css/images 等）
        app.mount("/app/assets", StaticFiles(directory=str(web_dist / "assets")), name="app-assets")

        # SPA fallback: /app, /app/, /app/xxx 路径返回 index.html
        @app.get("/app", response_class=HTMLResponse)
        @app.get("/app/", response_class=HTMLResponse)
        @app.get("/app/{spa_path:path}", response_class=HTMLResponse)
        async def app_spa(spa_path: str = ""):
            """Vue SPA fallback — 所有 /app 路径返回 index.html"""
            # 如果请求的是实际存在的文件（有扩展名且在 dist 中），直接返回文件
            if spa_path and "." in spa_path.split("/")[-1]:
                file_path = web_dist / spa_path
                if file_path.exists() and file_path.is_file():
                    return FileResponse(file_path)
            index_file = web_dist / "index.html"
            if index_file.exists():
                return HTMLResponse(content=index_file.read_text(encoding="utf-8"))
            return HTMLResponse(content="<h1>Vue SPA not found</h1>", status_code=404)

    # ─── 首页 ─────────────────────────────────────────────────

    @app.get("/", response_class=HTMLResponse)
    async def index():
        index_file = static_dir / "index.html"
        if index_file.exists():
            return HTMLResponse(content=index_file.read_text(encoding="utf-8"))
        return HTMLResponse(content="<h1>TS2 Server</h1><p>Frontend not found</p>")

    # ─── 系统 API ────────────────────────────────────────────

    @app.post("/api/system/version")
    @app.get("/api/system/version")
    async def system_version():
        return ok(data={
            "version": "1.0.0",
            "runtime": "python",
            "port": app.state.port,
            "local_ip": app.state.local_ip,
            "network": app.state.network_settings.to_dict(),
        })

    # ─── 网络设置 API（参考思源笔记网络配置）──────────────

    @app.post("/api/system/getNetworkSettings")
    async def get_network_settings():
        """获取网络设置"""
        ns: NetworkSettings = app.state.network_settings
        access_info = check_network_access(ns.port)
        return ok(data={
            "settings": ns.to_dict(),
            "access": access_info,
        })

    @app.post("/api/system/setNetworkSettings")
    async def set_network_settings(req: Request):
        """更新网络设置"""
        ns: NetworkSettings = app.state.network_settings
        body = await req.json()

        if "allow_lan" in body:
            ns.allow_lan = bool(body["allow_lan"])
        if "allow_public_network" in body:
            ns.allow_public_network = bool(body["allow_public_network"])
        if "allow_usb" in body:
            ns.allow_usb = bool(body["allow_usb"])

        ns.save()
        return ok(data=ns.to_dict())

    @app.post("/api/system/configureFirewall")
    async def api_configure_firewall(req: Request):
        """配置防火墙规则（参考思源笔记自动配置防火墙）"""
        body = await req.json()
        allow = body.get("allow", True)
        port = body.get("port", app.state.port)

        success, message = configure_firewall(port, allow)

        if success:
            ns: NetworkSettings = app.state.network_settings
            ns.firewall_configured = allow
            ns.save()

        return ok(data={"success": success, "message": message})

    @app.post("/api/system/setNetworkPrivate")
    async def api_set_network_private():
        """将网络配置文件设为专用（参考思源笔记自动设置网络类别）"""
        success, message = set_network_profile_private()
        return ok(data={"success": success, "message": message})

    @app.post("/api/system/checkNetworkAccess")
    async def api_check_network_access():
        """检查网络访问状态"""
        port = app.state.port
        access_info = check_network_access(port)
        return ok(data=access_info)

    # ─── FRP 隧道 API ───────────────────────────────────────

    @app.get("/api/tunnel/status")
    async def tunnel_status():
        """获取 frp 隧道状态"""
        tm: "TunnelManager" = app.state.tunnel_manager
        return ok(data=tm.get_status())

    @app.post("/api/tunnel/start")
    async def tunnel_start():
        """启动 frp 隧道"""
        tm: "TunnelManager" = app.state.tunnel_manager
        result = tm.start()
        if result.get("success"):
            return ok(data=result)
        return err(msg=result.get("message", "启动失败"))

    @app.post("/api/tunnel/stop")
    async def tunnel_stop():
        """停止 frp 隧道"""
        tm: "TunnelManager" = app.state.tunnel_manager
        result = tm.stop()
        return ok(data=result)

    @app.post("/api/tunnel/restart")
    async def tunnel_restart():
        """重启 frp 隧道"""
        tm: "TunnelManager" = app.state.tunnel_manager
        result = tm.restart()
        if result.get("success"):
            return ok(data=result)
        return err(msg=result.get("message", "重启失败"))

    @app.get("/api/tunnel/settings")
    async def tunnel_settings_get():
        """获取 frp 隧道配置"""
        tm: "TunnelManager" = app.state.tunnel_manager
        return ok(data=tm.get_settings())

    @app.post("/api/tunnel/settings")
    async def tunnel_settings_update(req: Request):
        """更新 frp 隧道配置"""
        tm: "TunnelManager" = app.state.tunnel_manager
        body = await req.json()
        # 敏感字段过滤（不返回完整 token）
        result = tm.update_settings(**body)
        # 隐藏 token 前缀
        if result.get("token"):
            result["token_preview"] = result["token"][:4] + "****"
        return ok(data=result)

    @app.post("/api/tunnel/bootProgress")
    async def system_boot_progress():
        return ok(data={"progress": 100, "details": "TS2 Server is ready"})

    @app.post("/api/system/currentTime")
    async def system_current_time():
        return ok(data={"time": time.time(), "date": time.strftime("%Y-%m-%d %H:%M:%S")})

    @app.get("/api/system/stats")
    async def system_stats():
        engine: FileSyncEngine = app.state.sync_engine
        return ok(data={
            "workspace": str(engine.workspace_dir),
            "file_stats": engine.get_file_stats(),
            "ws_sessions": app.state.ws_manager.get_session_count(),
            "uptime": time.time() - app.state.start_time,
            "local_ip": app.state.local_ip,
        })

    @app.get("/api/system/health")
    async def system_health():
        return {
            "status": "ok",
            "version": "1.0.0",
            "uptime": time.time() - app.state.start_time,
            "workspace": app.state.workspace_dir,
            "features": ["rag", "terminal", "jupyter", "sync"],
        }

    # ─── JupyterLab 管理 API ──────────────────────────────

    @app.post("/api/system/jupyterStart")
    async def jupyter_start():
        """启动 JupyterLab 服务（后台异步，不阻塞主服务）"""
        proc = app.state.jupyter_process
        if proc and proc.poll() is None:
            return ok(data={"starting": False, "url": app.state.jupyter_url})
        port = app.state.jupyter_port
        try:
            def _launch():
                p = subprocess.Popen(
                    [sys.executable, "-m", "jupyter", "lab", "--no-browser",
                     f"--port={port}", "--notebook-dir", app.state.workspace_dir,
                     "--LabApp.token=", "--LabApp.password="],
                    stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
                )
                app.state.jupyter_process = p
            threading.Thread(target=_launch, daemon=True).start()
            app.state.jupyter_url = f"http://localhost:{port}/lab"
            return ok(data={"starting": True, "url": app.state.jupyter_url})
        except FileNotFoundError:
            return err(msg="JupyterLab 未安装，请执行: pip install jupyterlab")
        except Exception as e:
            return err(msg=f"启动 JupyterLab 失败: {e}")

    @app.post("/api/system/jupyterStop")
    async def jupyter_stop():
        """停止 JupyterLab 服务"""
        proc = app.state.jupyter_process
        if proc and proc.poll() is None:
            proc.terminate()
            try:
                proc.wait(timeout=5)
            except subprocess.TimeoutExpired:
                proc.kill()
        app.state.jupyter_process = None
        app.state.jupyter_url = None
        return ok(data={"running": False})

    @app.get("/api/system/jupyterStatus")
    async def jupyter_status():
        """获取 JupyterLab 运行状态"""
        proc = app.state.jupyter_process
        running = proc is not None and proc.poll() is None
        return ok(data={"running": running, "url": app.state.jupyter_url})

    @app.post("/api/system/jupyterWaitReady")
    async def jupyter_wait_ready():
        """轮询等待 JupyterLab HTTP 服务就绪（最多 30s）"""
        port = app.state.jupyter_port
        for i in range(30):
            try:
                reader, writer = await asyncio.wait_for(
                    asyncio.open_connection("127.0.0.1", port), timeout=2)
                writer.write(b"GET /lab HTTP/1.0\r\n\r\n")
                await reader.read(64)
                writer.close()
                break
            except (ConnectionRefusedError, OSError, asyncio.TimeoutError):
                await asyncio.sleep(1)
        else:
            return err(msg="JupyterLab 启动超时")
        app.state.jupyter_url = f"http://localhost:{port}/lab"
        return ok(data={"ready": True, "url": app.state.jupyter_url})

    @app.post("/api/system/openInJupyter")
    async def open_in_jupyter(req: Request):
        """获取在 JupyterLab 中打开文件的 URL"""
        body = await req.json()
        file_path = body.get("path", "")
        rel_path = os.path.relpath(file_path, app.state.workspace_dir) if file_path else ""
        jupyter_base = app.state.jupyter_url or f"http://localhost:{app.state.jupyter_port}/lab"
        url = f"{jupyter_base}/tree/{rel_path}" if rel_path else jupyter_base
        return ok(data={"url": url, "file": rel_path})

    @app.get("/api/system/openWithOptions")
    async def open_with_options():
        """获取可用的 IDE 打开方式列表"""
        options = [
            {"id": "jupyter", "name": "JupyterLab", "icon": "📓", "desc": "Python/数据科学"},
            {"id": "builtin", "name": "内置编辑器", "icon": "📝", "desc": "Vditor/纯文本"},
            {"id": "external", "name": "外部编辑器", "icon": "💻", "desc": "系统默认应用"},
            {"id": "vscode", "name": "VS Code", "icon": "💻", "desc": "vscode://file/{path}"},
            {"id": "cursor", "name": "Cursor", "icon": "🖥️", "desc": "cursor://file/{path}"},
        ]
        return ok(data={"options": options})

    # ─── 文件 API（参考思源笔记 /api/file/*）──────────────

    # ─── 独立管线：文件IO、数据查询、Agent 各自独立线程池 ──────
    _file_executor = concurrent.futures.ThreadPoolExecutor(max_workers=4, thread_name_prefix="file-io")
    _data_executor = concurrent.futures.ThreadPoolExecutor(max_workers=2, thread_name_prefix="data")
    _agent_executor = concurrent.futures.ThreadPoolExecutor(max_workers=4, thread_name_prefix="agent")
    _push_executor = concurrent.futures.ThreadPoolExecutor(max_workers=2, thread_name_prefix="push")

    async def _run_file(method, *args):
        """文件IO管线：读文件、写文件、目录操作、同步"""
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(_file_executor, method, *args)

    async def _run_data(method, *args):
        """数据管线：任务、课程、项目、书签等轻量查询"""
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(_data_executor, method, *args)

    def _build_multimodal_message(text: str, attachments: Optional[List[Dict[str, str]]] = None):
        """构建多模态消息：纯文本或 content parts 列表。

        如果有附件，返回 OpenAI 格式的 content parts：
        [{"type":"text",...}, {"type":"image_url",...}, {"type":"video_url",...}]
        如果没有附件，返回纯文本字符串。
        """
        if not attachments:
            return text
        from ..media_utils import build_multimodal_content
        return build_multimodal_content(text, attachments)

    async def _run_agent(method, *args, **kwargs):
        """Agent管线：LLM调用、工具执行（支持关键字参数传递 session_id）"""
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(_agent_executor, lambda: method(*args, **kwargs))

    async def _run_push(method, *args):
        """推送管线：聚合推送数据，独立于数据查询"""
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(_push_executor, method, *args)

    @app.post("/api/file/getFile")
    async def file_get(req: FileReadRequest):
        """获取文件内容"""
        if not check_path_access(req.path, "read"):
            return err(403, "路径不在允许的读取目录中")
        engine: FileSyncEngine = app.state.sync_engine
        result = await _run_file(engine.get_file, req.path)
        if result is None:
            return err(code=404, msg=f"File not found: {req.path}")
        content, entry = result
        return ok(data={"content": content, "entry": entry.to_dict()})

    @app.post("/api/file/putFile")
    async def file_put(req: FileWriteRequest):
        """创建或更新文件"""
        if not check_path_access(req.path, "write"):
            return err(403, "路径不在允许的写入目录中")
        engine: FileSyncEngine = app.state.sync_engine
        entry = await _run_file(engine.put_file, req.path, req.content)
        if entry is None:
            return err(msg=f"Failed to write file: {req.path}")
        return ok(data=entry.to_dict())

    @app.post("/api/file/removeFile")
    async def file_remove(req: FileRemoveRequest):
        """删除文件"""
        if not check_path_access(req.path, "write"):
            return err(403, "路径不在允许的写入目录中")
        engine: FileSyncEngine = app.state.sync_engine
        if await _run_file(engine.remove_file, req.path):
            return ok()
        return err(msg=f"Failed to remove: {req.path}")

    @app.post("/api/file/renameFile")
    async def file_rename(req: FileRenameRequest):
        """重命名/移动文件"""
        if not check_path_access(req.old_path, "write") or not check_path_access(req.new_path, "write"):
            return err(403, "路径不在允许的写入目录中")
        engine: FileSyncEngine = app.state.sync_engine
        if await _run_file(engine.rename_file, req.old_path, req.new_path):
            return ok()
        return err(msg=f"Failed to rename: {req.old_path} -> {req.new_path}")

    @app.post("/api/file/readDir")
    async def file_read_dir(req: DirReadRequest):
        """读取目录内容"""
        if not check_path_access(req.path, "read"):
            return err(403, "路径不在允许的读取目录中")
        engine: FileSyncEngine = app.state.sync_engine
        entries = await _run_file(engine.read_dir, req.path)
        return ok(data=[e.to_dict() for e in entries])

    @app.post("/api/file/createDir")
    async def file_create_dir(req: DirCreateRequest):
        """创建目录"""
        if not check_path_access(req.path, "write"):
            return err(403, "路径不在允许的写入目录中")
        engine: FileSyncEngine = app.state.sync_engine
        if await _run_file(engine.create_dir, req.path):
            return ok()
        return err(msg=f"Failed to create dir: {req.path}")

    @app.post("/api/file/scanTree")
    async def file_scan_tree(req: ScanRequest):
        """扫描文件树"""
        if not check_path_access(req.subdir, "read"):
            return err(403, "路径不在允许的读取目录中")
        engine: FileSyncEngine = app.state.sync_engine
        entries = await _run_file(engine.scan_file_tree, req.subdir)
        return ok(data=[e.to_dict() for e in entries])

    @app.post("/api/file/search")
    async def file_search(req: SearchRequest):
        """搜索文件（支持排序/筛选，对标 Explorer）"""
        if not check_path_access(req.subdir, "read"):
            return err(403, "路径不在允许的读取目录中")
        engine: FileSyncEngine = app.state.sync_engine
        results = await _run_file(engine.search_files, req.query, req.subdir,
                                  req.sort_by, req.order, req.type_filter)
        return ok(data=[e.to_dict() for e in results])

    @app.post("/api/file/stat")
    async def file_stat(req: FileReadRequest):
        """返回文件元信息（mtimeMs/size/exists），供 texpile 等编辑器使用"""
        if not check_path_access(req.path, "read"):
            return err(403, "路径不在允许的读取目录中")
        engine: FileSyncEngine = app.state.sync_engine
        # 复用 get_file，但只取元信息不读 content（开销略高但实现简单）
        result = await _run_file(engine.get_file, req.path)
        if result is None:
            return ok(data={"exists": False, "mtimeMs": 0, "size": 0})
        _content, entry = result
        return ok(data={
            "exists": True,
            "mtimeMs": int(entry.modified * 1000),
            "size": entry.size,
        })

    # ─── Texpile LaTeX 编译 API（可拔插）───────────────────────
    # 集成模式下，TS2 后端代替 Electron 主进程：
    #   /api/texpile/compile  → draft 引擎全量编译（lualatex + shipout 钩子提取 records）
    #   /api/texpile/typeset  → draft daemon 单段排版（常驻 lualatex，~1-2ms/段）
    #   /api/texpile/stop     → 终止所有运行中的编译/daemon
    #   /api/texpile/synctex  → SyncTeX 源↔PDF 跳转（调用 synctex CLI）
    # Lua 提取脚本位于 mcp/server/draft/（walker.lua / page-extract.lua / texd-loop.lua），
    # 原样从 texpile 原版 electron/lua/ 复制，未做修改。

    import asyncio.subprocess as _asub  # noqa: F401  确保 asyncio.subprocess 加载
    import time as _time_texpile
    from .draft import compile_draft, typeset_paragraph, stop_draft, get_engine_dir

    def _resolve_main_file(engine_sync: "FileSyncEngine", root_abs: Path, main_file: str) -> Path:
        """把 texpile 传来的 mainFile 解析成绝对路径。

        texpile 前端始终发送 **相对 root** 的 mainFile（源码 Cu(mainFile, root)），
        对嵌套项目（tex 位于 root 的子目录）也是如此，例如 root='Notes'、
        mainFile='chapters/main.tex'。因此一律按 root_abs / main_file 解析，
        绝不当作相对 workspace 的路径处理（否则子目录场景会找错目录，
        导致 lualatex "在当前目录找不到 tex 文件"）。
        仅当 mainFile 本身是绝对路径时才直接使用。
        """
        mf = (main_file or "").replace("\\", "/").lstrip("/")
        p = Path(main_file)
        if p.is_absolute():
            return p.resolve()
        return (root_abs / mf).resolve()

    @app.post("/api/texpile/compile")
    async def texpile_compile(req: Request):
        """Draft 全量编译：调 lualatex + shipout 钩子提取每页 records。
        Body: { root, mainFile, engine? }
        返回 DraftResult: { ok, ms, passes, count, paperW, paperH, colW, marginX, marginY, pages }
        或 { ok:false, error, ms, log?, superseded? }
        """
        try:
            body = await req.json()
        except Exception:
            body = {}
        root = body.get("root") or ""
        main_file = body.get("mainFile") or ""
        engine = body.get("engine") or "lualatex"

        if not root or not main_file:
            return err(msg="root and mainFile are required")

        # 先解析成绝对路径再鉴权：root 是相对 workspace 的路径，
        # check_path_access 直接对相对路径 resolve() 会用服务端 cwd，可能误判。
        engine_sync: FileSyncEngine = app.state.sync_engine
        root_abs = engine_sync._absolute_path(root)
        if not check_path_access(str(root_abs), "read") or not check_path_access(str(root_abs), "write"):
            return err(403, "工作目录不在允许的访问范围内")
        if not root_abs.exists() or not root_abs.is_dir():
            return err(msg=f"工作目录不存在: {root}")

        # texpile 始终发送相对 root 的 mainFile（前端 Cu(mainFile, root) 计算），
        # 故一律按 root_abs / main_file 解析；仅在绝对路径时才直接使用。
        main_abs = _resolve_main_file(engine_sync, root_abs, main_file)
        if not main_abs.exists():
            return err(msg=f"主文件不存在: {main_file}")

        # draft 模块需要 root 的绝对路径，main_file 相对 root
        try:
            main_rel = str(main_abs.relative_to(root_abs)).replace("\\", "/")
        except ValueError:
            main_rel = main_file.replace("\\", "/")

        try:
            result = await compile_draft(
                root=str(root_abs),
                main_file=main_rel,
                engine=engine,
                engine_dir=get_engine_dir(),
            )
            return ok(data=result)
        except Exception as e:
            return ok(data={"ok": False, "error": str(e), "ms": 0})

    @app.post("/api/texpile/stop")
    async def texpile_stop():
        """终止运行中的 draft 编译和 daemon。"""
        try:
            result = await stop_draft()
            return ok(data=result)
        except Exception as e:
            return ok(data={"ok": False, "error": str(e)})

    @app.post("/api/texpile/synctex")
    async def texpile_synctex(req: Request):
        """SyncTeX 转发：源↔PDF 跳转。调用系统 synctex CLI。
        Body: { root, file, line, side?: 'src'|'pdf', pdf? }
        返回: { ok, output?, page?, x?, y?, file?, line?, column? }
        """
        try:
            body = await req.json()
        except Exception:
            body = {}
        root = body.get("root") or ""
        if not root or not check_path_access(root, "read"):
            return err(403, "工作目录不在允许的访问范围内")
        engine_sync: FileSyncEngine = app.state.sync_engine
        root_abs = engine_sync._absolute_path(root)
        synctex_path = root_abs / "_draft" / "draft.synctex.gz"
        if not synctex_path.exists():
            # fallback: output/draft.synctex.gz
            synctex_path = root_abs / "output" / "draft.synctex.gz"
        if not synctex_path.exists():
            return ok(data={"ok": False, "error": "synctex 文件不存在（请先编译）"})

        side = body.get("side") or "src"
        try:
            if side == "src":
                # 源 → PDF：synctex view -i line:col:file -o output.pdf -d synctex
                file_abs = engine_sync._absolute_path(body.get("file") or "")
                line = int(body.get("line") or 1)
                col = int(body.get("column") or 0)
                args = ["view", "-i", f"{line}:{col}:{file_abs}", "-o",
                        str(root_abs / "_draft" / "draft.pdf")]
            else:
                # PDF → 源：synctex edit -o page:x:y:file.pdf
                page = int(body.get("page") or 1)
                x = float(body.get("x") or 0)
                y = float(body.get("y") or 0)
                args = ["edit", "-o", f"{page}:{x}:{y}:{root_abs / '_draft' / 'draft.pdf'}"]
            proc = await asyncio.create_subprocess_exec(
                "synctex", *args, "-d", str(root_abs / "_draft"),
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
            )
            try:
                stdout_b, stderr_b = await asyncio.wait_for(proc.communicate(), timeout=10)
            except asyncio.TimeoutError:
                proc.kill()
                return ok(data={"ok": False, "error": "synctex 超时"})
            output = stdout_b.decode("utf-8", errors="replace") if stdout_b else ""
            return ok(data={"ok": proc.returncode == 0, "output": output})
        except FileNotFoundError:
            return ok(data={"ok": False, "error": "synctex 未安装（请安装 TeX Live / MiKTeX）"})
        except Exception as e:
            return ok(data={"ok": False, "error": str(e)})

    @app.post("/api/texpile/typeset")
    async def texpile_typeset(req: Request):
        """增量段落排版（draft daemon）。常驻 lualatex 进程，单段 ~1-2ms。
        Body: { root, mainFile, text, hsize? }
        返回 ParagraphResult: { ok, records, stats, hsize, textheight }
        """
        try:
            body = await req.json()
        except Exception:
            body = {}
        root = body.get("root") or ""
        main_file = body.get("mainFile") or ""
        text = body.get("text") or ""
        hsize = body.get("hsize")

        if not root or not main_file:
            return err(msg="root and mainFile are required")
        engine_sync: FileSyncEngine = app.state.sync_engine
        root_abs = engine_sync._absolute_path(root)
        if not check_path_access(str(root_abs), "read"):
            return err(403, "工作目录不在允许的访问范围内")
        # texpile 始终发送相对 root 的 mainFile（见 texpile Cu(mainFile, root)）
        main_abs = _resolve_main_file(engine_sync, root_abs, main_file)
        try:
            main_rel = str(main_abs.relative_to(root_abs)).replace("\\", "/")
        except ValueError:
            main_rel = main_file.replace("\\", "/")

        try:
            result = await typeset_paragraph(
                root=str(root_abs),
                main_file=main_rel,
                text=text,
                hsize=hsize,
                engine_dir=get_engine_dir(),
            )
            return ok(data=result)
        except Exception as e:
            return ok(data={"ok": False, "error": str(e)})

    # ─── 文件下载 API ───────────────────────────────────────

    @app.get("/api/file/download/{file_path:path}")
    @app.head("/api/file/download/{file_path:path}")
    async def file_download(file_path: str, preview: bool = False, request: Request = None):
        """下载文件（二进制流，支持手机端下载），HEAD 请求仅检查文件是否存在（FastAPI 自动剥离 body）

        preview=True 时返回 inline 响应（浏览器内预览，不触发下载）
        """
        if not check_path_access(file_path, "read"):
            # 不在工作区目录：尝试资源索引放行（须鉴权）
            if not _is_resource_path(file_path, app.state.workspace_dir):
                raise HTTPException(status_code=403, detail="路径不在允许的读取目录中")
            if not await check_auth(request, load_api_config()):
                raise HTTPException(status_code=403, detail="未授权，请提供有效 token 或授权码")
        engine: FileSyncEngine = app.state.sync_engine
        # 支持绝对路径（如资源索引中的外部路径），直接检查是否存在
        p = Path(file_path)
        if p.is_absolute():
            abs_path = p.resolve()
        else:
            abs_path = engine._absolute_path(file_path)

        if not abs_path.exists():
            raise HTTPException(status_code=404, detail="File not found")
        if abs_path.is_dir():
            # 目录打包为 zip
            import tempfile
            import zipfile
            tmp = tempfile.NamedTemporaryFile(suffix=".zip", delete=False)
            try:
                with zipfile.ZipFile(tmp.name, 'w', zipfile.ZIP_DEFLATED) as zf:
                    for f in abs_path.rglob("*"):
                        if f.is_file() and not engine._should_ignore(f):
                            arcname = str(f.relative_to(abs_path))
                            zf.write(f, arcname)
                return FileResponse(
                    tmp.name,
                    media_type="application/zip",
                    filename=f"{abs_path.name}.zip",
                )
            except Exception as e:
                os.unlink(tmp.name)
                raise HTTPException(status_code=500, detail=str(e))

        # 单文件下载
        mime_map = {
            ".pdf": "application/pdf",
            ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
            ".gif": "image/gif", ".svg": "image/svg+xml",
            ".webp": "image/webp", ".bmp": "image/bmp", ".ico": "image/x-icon",
            ".tif": "image/tiff", ".tiff": "image/tiff",
            ".avif": "image/avif", ".jfif": "image/jpeg", ".apng": "image/apng",
            ".heic": "image/heic", ".heif": "image/heif",
            ".mp4": "video/mp4", ".mp3": "audio/mpeg", ".wav": "audio/wav",
            ".md": "text/markdown", ".txt": "text/plain", ".rmd": "text/plain",
            ".py": "text/x-python", ".json": "application/json",
            ".html": "text/html", ".htm": "text/html", ".css": "text/css",
            ".zip": "application/zip", ".tar": "application/x-tar",
            ".gz": "application/gzip",
        }
        media_type = mime_map.get(abs_path.suffix.lower(), "application/octet-stream")
        is_media = media_type.startswith(("image/", "video/", "audio/"))
        if preview or is_media:
            return FileResponse(abs_path, media_type=media_type)
        return FileResponse(abs_path, media_type=media_type, filename=abs_path.name)

    # ─── 文件上传 API（支持手机端上传）─────────────────────

    @app.post("/api/file/upload")
    async def file_upload(
        files: List[UploadFile] = File(...),
        path: str = Form(""),
    ):
        """
        上传文件（multipart/form-data，支持多文件）

        - files: 上传的文件列表
        - path: 目标目录（相对于工作区，默认根目录）
        """
        if not check_path_access(path if path else os.getcwd(), "write"):
            return err(403, "路径不在允许的写入目录中")
        engine: FileSyncEngine = app.state.sync_engine
        target_dir = engine._absolute_path(path) if path else engine.workspace_dir

        # 确保目标目录存在
        target_dir.mkdir(parents=True, exist_ok=True)

        uploaded = []
        for upload_file in files:
            # 安全检查
            target_path = target_dir / upload_file.filename
            try:
                target_path.relative_to(engine.workspace_dir)
            except ValueError:
                continue

            # 写入文件
            try:
                content = await upload_file.read()
                target_path.write_bytes(content)

                rel_path = engine._relative_path(target_path)
                stat = target_path.stat()
                entry_dict = {
                    "path": rel_path,
                    "name": target_path.name,
                    "is_dir": False,
                    "size": stat.st_size,
                    "modified": stat.st_mtime,
                    "ext": target_path.suffix.lower(),
                }
                uploaded.append(entry_dict)

                # 触发变更通知
                if sync_engine._on_change_callback:
                    event = SyncEvent(
                        event_type="created",
                        path=rel_path,
                        timestamp=time.time(),
                        size=stat.st_size,
                    )
                    asyncio.create_task(sync_engine._on_change_callback(event))

            except Exception as e:
                logger.error(f"Upload write error: {e}")

        return ok(data={
            "uploaded": uploaded,
            "count": len(uploaded),
            "target_dir": path,
        })

    # ─── 批量下载 API ──────────────────────────────────────

    @app.post("/api/file/batchDownload")
    async def file_batch_download(request: Request):
        """批量下载文件（打包为 zip）"""
        body = await request.json()
        paths: List[str] = body.get("paths", [])
        for p in paths:
            if not check_path_access(p, "read"):
                return err(403, f"路径不在允许的读取目录中: {p}")

        if not paths:
            return err(msg="No files specified")

        engine: FileSyncEngine = app.state.sync_engine
        import tempfile
        import zipfile

        tmp = tempfile.NamedTemporaryFile(suffix=".zip", delete=False)
        try:
            with zipfile.ZipFile(tmp.name, 'w', zipfile.ZIP_DEFLATED) as zf:
                for rel_path in paths:
                    abs_path = engine._absolute_path(rel_path)
                    if abs_path.exists() and abs_path.is_file():
                        zf.write(abs_path, abs_path.name)
            return FileResponse(
                tmp.name,
                media_type="application/zip",
                filename=f"ts2_download_{int(time.time())}.zip",
            )
        except Exception as e:
            os.unlink(tmp.name)
            raise HTTPException(status_code=500, detail=str(e))

    # ─── Office 转 PDF API ─────────────────────────────────

    @app.post("/api/file/convert-to-pdf")
    async def file_convert_to_pdf(request: Request):
        """将 Office 文件（docx/xlsx/pptx）转换为 PDF 并返回文件"""
        body = await request.json()
        file_path: str = body.get("file_path", "")
        if not file_path:
            return err(msg="缺少 file_path")

        abs_path = sync_engine._absolute_path(file_path)
        if not os.path.isfile(abs_path):
            return err(msg=f"文件不存在: {file_path}")

        ext = os.path.splitext(abs_path)[1].lower()
        if ext not in (".docx", ".xlsx", ".pptx"):
            return err(msg=f"不支持的文件类型: {ext}")

        try:
            pdf_path = await asyncio.get_event_loop().run_in_executor(
                None, _convert_office_to_pdf, str(abs_path)
            )
            return FileResponse(
                pdf_path,
                media_type="application/pdf",
                filename=os.path.splitext(os.path.basename(abs_path))[0] + ".pdf",
            )
        except Exception as e:
            return err(msg=f"转换失败: {e}")

    # ─── 同步 API（参考思源笔记 /api/sync/*）──────────────

    @app.post("/api/sync/performSync")
    async def sync_perform(req: SyncRequest):
        """
        执行增量同步（参考 better-sync syncHandler + siyuan-android syncData）
        
        流程：
        1. 检查同步锁（参考 siyuan-android syncing 互斥）
        2. 扫描变更（参考 better-sync scanDirectory）
        3. 冲突检测（参考 better-sync detectConflict）
        4. 创建冲突副本（参考 better-sync createConflictFile）
        5. 通过 WebSocket 推送变更
        6. 记录同步历史（参考 better-sync SyncHistory）
        """
        engine: FileSyncEngine = app.state.sync_engine

        # 检查同步锁（参考 siyuan-android: if (syncing) return）
        if engine.get_sync_status() == SyncStatus.InProgress:
            return ok(data={
                "synced": 0,
                "changes": [],
                "status": "InProgress",
                "message": "Sync already in progress",
                "direction": req.direction,
                "timestamp": time.time(),
            })

        # 执行增量同步
        changes, status = engine.perform_incremental_sync(mobile_switch=req.mobile_switch)

        # 转换变更列表为 API 格式
        api_changes = []
        for c in changes:
            change_dict = {
                "type": c.type,
                "path": c.path,
                "hash": c.hash,
                "size": c.size,
                "operation": c.operation,
            }
            if c.conflict_path:
                change_dict["conflict_path"] = c.conflict_path
            api_changes.append(change_dict)

        # 通过 WebSocket 推送变更
        if api_changes and hasattr(app.state, 'ws_manager'):
            for change in api_changes:
                await app.state.ws_manager.notify_path_change(
                    path=change["path"],
                    change_type=change["type"],
                    data=change,
                )
            if any(c["type"] in ("created", "deleted", "renamed", "conflict") for c in api_changes):
                await app.state.ws_manager.push_reload_filetree()

        return ok(data={
            "synced": len(api_changes),
            "changes": api_changes,
            "status": status.name,
            "conflicts": sum(1 for c in api_changes if c["type"] == "conflict"),
            "direction": req.direction,
            "mobile_switch": req.mobile_switch,
            "timestamp": time.time(),
        })

    @app.post("/api/sync/getSyncInfo")
    async def sync_info():
        """获取同步状态（参考 better-sync 同步状态管理）"""
        engine: FileSyncEngine = app.state.sync_engine
        return ok(data={
            "workspace": str(engine.workspace_dir),
            "file_stats": engine.get_file_stats(),
            "watching": engine._watching,
            "sync_status": engine.get_sync_status().name,
            "last_sync_time": engine._last_sync_time,
            "sync_history": engine.get_sync_history(limit=10),
        })

    @app.post("/api/sync/getSyncHistory")
    async def sync_history(request: Request):
        """获取同步历史（参考 better-sync SyncHistory）"""
        engine: FileSyncEngine = app.state.sync_engine
        body = await request.json()
        limit = body.get("limit", 20)
        return ok(data=engine.get_sync_history(limit=limit))

    @app.post("/api/sync/startWatch")
    async def sync_start_watch():
        """启动文件监听"""
        engine: FileSyncEngine = app.state.sync_engine
        await engine.start_watching()
        return ok(data={"watching": True})

    @app.post("/api/sync/stopWatch")
    async def sync_stop_watch():
        """停止文件监听"""
        engine: FileSyncEngine = app.state.sync_engine
        await engine.stop_watching()
        return ok(data={"watching": False})

    # ─── 通知 API（参考思源笔记 /api/notification/*）─────

    @app.post("/api/notification/pushMsg")
    async def push_msg(request: Request):
        """推送消息"""
        body = await request.json()
        msg = body.get("msg", "")
        timeout = body.get("timeout", 5000)
        await app.state.ws_manager.push_msg(msg, timeout)
        return ok()

    # ─── 代码执行 ─────────────────────────────────────────

    @app.post("/api/exec/run")
    async def exec_run(request: Request):
        """执行代码，返回 stdout/stderr/exit_code，支持停止"""
        body = await request.json()
        language = body.get("language", "python")
        code = body.get("code", "")
        file_path = body.get("file_path", "")

        JAVA_NEEDS_COMPILE = True  # Java 需 javac + java 两步

        lang_map = {
            "python": [sys.executable],
            "py": [sys.executable],
            "python3": [sys.executable],
            "javascript": ["node"],
            "js": ["node"],
            "typescript": ["npx", "--yes", "tsx"],
            "ts": ["npx", "--yes", "tsx"],
            "r": ["Rscript"],
            "bash": ["bash"],
            "sh": ["sh"],
            "powershell": ["powershell", "-NoProfile", "-Command"],
            "pwsh": ["pwsh", "-NoProfile", "-Command"],
            "cmd": ["cmd", "/c"],
            "java": ["javac"],
            "kt": ["kotlinc"],
            "kotlin": ["kotlinc"],
            "c": ["gcc"],
            "cpp": ["g++"],
            "go": ["go", "run"],
            "rs": ["rustc"],
            "rust": ["rustc"],
        }

        cmd_base = lang_map.get(language, [sys.executable])

        # Build execution environment: force UTF-8 everywhere
        exec_env = os.environ.copy()
        exec_env["PYTHONUNBUFFERED"] = "1"
        exec_env["PYTHONIOENCODING"] = "utf-8"
        if "LANG" not in exec_env:
            exec_env["LANG"] = "C.UTF-8"
        conda_prefix = os.environ.get("CONDA_PREFIX") or os.environ.get("CONDA_EXE", "")
        if conda_prefix:
            conda_base = Path(conda_prefix).parent if "CONDA_EXE" in os.environ else Path(conda_prefix)
            scripts_dir = str(conda_base / "Scripts")
            condabin_dir = str(conda_base / "condabin")
            if "PATH" in exec_env:
                exec_env["PATH"] = f"{scripts_dir};{condabin_dir};{exec_env['PATH']}"

        # ── 编译型语言辅助函数 ──
        async def _run_java(source_path: Path, cwd: str):
            """javac 编译 → java 运行"""
            class_dir = source_path.parent
            class_name = source_path.stem
            compile_proc = await asyncio.create_subprocess_exec(
                "javac", "-d", str(class_dir), str(source_path),
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
                cwd=cwd, env=exec_env,
            )
            cout, cerr = await compile_proc.communicate()
            if compile_proc.returncode != 0:
                return cout + b"\n[Java Compile Error]\n" + cerr
            run_proc = await asyncio.create_subprocess_exec(
                "java", "-cp", str(class_dir), class_name,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
                cwd=cwd, env=exec_env,
            )
            rout, rerr = await run_proc.communicate()
            return rout + b"\n" + rerr

        COMPILE_LANGS = {"java", "c", "cpp", "c++", "rs", "rust", "go"}

        try:
            if file_path:
                full_path = (Path(app.state.workspace_dir) / file_path).resolve()
                if language in COMPILE_LANGS:
                    display_cmd = f"{cmd_base[0]} {file_path}"
                else:
                    display_cmd = f"{cmd_base[0]} {file_path}"
                if language == "java":
                    combined = await _run_java(full_path, str(full_path.parent))
                    stdout = combined.decode("utf-8", errors="replace")
                    return ok(data={"exit_code": 0, "stdout": stdout, "stderr": "", "command": display_cmd})
                proc = await asyncio.create_subprocess_exec(
                    *cmd_base, str(full_path),
                    stdout=asyncio.subprocess.PIPE,
                    stderr=asyncio.subprocess.PIPE,
                    cwd=full_path.parent,
                    env=exec_env,
                )
            elif code:
                if language in ("powershell", "pwsh"):
                    proc = await asyncio.create_subprocess_exec(
                        *cmd_base, code,
                        stdout=asyncio.subprocess.PIPE,
                        stderr=asyncio.subprocess.PIPE,
                        cwd=app.state.workspace_dir,
                        env=exec_env,
                    )
                    display_cmd = f"{cmd_base[0]} <inline>"
                elif language in ("bash", "sh", "cmd"):
                    proc = await asyncio.create_subprocess_exec(
                        *cmd_base, code,
                        stdout=asyncio.subprocess.PIPE,
                        stderr=asyncio.subprocess.PIPE,
                        cwd=app.state.workspace_dir,
                        env=exec_env,
                    )
                    display_cmd = f"{cmd_base[0]} <inline>"
                elif language in COMPILE_LANGS:
                    # Write code to temp file, then compile & run
                    tmp_dir = Path(app.state.workspace_dir) / ".ts2_data" / "exec_tmp"
                    tmp_dir.mkdir(parents=True, exist_ok=True)
                    ext_map = {"java": ".java", "c": ".c", "cpp": ".cpp", "c++": ".cpp",
                               "rs": ".rs", "rust": ".rs", "go": ".go"}
                    suffix = ext_map.get(language, ".java")
                    tmp_file = tmp_dir / f"exec_{uuid.uuid4().hex}{suffix}"
                    tmp_file.write_text(code, encoding="utf-8")
                    if language == "java":
                        combined = await _run_java(tmp_file, str(app.state.workspace_dir))
                        stdout = combined.decode("utf-8", errors="replace")
                        return ok(data={"exit_code": 0, "stdout": stdout, "stderr": "", "command": f"{cmd_base[0]} <code>"})
                    display_cmd = f"{cmd_base[0]} <code>"
                    proc = await asyncio.create_subprocess_exec(
                        *cmd_base, str(tmp_file),
                        stdout=asyncio.subprocess.PIPE,
                        stderr=asyncio.subprocess.PIPE,
                        cwd=app.state.workspace_dir,
                        env=exec_env,
                    )
                else:
                    # Write code to a UTF-8 temp file to avoid GBK cmdline corruption
                    tmp_dir = Path(app.state.workspace_dir) / ".ts2_data" / "exec_tmp"
                    tmp_dir.mkdir(parents=True, exist_ok=True)
                    ext_map = {"python": ".py", "javascript": ".js", "js": ".js",
                               "typescript": ".ts", "ts": ".ts", "r": ".R",
                               "java": ".java", "c": ".c", "cpp": ".cpp", "c++": ".cpp",
                               "rs": ".rs", "rust": ".rs", "go": ".go"}
                    suffix = ext_map.get(language, ".py")
                    tmp_file = tmp_dir / f"exec_{uuid.uuid4().hex}{suffix}"
                    tmp_file.write_text(code, encoding="utf-8")
                    proc = await asyncio.create_subprocess_exec(
                        *cmd_base, str(tmp_file),
                        stdout=asyncio.subprocess.PIPE,
                        stderr=asyncio.subprocess.PIPE,
                        cwd=app.state.workspace_dir,
                        env=exec_env,
                    )
                    display_cmd = f"{cmd_base[0]} <code>"
            else:
                return ok(data={"exit_code": -1, "stdout": "", "stderr": "No code or file_path provided"})

            # Kill previous, store current
            if app.state._exec_proc and app.state._exec_proc.returncode is None:
                app.state._exec_proc.terminate()
            app.state._exec_proc = proc

            try:
                stdout_bytes, stderr_bytes = await asyncio.wait_for(proc.communicate(), timeout=120)
            except asyncio.CancelledError:
                if proc.returncode is None:
                    proc.terminate()
                raise

            stdout = stdout_bytes.decode("utf-8", errors="replace")
            stderr = stderr_bytes.decode("utf-8", errors="replace")

            # Clean up temp file
            if not file_path and code and language not in ("powershell", "pwsh", "bash", "sh", "cmd"):
                try:
                    tmp_file.unlink(missing_ok=True)
                except Exception:
                    pass

            if proc.returncode is None:
                exit_code = -1
            else:
                exit_code = proc.returncode

            return ok(data={
                "exit_code": exit_code,
                "stdout": stdout,
                "stderr": stderr,
                "command": display_cmd,
            })
        except asyncio.TimeoutError:
            if proc.returncode is None:
                proc.terminate()
            return ok(data={"exit_code": -1, "stdout": "", "stderr": "Execution timed out (>120s)"})
        except asyncio.CancelledError:
            return ok(data={"exit_code": -2, "stdout": "", "stderr": "Execution cancelled"})
        except FileNotFoundError:
            return ok(data={"exit_code": -1, "stdout": "", "stderr": f"Interpreter not found: {cmd_base[0]}"})
        except Exception as e:
            return ok(data={"exit_code": -1, "stdout": "", "stderr": str(e)})

    @app.post("/api/exec/stop")
    async def exec_stop():
        """强制停止当前正在执行的代码"""
        proc = app.state._exec_proc
        if proc and proc.returncode is None:
            try:
                proc.kill()
            except Exception:
                pass
            return ok(msg="已停止")
        return ok(msg="没有正在执行的任务")

    # ─── WebSocket（参考思源笔记 /ws）─────────────────────

    @app.websocket("/ws")
    async def websocket_endpoint(websocket: WebSocket):
        """
        WebSocket 端点

        连接参数：
        - app: 应用ID
        - id: 会话ID
        - type: 客户端类型 (main, filetree, editor, sync)

        消息格式（参考思源笔记）：
        {"cmd": "...", "reqId": 1, "param": {...}}

        事件推送格式：
        {"cmd": "...", "code": 0, "msg": "", "data": {...}}
        """
        # 鉴权检查
        config = load_api_config()
        if not _check_ws_auth(websocket, config):
            await websocket.close(code=4001, reason="未授权")
            return

        await websocket.accept()

        # 解析连接参数
        app_id = websocket.query_params.get("app", str(uuid.uuid4())[:8])
        session_id = websocket.query_params.get("id", str(uuid.uuid4())[:8])
        client_type = websocket.query_params.get("type", "main")

        # 注册会话
        session = await app.state.ws_manager.add_session(
            websocket, app_id, session_id, client_type
        )

        try:
            # 发送连接成功事件
            await websocket.send_text(_json_dumps({
                "cmd": "connected",
                "code": 0,
                "msg": "",
                "data": {
                    "app_id": app_id,
                    "session_id": session_id,
                    "type": client_type,
                }
            }, ensure_ascii=False))

            # 消息循环
            while True:
                raw = await websocket.receive_text()
                try:
                    msg = _json_loads(raw)
                    cmd = msg.get("cmd", "")
                    req_id = msg.get("reqId", 0)
                    param = msg.get("param", {})

                    # 处理命令
                    result = await _handle_ws_command(
                        app, cmd, param, app_id, session_id
                    )
                    result["reqId"] = req_id
                    await websocket.send_text(_json_dumps(result, ensure_ascii=False))

                except json.JSONDecodeError:
                    await websocket.send_text(_json_dumps({
                        "cmd": "error", "code": -1,
                        "msg": "Invalid JSON", "reqId": 0
                    }))
                except Exception as e:
                    logger.error(f"WS command error: {e}")
                    await websocket.send_text(_json_dumps({
                        "cmd": "error", "code": -1,
                        "msg": str(e), "reqId": msg.get("reqId", 0)
                    }))

        except WebSocketDisconnect:
            pass
        finally:
            await app.state.ws_manager.remove_session(app_id, session_id)

    # ─── 协同：创建物理副本（右键"协同打开副本"）────
    from .collab import create_copy
    from pydantic import BaseModel

    class CollabCopyRequest(BaseModel):
        path: str

    @app.post("/api/collab/createCopy")
    async def collab_create_copy(req: CollabCopyRequest, request: Request):
        config = load_api_config()
        if not check_auth(request, config):  # type: ignore[name-defined]
            return err(403, "未授权")
        try:
            result = create_copy(req.path, _read_collab_initial, _write_collab_persist)
            return ok(data=result)
        except Exception as e:
            logger.error(f"collab createCopy failed {req.path}: {e}")
            return err(msg=str(e))

    # ─── 协同 WebSocket（Loro 协同编辑，正式项目）────────

    from .collab import CollabManager

    def _read_collab_initial(p: str) -> Optional[str]:
        """协同文件初始文本（用于服务端权威 LoroDoc 初始化）。"""
        try:
            base = Path(str(app.state.workspace_dir)).resolve()
            fp = (base / p.lstrip("/")).resolve()
            if not str(fp).startswith(str(base)):
                return None
            if fp.is_file():
                return fp.read_text(encoding="utf-8", errors="replace")
        except Exception:
            pass
        return None

    def _write_collab_persist(p: str, text: str) -> None:
        """协同权威 LoroDoc 文本写回文件（持久化）。"""
        try:
            base = Path(str(app.state.workspace_dir)).resolve()
            fp = (base / p.lstrip("/")).resolve()
            if not str(fp).startswith(str(base)):
                return
            fp.parent.mkdir(parents=True, exist_ok=True)
            fp.write_text(text, encoding="utf-8")
        except Exception as e:
            logger.error(f"collab persist write failed {p}: {e}")

    collab_manager = CollabManager(
        load_text=_read_collab_initial,
        persist_text=_write_collab_persist,
    )
    app.state.collab_manager = collab_manager

    @app.websocket("/ws/collab/{path:path}")
    async def collab_ws(websocket: WebSocket, path: str):
        """协同编辑 WebSocket。纯转发：CRDT 在客户端，服务端做房间与快照协商。"""
        config = load_api_config()
        if not _check_ws_auth(websocket, config):
            await websocket.close(code=4001, reason="未授权")
            return
        await websocket.accept()
        await collab_manager.handle(
            websocket, path,
            send_json=lambda m: websocket.send_text(_json_dumps(m, ensure_ascii=False)),
        )

    # ─── 终端 WebSocket（xterm.js）────────────────────────

    @app.websocket("/api/terminal")
    async def terminal_ws(websocket: WebSocket):
        """xterm.js 终端 WebSocket — 桥接 shell stdio（使用 PTY 支持交互式 shell）

        Query params:
            cwd: 工作目录（绝对路径或相对 workspace_dir 的路径，可选）
            cols/rows: 初始 PTY 尺寸（可选，默认 100x30）
        """
        # 鉴权检查
        config = load_api_config()
        if not _check_ws_auth(websocket, config):
            await websocket.close(code=4001, reason="未授权")
            return

        await websocket.accept()
        loop = asyncio.get_event_loop()

        # ── 解析 cwd（支持 texpile 在工作目录启 shell）──
        cwd_param = websocket.query_params.get("cwd", "").strip()
        cwd_abs: Optional[Path] = None
        if cwd_param:
            try:
                engine: FileSyncEngine = app.state.sync_engine
                cwd_abs = engine._absolute_path(cwd_param)
                # 如果路径已经是绝对路径且存在，直接用
                if not cwd_abs.exists() or not cwd_abs.is_dir():
                    p = Path(cwd_param)
                    if p.is_absolute() and p.exists() and p.is_dir():
                        cwd_abs = p
                    else:
                        cwd_abs = None
            except Exception:
                # 退化：尝试当绝对路径用
                p = Path(cwd_param)
                if p.is_absolute() and p.exists() and p.is_dir():
                    cwd_abs = p
                else:
                    cwd_abs = None

        # ── 检测 conda 环境 ──
        conda_prefix = os.environ.get("CONDA_PREFIX", "")
        conda_base = os.environ.get("CONDA_EXE", "")
        if conda_base:
            conda_base = str(Path(conda_base).parent)
        elif conda_prefix:
            conda_base = str(Path(conda_prefix).parent)
        # 使用 cmd.exe 而非 pwsh.exe — cmd 原生支持 conda.bat
        shell = "cmd.exe" if os.name == "nt" else "bash"
        proc_env = os.environ.copy()
        if conda_base:
            scripts_dir = os.path.join(conda_base, "Scripts")
            condabin_dir = os.path.join(conda_base, "condabin")
            proc_env["PATH"] = f"{scripts_dir};{condabin_dir};{proc_env.get('PATH', '')}"

        # 工作目录：优先用 cwd_param，否则用 workspace_dir
        spawn_cwd = str(cwd_abs) if cwd_abs else str(app.state.workspace_dir)

        # ── 尝试 PTY 路径 ──
        # 初始尺寸：texpileTerminal spawn 时传入 cols/rows；xterm.js 端用默认值
        try:
            init_cols = int(websocket.query_params.get("cols", "100"))
            init_rows = int(websocket.query_params.get("rows", "30"))
        except Exception:
            init_cols, init_rows = 100, 30
        if init_cols < 1 or init_rows < 1:
            init_cols, init_rows = 100, 30

        import winpty as _winpty
        import shutil as _shutil
        # pywinpty 3.x 接口要点：
        #   - spawn 的 appname 必须是完整路径（'cmd.exe' 会 os error 2）
        #   - write 要 str（bytes 报 TypeError）
        #   - read(blocking=False) 无 length 参数，返回 str（已解码）
        #   - spawn 的 env 是 "\0" 分隔的字符串（非 dict）
        env_str = '\0'.join(f'{k}={v}' for k, v in proc_env.items()) + '\0'
        # 解析 shell 完整路径：cmd.exe → C:\Windows\System32\cmd.exe
        shell_full = _shutil.which(shell) or shell
        pty = None
        def _create_pty():
            p = _winpty.PTY(init_cols, init_rows)
            p.spawn(shell_full, cwd=spawn_cwd, env=env_str)
            return p
        try:
            pty = await loop.run_in_executor(None, _create_pty)
            logger.info(f"PTY spawn OK: shell={shell_full} cwd={spawn_cwd} cols={init_cols} rows={init_rows}")
        except Exception as _e:
            logger.warning(f"PTY spawn failed, falling back to subprocess: {_e}")
            pty = None

        if pty is not None:
            async def _run(method, *args):
                return await loop.run_in_executor(None, getattr(pty, method), *args)

            # 第一条消息：告知前端 shell 名（texpile 用来选 sentinel 语法）
            try:
                await websocket.send_text("__shell__:" + shell)
            except Exception:
                pass

            async def read_pty():
                while True:
                    try:
                        # pywinpty 3.x: read(blocking=False) 返回 str，无数据时返回 ''
                        raw = await _run("read", False)
                        if not raw:
                            await asyncio.sleep(0.02)
                            continue
                        if isinstance(raw, str):
                            await websocket.send_bytes(raw.encode("utf-8", errors="replace"))
                        else:
                            await websocket.send_bytes(raw)
                    except EOFError:
                        logger.info("terminal read_pty: EOF")
                        break
                    except Exception as e:
                        logger.warning(f"terminal read_pty exit: {type(e).__name__}: {e}")
                        break

            async def write_pty():
                while True:
                    try:
                        msg = await websocket.receive_text()
                        if msg.startswith("__resize__"):
                            parts = msg.split(":")
                            if len(parts) == 3:
                                try:
                                    await _run("set_size", int(parts[1]), int(parts[2]))
                                except Exception:
                                    pass
                            continue
                        if await _run("isalive"):
                            # pywinpty 3.x: write 要 str（不是 bytes）
                            await _run("write", msg)
                        else:
                            logger.info("terminal write_pty: PTY not alive")
                            break
                    except WebSocketDisconnect:
                        break
                    except Exception as e:
                        logger.warning(f"terminal write_pty error: {type(e).__name__}: {e}")
                        break

            tasks = [asyncio.create_task(read_pty()), asyncio.create_task(write_pty())]
            try:
                await asyncio.gather(*tasks)
            finally:
                for t in tasks:
                    t.cancel()
                try:
                    if await _run("isalive"):
                        await _run("write", "\x03")
                except Exception:
                    pass
            logger.info("terminal PTY session ended")
            return

        # ── Fallback: 无 PTY ──
        try:
            proc = await asyncio.create_subprocess_exec(
                shell,
                stdin=asyncio.subprocess.PIPE,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.STDOUT,
                cwd=spawn_cwd,
                env=proc_env,
            )
        except FileNotFoundError:
            shell = "powershell.exe" if os.name == "nt" else "sh"
            proc = await asyncio.create_subprocess_exec(
                shell,
                stdin=asyncio.subprocess.PIPE,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.STDOUT,
                cwd=spawn_cwd,
                env=proc_env,
            )

        # 第一条消息：告知前端 shell 名（fallback 路径也要告知）
        try:
            await websocket.send_text("__shell__:" + shell)
        except Exception:
            pass

        async def read_proc():
            while True:
                data = await proc.stdout.read(4096)
                if not data:
                    break
                try:
                    await websocket.send_bytes(data)
                except Exception:
                    break

        async def write_proc():
            # fallback 路径下 stdin 是 pipe，cmd.exe 不回显输入。
            # 后端收到用户输入时先回显给前端，再写入 stdin，让用户能看到自己打的内容。
            # （PTY 路径不需要这层：console 回显由 PTY 自己处理）
            while True:
                try:
                    msg = await websocket.receive_text()
                    if msg.startswith("__resize__"):
                        continue
                    if proc.stdin and not proc.stdin.is_closing():
                        # 回显：把用户输入原样发回前端（bytes）
                        try:
                            await websocket.send_bytes(msg.encode("utf-8"))
                        except Exception:
                            pass
                        proc.stdin.write(msg.encode("utf-8"))
                        await proc.stdin.drain()
                except WebSocketDisconnect:
                    break
                except Exception:
                    break

        tasks = [asyncio.create_task(read_proc()), asyncio.create_task(write_proc())]
        try:
            await asyncio.gather(*tasks)
        finally:
            for t in tasks:
                t.cancel()
            if proc.returncode is None:
                proc.terminate()
                try:
                    await asyncio.wait_for(proc.wait(), timeout=3)
                except asyncio.TimeoutError:
                    proc.kill()

    async def _handle_ws_command(app, cmd: str, param: dict,
                                  app_id: str, session_id: str) -> dict:
        """处理 WebSocket 命令"""
        engine: FileSyncEngine = app.state.sync_engine
        manager: WebSocketManager = app.state.ws_manager

        if cmd == "subscribe":
            path = param.get("path", "")
            await manager.subscribe(app_id, session_id, path)
            return {"cmd": "subscribe", "code": 0, "data": {"path": path}}

        elif cmd == "unsubscribe":
            path = param.get("path", "")
            await manager.unsubscribe(app_id, session_id, path)
            return {"cmd": "unsubscribe", "code": 0, "data": {"path": path}}

        elif cmd == "getFile":
            fpath = param.get("path", "")
            if not check_path_access(fpath, "read"):
                return {"cmd": "getFile", "code": 403, "msg": "路径不在允许的读取目录中"}
            result = engine.get_file(fpath)
            if result:
                content, entry = result
                return {"cmd": "getFile", "code": 0, "data": {"content": content, "entry": entry.to_dict()}}
            return {"cmd": "getFile", "code": 404, "msg": "File not found"}

        elif cmd == "putFile":
            fpath = param.get("path", "")
            if not check_path_access(fpath, "write"):
                return {"cmd": "putFile", "code": 403, "msg": "路径不在允许的写入目录中"}
            entry = engine.put_file(fpath, param.get("content", ""))
            if entry:
                return {"cmd": "putFile", "code": 0, "data": entry.to_dict()}
            return {"cmd": "putFile", "code": -1, "msg": "Write failed"}

        elif cmd == "readDir":
            dpath = param.get("path", "")
            if dpath and not check_path_access(dpath, "read"):
                return {"cmd": "readDir", "code": 403, "msg": "路径不在允许的读取目录中"}
            entries = engine.read_dir(dpath)
            return {"cmd": "readDir", "code": 0, "data": [e.to_dict() for e in entries]}

        elif cmd == "search":
            entries = engine.search_files(param.get("query", ""), param.get("subdir", ""))
            return {"cmd": "search", "code": 0, "data": [e.to_dict() for e in entries]}

        elif cmd == "ping":
            return {"cmd": "pong", "code": 0, "data": {"time": time.time()}}

        else:
            return {"cmd": cmd, "code": -1, "msg": f"Unknown command: {cmd}"}

    # ─── 数据 API（任务、书签、项目、课程）──────────────────

    # ── 同步辅助函数（在线程池中执行，避免阻塞事件循环）──

    def _read_tasks_data(workspace_dir: str):
        """同步读取任务数据（在线程池中执行）"""
        tb_path = Path(workspace_dir) / "task_board.json"
        if tb_path.exists():
            try:
                content = tb_path.read_text(encoding="utf-8")
                return _json_loads(content)
            except Exception as e:
                logger.warning(f"Read task_board.json failed: {e}")

        db_path = Path(workspace_dir) / "data" / "automation.db"
        if not db_path.exists():
            return []
        try:
            conn = sqlite3.connect(str(db_path))
            conn.row_factory = sqlite3.Row
            cursor = conn.cursor()
            cursor.execute("SELECT name FROM sqlite_master WHERE type='table'")
            tables = [r[0] for r in cursor.fetchall()]
            task_table = None
            for t in tables:
                if 'task' in t.lower() or 'todo' in t.lower():
                    task_table = t
                    break
            if not task_table and tables:
                task_table = tables[0]
            if task_table:
                cursor.execute(f"SELECT * FROM {task_table} LIMIT 500")
                rows = cursor.fetchall()
                result = [dict(r) for r in rows]
            else:
                result = []
            conn.close()
            return result
        except Exception as e:
            return []

    def _read_bookmarks_data(workspace_dir: str):
        """同步读取书签数据"""
        bm_path = Path(workspace_dir) / "bookmarks.json"
        if not bm_path.exists():
            return []
        try:
            content = bm_path.read_text(encoding="utf-8")
            return _json_loads(content)
        except Exception:
            return []

    def _add_bookmark_data(workspace_dir: str, bookmark: dict):
        """同步添加书签"""
        bm_path = Path(workspace_dir) / "bookmarks.json"
        bookmarks = _read_bookmarks_data(workspace_dir)
        if "id" not in bookmark or not bookmark["id"]:
            bookmark["id"] = str(uuid.uuid4())
        bookmarks.append(bookmark)
        bm_path.write_text(_json_dumps(bookmarks, ensure_ascii=False, indent=2), encoding="utf-8")
        return bookmark, None

    def _delete_bookmark_data(workspace_dir: str, bookmark_id: str):
        """同步删除书签"""
        bm_path = Path(workspace_dir) / "bookmarks.json"
        bookmarks = _read_bookmarks_data(workspace_dir)
        filtered = [b for b in bookmarks if b.get("id") != bookmark_id and b.get("url") != bookmark_id]
        if len(filtered) == len(bookmarks):
            return None, "未找到该书签"
        bm_path.write_text(_json_dumps(filtered, ensure_ascii=False, indent=2), encoding="utf-8")
        return None, None

    def _read_projects_data(workspace_dir: str):
        """同步读取项目列表"""
        workspace = Path(workspace_dir)
        projects = []
        for p in workspace.rglob(".ts2_project.json"):
            try:
                data = _json_loads(p.read_text(encoding="utf-8"))
                proj_dir = p.parent
                file_count = sum(1 for f in proj_dir.rglob("*") if f.is_file() and not f.name.startswith("."))
                data["file_count"] = file_count
                data["path"] = str(p.parent.relative_to(workspace)).replace("\\", "/")
                projects.append(data)
            except Exception:
                continue
        root_proj = workspace / ".ts2_project.json"
        if root_proj.exists():
            try:
                data = _json_loads(root_proj.read_text(encoding="utf-8"))
                data["path"] = ""
                projects.insert(0, data)
            except Exception:
                pass
        return projects

    def _read_projects_json(workspace_dir: str):
        """读取 projects.json（集中式项目列表，用于同步）"""
        pj_path = Path(workspace_dir) / "projects.json"
        if not pj_path.exists():
            return []
        try:
            return _json_loads(pj_path.read_text(encoding="utf-8"))
        except Exception:
            return []

    def _read_courses_data(workspace_dir: str):
        """同步读取课程数据（多源合并，同 course_tracker.py 逻辑）"""
        excluded = {"courses_structured_progress.json", "resource_index.json",
                     "courses_structured_progress.bak.json"}
        scan_dirs = [Path(workspace_dir), Path.home() / ".ts2", Path(workspace_dir).parent]
        all_courses = []
        seen_keys = set()
        for scan_dir in scan_dirs:
            if not scan_dir.exists():
                continue
            for f in sorted(scan_dir.glob("*.json")):
                if f.name in excluded or "_progress" in f.name:
                    continue
                try:
                    content = f.read_text(encoding="utf-8")
                    data = _json_loads(content)
                except Exception:
                    continue
                if not isinstance(data, dict):
                    continue
                courses = data.get("courses")
                if not isinstance(courses, list) or not courses:
                    continue
                for course in courses:
                    cid = course.get("note_id", "")
                    title = (course.get("course_title", "") or "").strip()
                    key = cid or title
                    if key and key in seen_keys:
                        continue
                    if key:
                        seen_keys.add(key)
                    if isinstance(course.get("lessons"), list):
                        for lesson in course["lessons"]:
                            if "estimated_hours" not in lesson or lesson["estimated_hours"] is None:
                                lesson["estimated_hours"] = 1.0
                            elif not isinstance(lesson["estimated_hours"], (int, float)):
                                try:
                                    lesson["estimated_hours"] = float(lesson["estimated_hours"])
                                except (ValueError, TypeError):
                                    lesson["estimated_hours"] = 1.0
                    all_courses.append(course)
        return {"courses": all_courses, "metadata": {"source_count": len(scan_dirs)}}

    def _read_agent_status(workspace_dir: str):
        """同步读取 Agent 状态"""
        agent = _get_web_agent(workspace_dir)
        if agent is None:
            return {"available": False, "tools": 0}
        return {
            "available": True,
            "tools": len(agent.tools),
            "model": agent.config.model_id or "unknown",
            "messages": len(agent.messages),
        }

    def _read_projects_dir(workspace_dir: str, req_path: str):
        """同步读取项目目录内容"""
        workspace = Path(workspace_dir)
        target = workspace / req_path if req_path else workspace

        try:
            target.resolve().relative_to(workspace.resolve())
        except ValueError:
            return None, "Access denied: path outside workspace"

        if not target.exists() or not target.is_dir():
            return None, "Directory not found"

        entries = []
        try:
            for item in sorted(target.iterdir()):
                if item.name.startswith(".") or item.name == "__pycache__":
                    continue
                rel_path = str(item.relative_to(workspace)).replace("\\", "/")
                entry = {
                    "path": rel_path,
                    "name": item.name,
                    "is_dir": item.is_dir(),
                    "ext": item.suffix.lower() if item.suffix else "",
                }
                if not item.is_dir():
                    try:
                        stat = item.stat()
                        entry["size"] = stat.st_size
                        entry["modified"] = stat.st_mtime
                    except (OSError, PermissionError):
                        continue
                entries.append(entry)
        except PermissionError:
            return None, "Permission denied"

        return entries, None

    def _read_project_file(workspace_dir: str, req_path: str):
        """同步读取项目源代码文件内容"""
        workspace = Path(workspace_dir)
        target = workspace / req_path

        try:
            target.resolve().relative_to(workspace.resolve())
        except ValueError:
            return None, "Access denied: path outside workspace"

        if not target.exists() or not target.is_file():
            return None, "File not found"

        try:
            size = target.stat().st_size
            if size > 2 * 1024 * 1024:
                return None, "File too large (max 2MB)"
        except OSError:
            return None, "Cannot read file stats"

        try:
            content = target.read_text(encoding="utf-8")
            is_binary = False
        except UnicodeDecodeError:
            import base64
            raw = target.read_bytes()
            content = base64.b64encode(raw).decode("ascii")
            is_binary = True

        stat = target.stat()
        ext = target.suffix.lower()

        return {
            "path": req_path,
            "name": target.name,
            "content": content,
            "is_binary": is_binary,
            "size": stat.st_size,
            "ext": ext,
            "modified": stat.st_mtime,
        }, None

    def _write_project_file(workspace_dir: str, req_path: str, req_content: str):
        """同步写入项目源代码文件"""
        workspace = Path(workspace_dir)
        target = workspace / req_path

        try:
            target.resolve().relative_to(workspace.resolve())
        except ValueError:
            return None, "Access denied: path outside workspace"

        if ".." in req_path:
            return None, "Access denied: invalid path"

        target.parent.mkdir(parents=True, exist_ok=True)

        try:
            target.write_text(req_content, encoding="utf-8")
            stat = target.stat()
            return {
                "path": req_path,
                "name": target.name,
                "size": stat.st_size,
                "modified": stat.st_mtime,
            }, None
        except (OSError, PermissionError) as e:
            return None, f"Write failed: {e}"

    def _update_task_data(workspace_dir: str, req_id: str, updates: dict):
        """同步更新任务"""
        tb_path = Path(workspace_dir) / "task_board.json"
        if not tb_path.exists():
            return None, "task_board.json not found"
        try:
            content = tb_path.read_text(encoding="utf-8")
            tasks = _json_loads(content)
        except Exception as e:
            return None, f"Read failed: {e}"

        updated = None
        for task in tasks:
            if task.get("id") == req_id:
                for key, value in updates.items():
                    if value is not None:
                        task[key] = value
                updated = task
                break

        if updated is None:
            return None, f"Task not found: {req_id}"

        try:
            tb_path.write_text(_json_dumps(tasks, ensure_ascii=False, indent=2), encoding="utf-8")
        except Exception as e:
            return None, f"Write failed: {e}"

        return updated, None

    def _create_task_data(workspace_dir: str, new_task: dict):
        """同步创建任务"""
        tb_path = Path(workspace_dir) / "task_board.json"
        tasks = []
        if tb_path.exists():
            try:
                content = tb_path.read_text(encoding="utf-8")
                tasks = _json_loads(content)
            except Exception:
                tasks = []

        tasks.append(new_task)

        try:
            tb_path.write_text(_json_dumps(tasks, ensure_ascii=False, indent=2), encoding="utf-8")
        except Exception as e:
            return None, f"Write failed: {e}"

        return new_task, None

    def _delete_task_data(workspace_dir: str, req_id: str):
        """同步删除任务"""
        tb_path = Path(workspace_dir) / "task_board.json"
        if not tb_path.exists():
            return None, "task_board.json not found"
        try:
            content = tb_path.read_text(encoding="utf-8")
            tasks = _json_loads(content)
        except Exception as e:
            return None, f"Read failed: {e}"

        original_len = len(tasks)
        tasks = [t for t in tasks if t.get("id") != req_id]
        if len(tasks) == original_len:
            return None, f"Task not found: {req_id}"

        try:
            tb_path.write_text(_json_dumps(tasks, ensure_ascii=False, indent=2), encoding="utf-8")
        except Exception as e:
            return None, f"Write failed: {e}"

        return True, None

    def _read_course_progress(workspace_dir: str, course_id: str):
        """同步读取课程进度"""
        progress_path = Path(workspace_dir) / "data" / "progress" / f"{course_id}.json"
        if not progress_path.exists():
            return {}
        try:
            content = progress_path.read_text(encoding="utf-8")
            return _json_loads(content)
        except Exception:
            return {}

    def _update_lesson_status(workspace_dir: str, course_id: str, lesson_number: int, status: str):
        """同步更新课时状态"""
        progress_dir = Path(workspace_dir) / "data" / "progress"
        progress_path = progress_dir / f"{course_id}.json"

        progress_data = {}
        if progress_path.exists():
            try:
                content = progress_path.read_text(encoding="utf-8")
                progress_data = _json_loads(content)
            except Exception:
                progress_data = {}

        lessons = progress_data.get("lessons", {})
        lessons[str(lesson_number)] = status
        progress_data["lessons"] = lessons

        try:
            progress_dir.mkdir(parents=True, exist_ok=True)
            progress_path.write_text(_json_dumps(progress_data, ensure_ascii=False, indent=2), encoding="utf-8")
        except Exception as e:
            return None, f"Write progress file failed: {e}"

        return progress_data, None

    def _update_lesson_data(workspace_dir: str, course_id: str, lesson_number: int, updates: dict):
        """更新课时信息（lesson_title, central_question, description, estimated_hours 等）"""
        scan_dirs = [Path(workspace_dir), Path.home() / ".ts2", Path(workspace_dir).parent]
        for scan_dir in scan_dirs:
            if not scan_dir.exists():
                continue
            for f in sorted(scan_dir.glob("*.json")):
                if "_progress" in f.name or f.name in {"courses_structured_progress.json",
                                                        "resource_index.json",
                                                        "courses_structured_progress.bak.json"}:
                    continue
                try:
                    content = f.read_text(encoding="utf-8")
                    data = _json_loads(content)
                except Exception:
                    continue
                if not isinstance(data, dict):
                    continue
                courses = data.get("courses")
                if not isinstance(courses, list):
                    continue
                for course in courses:
                    cid = course.get("note_id", "") or course.get("id", "")
                    if str(cid) != str(course_id):
                        continue
                    lessons = course.get("lessons", [])
                    for lesson in lessons:
                        ln = lesson.get("lesson_number", lesson.get("number", 0))
                        if ln == lesson_number:
                            lesson.update(updates)
                            f.write_text(_json_dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
                            return lesson, None
        return None, "Lesson not found"

    def _update_course_data(workspace_dir: str, course_id: str, updates: dict):
        """更新课程信息（course_title, domain, total_hours 等）"""
        scan_dirs = [Path(workspace_dir), Path.home() / ".ts2", Path(workspace_dir).parent]
        for scan_dir in scan_dirs:
            if not scan_dir.exists():
                continue
            for f in sorted(scan_dir.glob("*.json")):
                if "_progress" in f.name or f.name in {"courses_structured_progress.json",
                                                        "resource_index.json",
                                                        "courses_structured_progress.bak.json"}:
                    continue
                try:
                    content = f.read_text(encoding="utf-8")
                    data = _json_loads(content)
                except Exception:
                    continue
                if not isinstance(data, dict):
                    continue
                courses = data.get("courses")
                if not isinstance(courses, list):
                    continue
                for course in courses:
                    cid = course.get("note_id", "") or course.get("id", "")
                    if str(cid) != str(course_id):
                        continue
                    course.update(updates)
                    f.write_text(_json_dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
                    return course, None
        return None, "Course not found"

    # ── 数据 API 路由 ──

    @app.post("/api/data/tasks")
    async def data_tasks():
        """获取任务列表（优先从 task_board.json，回退到 automation.db）"""
        data = await _run_data(_read_tasks_data, app.state.workspace_dir)
        return ok(data=data)

    @app.post("/api/data/bookmarks")
    async def data_bookmarks():
        """获取书签数据"""
        data = await _run_data(_read_bookmarks_data, app.state.workspace_dir)
        return ok(data=data)

    @app.post("/api/data/bookmarks/add")
    async def data_bookmarks_add(req: BookmarkAddRequest):
        """添加书签"""
        bookmark = {
            "name": req.name,
            "url": req.url,
            "category": req.category,
            "icon": req.icon or "🔖",
            "color": "#3498db",
            "children": [],
        }
        result, error = await _run_data(_add_bookmark_data, app.state.workspace_dir, bookmark)
        if error:
            return err(msg=error)
        return ok(data=result)

    @app.post("/api/data/bookmarks/delete")
    async def data_bookmarks_delete(req: BookmarkDeleteRequest):
        """删除书签"""
        result, error = await _run_data(_delete_bookmark_data, app.state.workspace_dir, req.id)
        if error:
            return err(msg=error)
        return ok(data={"id": req.id})

    @app.post("/api/data/projects")
    async def data_projects(request: Request):
        """获取项目列表"""
        if not _check_source_auth(request):
            return err(403, "源码浏览器需要单独授权，请输入源码授权码")
        data = await _run_data(_read_projects_data, app.state.workspace_dir)
        return ok(data=data)

    @app.post("/api/data/projects/readDir")
    async def data_projects_read_dir(req: DirReadRequest, request: Request):
        """读取项目目录内容（不受 EXPOSED_DIRS 限制）"""
        if not _check_source_auth(request):
            return err(403, "源码浏览器需要单独授权，请输入源码授权码")
        result, error = await _run_data(_read_projects_dir, app.state.workspace_dir, req.path)
        if error:
            if "Access denied" in error:
                return err(code=403, msg=error)
            if "not found" in error.lower():
                return err(code=404, msg=error)
            if "Permission" in error:
                return err(code=403, msg=error)
            return err(msg=error)
        return ok(data=result)

    @app.post("/api/data/projects/readFile")
    async def data_projects_read_file(req: FileReadRequest, request: Request):
        """读取项目源代码文件内容（不受 EXPOSED_DIRS 限制）"""
        if not _check_source_auth(request):
            return err(403, "源码浏览器需要单独授权，请输入源码授权码")
        result, error = await _run_file(_read_project_file, app.state.workspace_dir, req.path)
        if error:
            if "Access denied" in error:
                return err(code=403, msg=error)
            if "not found" in error.lower():
                return err(code=404, msg=error)
            if "too large" in error.lower():
                return err(code=413, msg=error)
            return err(msg=error)
        return ok(data=result)

    @app.post("/api/data/projects/writeFile")
    async def data_projects_write_file(req: FileWriteRequest, request: Request):
        """写入项目源代码文件（不受 EXPOSED_DIRS 限制）"""
        if not _check_source_auth(request, require_write=True):
            return err(403, "源码写入需要源码授权码")
        result, error = await _run_file(_write_project_file, app.state.workspace_dir, req.path, req.content)
        if error:
            if "Access denied" in error:
                return err(code=403, msg=error)
            if "Write failed" in error:
                return err(code=500, msg=error)
            return err(msg=error)
        return ok(data=result)

    # ─── RMD 编译工具 API ──────────────────────────────────────

    class KnitRmdRequest(BaseModel):
        path: str
        output_format: Optional[str] = "html_document"

    @app.post("/api/tools/knit-rmd")
    async def tools_knit_rmd(req: KnitRmdRequest):
        """编译 Rmd/Markdown 文件（参考 course_tracker._add_rmd_compile_buttons 方法）"""
        import subprocess
        workspace = Path(app.state.workspace_dir)
        rmd_path = (workspace / req.path).resolve()
        if not rmd_path.exists():
            return err(msg=f"文件不存在: {req.path}")
        ext = rmd_path.suffix.lower()
        if ext not in (".rmd", ".rmarkdown", ".md"):
            return err(msg=f"不支持的文件格式: {ext}")

        output_format = req.output_format or "html_document"
        # 使用 posix 路径，与 course_tracker 一致
        rmd_posix = rmd_path.as_posix()

        try:
            result = await asyncio.get_event_loop().run_in_executor(
                None,
                lambda: subprocess.run(
                    ["Rscript", "-e", f"rmarkdown::render('{rmd_posix}', output_format='{output_format}')"],
                    cwd=str(rmd_path.parent),
                    capture_output=True, text=True, encoding='utf-8', errors='replace', timeout=300
                )
            )
            if result.returncode == 0:
                output_ext = {"pdf_document": ".pdf", "html_document": ".html", "word_document": ".docx",
                              "ioslides_presentation": ".html", "slidy_presentation": ".html",
                              "beamer_presentation": ".pdf"}.get(output_format, ".html")
                output_path = rmd_path.with_suffix(output_ext)
                rel_output = str(output_path.relative_to(workspace)) if output_path.exists() else None
                return ok(data={
                    "success": True,
                    "output_file": rel_output,
                    "stdout": result.stdout,
                    "stderr": result.stderr
                })
            else:
                err_msg = (result.stderr or result.stdout or "未知错误")[:500]
                return ok(data={"success": False, "error": err_msg, "stdout": result.stdout, "stderr": result.stderr})
        except subprocess.TimeoutExpired:
            return ok(data={"success": False, "error": "编译超时 (超过5分钟)"})
        except Exception as e:
            return err(msg=f"编译失败: {e}")

    @app.get("/api/tools/rmd-files")
    async def tools_list_rmd_files():
        """列出所有 Rmd 文件"""
        try:
            from ..ws2_tools import NotesManager
            nm = NotesManager(Path(app.state.workspace_dir))
            if not nm.ready:
                return ok(data={"files": [], "courses": []})
            courses = nm.list_courses()
            all_files = []
            for course in courses:
                rmds = nm.list_rmd_files(course)
                for f in rmds:
                    all_files.append({"course": course, "filename": f, "path": f"notes/{course}/{f}"})
            return ok(data={"files": all_files, "courses": courses})
        except Exception as e:
            return ok(data={"files": [], "courses": [], "error": str(e)})

    # ─── Bilibili 代理（浏览器端绕过 CORS） ─────────────────────

    class BiliProxyRequest(BaseModel):
        url: str
        method: str = "GET"
        headers: Dict[str, str] = {}
        body: Optional[str] = None

    @app.post("/api/extractor/biliProxy")
    @app.get("/api/extractor/biliProxy")
    async def bili_proxy(req: Optional[BiliProxyRequest] = None, url: str = "", method: str = "GET"):
        """代理 Bilibili API 请求，解决浏览器 CORS 限制"""
        import httpx
        target = req.url if req else url
        if not target:
            return err(msg="Missing 'url' parameter")
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/130.0.0.0 Safari/537.36",
            "Referer": "https://www.bilibili.com",
            "Origin": "https://www.bilibili.com",
        }
        if req and req.headers:
            headers.update(req.headers)
        http_method = (req.method if req else method).upper()
        async with httpx.AsyncClient(timeout=30.0) as client:
            try:
                if http_method == "POST":
                    resp = await client.post(target, headers=headers, content=req.body if req else None)
                else:
                    resp = await client.get(target, headers=headers)
                try:
                    return JSONResponse(content=resp.json())
                except Exception:
                    return JSONResponse(content={"code": -1, "msg": f"Non-JSON response: {resp.text[:500]}", "data": None})
            except Exception as e:
                return err(msg=f"Bilibili proxy failed: {e}")

    @app.get("/api/extractor/imageProxy")
    async def image_proxy(url: str = ""):
        """代理 Bilibili 图片，解决浏览器 CORS/Referer 限制"""
        if not url:
            return JSONResponse(content={"error": "Missing url"}, status_code=400)
        import httpx
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/130.0.0.0 Safari/537.36",
            "Referer": "https://www.bilibili.com",
        }
        async with httpx.AsyncClient(timeout=15.0) as client:
            try:
                resp = await client.get(url, headers=headers)
                content_type = resp.headers.get("content-type", "image/webp")
                return Response(content=resp.content, media_type=content_type)
            except Exception as e:
                return JSONResponse(content={"error": str(e)}, status_code=502)

    @app.get("/api/browser/proxy")
    async def browser_proxy(url: str = ""):
        """代理网页浏览，绕过 X-Frame-Options 限制"""
        if not url:
            return HTMLResponse(content="<html><body>错误：缺少 url 参数</body></html>", status_code=200)
        import httpx
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/130.0.0.0 Safari/537.36",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
            "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.8",
        }
        async with httpx.AsyncClient(timeout=60.0, follow_redirects=True, verify=False) as client:
            try:
                resp = await client.get(url, headers=headers)
                ct = resp.headers.get("content-type", "").lower()
                if "text/html" in ct or "application/xhtml" in ct:
                    body = resp.text
                    base_tag = f'<base href="{url}">'
                    # 移除常见 frame-busting 脚本
                    body = re.sub(
                        r'<script[^>]*>[\s\S]*?(?:'
                        r'window\.top\s*(?:!==?|!=)\s*window|'
                        r'top\s*(?:!==?|!=)\s*self|'
                        r'(?:self|window)\s*(?:!==?|!=)\s*top|'
                        r'top\.location\s*[=.]|'
                        r'parent\.location\s*[=.]|'
                        r'window\.top\.location\s*[=.]|'
                        r'top\[[\'"]location[\'"]\]\s*[=.]|'
                        r'parent\[[\'"]location[\'"]\]\s*[=.]'
                        r')[\s\S]*?</script>',
                        '<script>/*fb-stripped*/</script>',
                        body,
                        flags=re.IGNORECASE,
                        count=30
                    )
                    proxy_inject = (
                        '<script>'
                        '(function(){'
                        'var P="/api/browser/proxy?url=";'
                        'document.addEventListener("click",function(e){'
                        'var a=e.target.closest("a");'
                        'if(a&&a.href&&!a.href.startsWith(P)&&a.target!=="_blank"){'
                        'e.preventDefault();'
                        'window.location.href=P+encodeURIComponent(a.href);'
                        '}'
                        '});'
                        'document.addEventListener("submit",function(e){'
                        'var f=e.target;'
                        'if(f&&f.action&&!f.action.startsWith(P)){'
                        'e.preventDefault();'
                        'var fd=new FormData(f);var q=[];'
                        'fd.forEach(function(v,k){q.push(encodeURIComponent(k)+"="+encodeURIComponent(v));});'
                        'var sep=f.action.includes("?")?"&":"?";'
                        'window.location.href=P+encodeURIComponent(f.action+sep+q.join("&"));'
                        '}'
                        '});'
                        '})();'
                        '</script>'
                    )
                    if "<head>" in body:
                        body = body.replace("<head>", f"<head>{base_tag}{proxy_inject}", 1)
                    elif "<html>" in body:
                        body = body.replace("<html>", f"<html><head>{base_tag}{proxy_inject}</head>", 1)
                    else:
                        body = f"<head>{base_tag}{proxy_inject}</head>{body}"
                    return HTMLResponse(content=body, status_code=200)
                else:
                    return Response(content=resp.content, media_type=ct or "application/octet-stream", status_code=200)
            except Exception as e:
                err_html = f"<html><body>代理错误：{str(e)}</body></html>"
                return HTMLResponse(content=err_html, status_code=200)

    @app.post("/api/data/courses")
    async def data_courses():
        """获取课程数据"""
        data = await _run_data(_read_courses_data, app.state.workspace_dir)
        return ok(data=data)

    # ─── 任务 CRUD API ──────────────────────────────────────

    @app.post("/api/data/tasks/update")
    async def task_update(req: TaskUpdateRequest):
        """更新任务（task_board.json）"""
        updates = {
            "title": req.title,
            "description": req.description,
            "due_date": req.due_date,
            "priority": req.priority,
            "status": req.status,
            "start_time": req.start_time,
            "duration": req.duration,
            "recurrence": req.recurrence,
        }
        updated, error = await _run_data(_update_task_data, app.state.workspace_dir, req.id, updates)
        if error:
            if "not found" in error.lower() and "task_board" not in error:
                return err(code=404, msg=error)
            return err(msg=error)
        return ok(data=updated)

    @app.post("/api/data/tasks/create")
    async def task_create(req: TaskCreateRequest):
        """创建任务（task_board.json）"""
        new_task = {
            "id": str(uuid.uuid4()),
            "title": req.title,
            "description": req.description,
            "due_date": req.due_date,
            "priority": req.priority,
            "status": req.status,
            "start_time": req.start_time,
            "duration": req.duration,
            "recurrence": req.recurrence,
        }
        result, error = await _run_data(_create_task_data, app.state.workspace_dir, new_task)
        if error:
            return err(msg=error)
        return ok(data=result)

    @app.post("/api/data/tasks/delete")
    async def task_delete(req: TaskDeleteRequest):
        """删除任务（task_board.json）"""
        result, error = await _run_data(_delete_task_data, app.state.workspace_dir, req.id)
        if error:
            if "not found" in error.lower() and "task_board" not in error:
                return err(code=404, msg=error)
            return err(msg=error)
        return ok()

    # ─── 课程进度 API ──────────────────────────────────────

    @app.post("/api/data/courses/progress")
    async def course_progress(req: CourseProgressRequest):
        """获取课程进度（data/progress/{course_id}.json）"""
        data = await _run_data(_read_course_progress, app.state.workspace_dir, req.course_id)
        return ok(data=data)

    @app.post("/api/data/courses/lessonStatus")
    async def lesson_status(req: LessonStatusRequest):
        """更新课时状态（data/progress/{course_id}.json）"""
        result, error = await _run_data(_update_lesson_status, app.state.workspace_dir, req.course_id, req.lesson_number, req.status)
        if error:
            return err(msg=error)
        return ok(data=result)

    @app.post("/api/data/courses/updateLesson")
    async def update_lesson(req: UpdateLessonRequest):
        """更新课时信息（标题、中心问题、描述、预计时长等）"""
        result, error = await _run_data(_update_lesson_data, app.state.workspace_dir, req.course_id, req.lesson_number, req.updates)
        if error:
            return err(msg=error)
        return ok(data=result)

    @app.post("/api/data/courses/update")
    async def update_course(req: UpdateCourseRequest):
        """更新课程信息（标题、域、总学时等）"""
        result, error = await _run_data(_update_course_data, app.state.workspace_dir, req.course_id, req.updates)
        if error:
            return err(msg=error)
        return ok(data=result)

    # ─── 课程表 API ──────────────────────────────────────────

    @app.get("/api/data/timetable")
    async def timetable_get():
        """获取所有课程表"""
        from ..automation.course_simulation import TimetablePersistence
        persistence = TimetablePersistence()
        timetables = persistence.load()
        result = {}
        for tid, tt in timetables.items():
            result[tid] = tt.to_dict()
        return ok(data=result)

    @app.post("/api/data/timetable/create")
    async def timetable_create(req: TimetableCreateRequest):
        """创建课程表"""
        from ..automation.course_simulation import Timetable, TimetablePersistence
        import uuid as _uuid
        persistence = TimetablePersistence()
        timetables = persistence.load()
        tt_id = f"tt_{_uuid.uuid4().hex[:8]}"
        new_tt = Timetable(
            timetable_id=tt_id,
            name=req.name,
            semester_start=req.semester_start,
            semester_end=req.semester_end,
            slots=[],
            enabled=len(timetables) == 0,  # 第一个自动启用
        )
        timetables[tt_id] = new_tt
        persistence.save(timetables)
        return ok(data=new_tt.to_dict())

    @app.post("/api/data/timetable/setActive")
    async def timetable_set_active(req: TimetableSetActiveRequest):
        """设置激活课程表"""
        from ..automation.course_simulation import TimetablePersistence
        persistence = TimetablePersistence()
        timetables = persistence.load()
        if req.timetable_id not in timetables:
            return err(msg="课程表不存在")
        for tid, tt in timetables.items():
            tt.enabled = (tid == req.timetable_id)
        persistence.save(timetables)
        return ok(data=timetables[req.timetable_id].to_dict())

    @app.post("/api/data/timetable/delete")
    async def timetable_delete(req: TimetableDeleteRequest):
        """删除课程表"""
        from ..automation.course_simulation import TimetablePersistence
        persistence = TimetablePersistence()
        timetables = persistence.load()
        if req.timetable_id not in timetables:
            return err(msg="课程表不存在")
        del timetables[req.timetable_id]
        persistence.save(timetables)
        return ok()

    @app.post("/api/data/timetable/slot/add")
    async def timetable_slot_add(req: TimetableSlotCreateRequest):
        """添加课时到课程表"""
        from ..automation.course_simulation import TimetableSlot, TimetablePersistence, COURSE_COLORS
        persistence = TimetablePersistence()
        timetables = persistence.load()
        # 找到目标课程表
        tt = None
        if req.timetable_id and req.timetable_id in timetables:
            tt = timetables[req.timetable_id]
        else:
            # 默认使用激活的课程表
            for t in timetables.values():
                if t.enabled:
                    tt = t
                    break
            if not tt and timetables:
                tt = next(iter(timetables.values()))
        if not tt:
            return err(msg="没有可用的课程表")
        # 自动分配颜色
        color = req.color
        if not color:
            existing_colors = {s.color for s in tt.slots if s.color}
            for c in COURSE_COLORS:
                if c not in existing_colors:
                    color = c
                    break
            if not color:
                color = COURSE_COLORS[len(tt.slots) % len(COURSE_COLORS)]
        slot = TimetableSlot(
            slot_id=f"slot_{uuid.uuid4().hex[:8]}",
            course_id=req.course_name,
            course_name=req.course_name,
            day_of_week=req.day_of_week,
            start_time=req.start_time,
            end_time=req.end_time,
            location=req.location,
            teacher=req.teacher,
            period_idx=req.period_idx,
            color=color,
        )
        tt.slots.append(slot)
        persistence.save(timetables)
        return ok(data=slot.to_dict())

    @app.post("/api/data/timetable/slot/delete")
    async def timetable_slot_delete(req: TimetableSlotDeleteRequest):
        """删除课时"""
        from ..automation.course_simulation import TimetablePersistence
        persistence = TimetablePersistence()
        timetables = persistence.load()
        # 找到包含该 slot 的课程表
        for tt in timetables.values():
            before = len(tt.slots)
            tt.slots = [s for s in tt.slots if s.slot_id != req.slot_id]
            if len(tt.slots) < before:
                persistence.save(timetables)
                return ok()
        return err(msg="课时不存在")

    # ─── 关键路径检测 API ─────────────────────────────────────────

    @app.get("/api/tasks/critical-path")
    async def tasks_critical_path():
        """
        关键路径检测（Critical Path Method）
        基于任务的依赖关系和持续时间，计算关键路径
        返回：关键路径上的任务序列、项目最短完成时间、各任务的时间裕度
        """
        workspace_dir = app.state.workspace_dir
        tasks = _read_tasks_data(workspace_dir) or []

        if not tasks:
            return ok(data={"critical_path": [], "project_duration": 0, "tasks_schedule": []})

        # 构建任务图
        task_map = {}
        for t in tasks:
            tid = t.get("id")
            if not tid:
                continue
            task_map[tid] = {
                "id": tid,
                "title": t.get("title", ""),
                "status": t.get("status", ""),
                "priority": t.get("priority", ""),
                "due_date": t.get("due_date", ""),
                "duration": t.get("duration", 60),  # 默认 60 分钟
                "dependencies": t.get("dependencies", []),  # 前置任务 ID 列表
                "updated_at": t.get("updated_at", ""),
            }

        # 拓扑排序 + 前向传播（计算最早开始/结束时间）
        from collections import defaultdict, deque

        in_degree = defaultdict(int)
        graph = defaultdict(list)  # adjacency: dep -> [successors]
        for tid, t in task_map.items():
            deps = t["dependencies"]
            if not isinstance(deps, list):
                deps = []
                t["dependencies"] = deps
            in_degree.setdefault(tid, 0)
            for dep_id in deps:
                if dep_id in task_map:
                    graph[dep_id].append(tid)
                    in_degree[tid] = in_degree.get(tid, 0) + 1

        # 拓扑排序
        queue = deque([tid for tid in task_map if in_degree.get(tid, 0) == 0])
        topo_order = []
        while queue:
            tid = queue.popleft()
            topo_order.append(tid)
            for succ in graph.get(tid, []):
                in_degree[succ] -= 1
                if in_degree[succ] == 0:
                    queue.append(succ)

        # 如果有环，跳过有环的任务
        if len(topo_order) != len(task_map):
            # 移除不在拓扑序中的任务（有环）
            valid_ids = set(topo_order)
        else:
            valid_ids = set(task_map.keys())

        # 前向传播：计算最早开始时间(ES)和最早结束时间(EF)
        es = {}  # Earliest Start
        ef = {}  # Earliest Finish
        for tid in topo_order:
            t = task_map[tid]
            deps = [d for d in t["dependencies"] if d in valid_ids]
            if not deps:
                es[tid] = 0
            else:
                es[tid] = max(ef.get(d, 0) for d in deps)
            ef[tid] = es[tid] + t["duration"]

        # 项目总工期
        project_duration = max(ef.values()) if ef else 0

        # 后向传播：计算最迟开始时间(LS)和最迟结束时间(LF)
        ls = {}  # Latest Start
        lf = {}  # Latest Finish
        for tid in reversed(topo_order):
            t = task_map[tid]
            successors = [s for s in graph.get(tid, []) if s in valid_ids]
            if not successors:
                lf[tid] = project_duration
            else:
                lf[tid] = min(ls.get(s, project_duration) for s in successors)
            ls[tid] = lf[tid] - t["duration"]

        # 计算时间裕度（Total Float）
        schedule = []
        critical_path = []
        for tid in topo_order:
            t = task_map[tid]
            total_float = ls.get(tid, 0) - es.get(tid, 0)
            is_critical = total_float == 0 and t["status"] != "已完成"
            entry = {
                "id": tid,
                "title": t["title"],
                "status": t["status"],
                "priority": t["priority"],
                "due_date": t["due_date"],
                "duration": t["duration"],
                "dependencies": t["dependencies"],
                "earliest_start": es.get(tid, 0),
                "earliest_finish": ef.get(tid, 0),
                "latest_start": ls.get(tid, 0),
                "latest_finish": lf.get(tid, 0),
                "total_float": total_float,
                "is_critical": is_critical,
            }
            schedule.append(entry)
            if is_critical:
                critical_path.append(entry)

        return ok(data={
            "critical_path": critical_path,
            "project_duration": project_duration,
            "tasks_schedule": schedule,
            "total_tasks": len(task_map),
            "critical_tasks": len(critical_path),
        })

    # ─── 分页笔记 API ──────────────────────────────────────

    def _read_notebooks(workspace_dir: str):
        """读取所有笔记本列表"""
        nb_dir = Path(workspace_dir) / "data" / "notebooks"
        if not nb_dir.exists():
            return []
        result = []
        for f in sorted(nb_dir.glob("*.json")):
            try:
                nb = json.loads(f.read_text(encoding="utf-8"))
                nb["_file"] = f.name
                result.append(nb)
            except:
                pass
        return result

    def _read_notebook(workspace_dir: str, notebook_id: str):
        """读取单个笔记本"""
        nb_path = Path(workspace_dir) / "data" / "notebooks" / f"{notebook_id}.json"
        if not nb_path.exists():
            return None
        return json.loads(nb_path.read_text(encoding="utf-8"))

    def _write_notebook(workspace_dir: str, notebook_id: str, data: dict):
        """写入笔记本"""
        nb_dir = Path(workspace_dir) / "data" / "notebooks"
        nb_dir.mkdir(parents=True, exist_ok=True)
        nb_path = nb_dir / f"{notebook_id}.json"
        data["updatedAt"] = int(time.time() * 1000)
        nb_path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
        return data

    def _delete_notebook(workspace_dir: str, notebook_id: str):
        """删除笔记本"""
        nb_path = Path(workspace_dir) / "data" / "notebooks" / f"{notebook_id}.json"
        if nb_path.exists():
            nb_path.unlink()
            return True
        return False

    @app.get("/api/notebooks")
    async def list_notebooks():
        data = await _run_data(_read_notebooks, app.state.workspace_dir)
        return ok(data=data)

    @app.get("/api/notebooks/{notebook_id}")
    async def get_notebook(notebook_id: str):
        data = await _run_data(_read_notebook, app.state.workspace_dir, notebook_id)
        if data is None:
            return ok(data=None, msg="笔记本不存在")
        return ok(data=data)

    @app.post("/api/notebooks/{notebook_id}")
    async def save_notebook(notebook_id: str, req: dict):
        data = await _run_data(_write_notebook, app.state.workspace_dir, notebook_id, req)
        return ok(data=data)

    @app.delete("/api/notebooks/{notebook_id}")
    async def del_notebook(notebook_id: str):
        result = await _run_data(_delete_notebook, app.state.workspace_dir, notebook_id)
        return ok(data={"deleted": result})

    # ─── 课程系统增强 API（移植自 CourseSystem/ResourceMgr/WorkflowLogger）──

    # ─── 增量同步 API ──────────────────────────────────────

    def _sync_items(local_items: list, server_items: list, id_field: str = "id") -> dict:
        """
        通用增量去重同步算法：
        - 按 id 匹配，按 updated_at 判断新旧
        - 返回：pull（服务器更新的，客户端需要拉取）、push（客户端更新的，服务器需要接受）、conflicts（双方都修改了）
        - 本地有但服务器没有 → push（新增）
        - 服务器有但本地没有 → pull（新增）
        - 都有，服务器更新 → pull
        - 都有，本地更新 → push
        - 都有，同时更新 → conflict
        """
        from datetime import datetime as _dt

        def _parse_time(t) -> float:
            if not t:
                return 0.0
            if isinstance(t, (int, float)):
                return float(t)
            try:
                return _dt.strptime(str(t), "%Y-%m-%d %H:%M:%S").timestamp()
            except Exception:
                try:
                    return _dt.strptime(str(t), "%Y-%m-%dT%H:%M:%S").timestamp()
                except Exception:
                    return 0.0

        local_map = {item.get(id_field): item for item in local_items if item.get(id_field)}
        server_map = {item.get(id_field): item for item in server_items if item.get(id_field)}

        pull = []   # 客户端需要拉取的（服务器更新或新增）
        push = []   # 服务器需要接受的（客户端更新或新增）
        conflicts = []  # 冲突的

        # 服务器有但本地没有 → pull
        for sid, sitem in server_map.items():
            if sid not in local_map:
                pull.append(sitem)

        # 本地有但服务器没有 → push
        for lid, litem in local_map.items():
            if lid not in server_map:
                push.append(litem)

        # 都有 → 比较 updated_at
        for lid in local_map:
            if lid not in server_map:
                continue
            litem = local_map[lid]
            sitem = server_map[lid]
            l_time = _parse_time(litem.get("updated_at"))
            s_time = _parse_time(sitem.get("updated_at"))

            if l_time == s_time:
                continue  # 无变化
            elif s_time > l_time:
                pull.append(sitem)  # 服务器更新
            elif l_time > s_time:
                push.append(litem)  # 客户端更新
            else:
                # 时间相同但内容不同
                conflicts.append({"local": litem, "server": sitem})

        return {"pull": pull, "push": push, "conflicts": conflicts}

    def _apply_sync_push(workspace_dir: str, data_type: str, items: list) -> dict:
        """应用客户端推送的变更到服务器"""
        if data_type == "tasks":
            tb_path = Path(workspace_dir) / "task_board.json"
            if not tb_path.exists():
                return None, "task_board.json not found"
            try:
                tasks = _json_loads(tb_path.read_text(encoding="utf-8"))
            except Exception as e:
                return None, f"Read failed: {e}"

            task_map = {t.get("id"): t for t in tasks}
            updated_count = 0
            created_count = 0
            for item in items:
                item_id = item.get("id")
                if item_id in task_map:
                    # 更新
                    task_map[item_id].update(item)
                    updated_count += 1
                else:
                    # 新增
                    tasks.append(item)
                    created_count += 1

            tb_path.write_text(_json_dumps(tasks, ensure_ascii=False, indent=2), encoding="utf-8")
            return {"updated": updated_count, "created": created_count}, None

        elif data_type == "bookmarks":
            bm_path = Path(workspace_dir) / "bookmarks.json"
            try:
                bookmarks = _json_loads(bm_path.read_text(encoding="utf-8")) if bm_path.exists() else []
            except Exception:
                bookmarks = []

            bm_map = {b.get("id"): b for b in bookmarks}
            updated_count = 0
            created_count = 0
            for item in items:
                item_id = item.get("id")
                if item_id in bm_map:
                    bm_map[item_id].update(item)
                    updated_count += 1
                else:
                    bookmarks.append(item)
                    created_count += 1

            bm_path.write_text(_json_dumps(bookmarks, ensure_ascii=False, indent=2), encoding="utf-8")
            return {"updated": updated_count, "created": created_count}, None

        elif data_type == "projects":
            pj_path = Path(workspace_dir) / "projects.json"
            try:
                projects = _json_loads(pj_path.read_text(encoding="utf-8")) if pj_path.exists() else []
            except Exception:
                projects = []

            pj_map = {p.get("id"): p for p in projects}
            updated_count = 0
            created_count = 0
            for item in items:
                item_id = item.get("id")
                if item_id in pj_map:
                    pj_map[item_id].update(item)
                    updated_count += 1
                else:
                    projects.append(item)
                    created_count += 1

            pj_path.write_text(_json_dumps(projects, ensure_ascii=False, indent=2), encoding="utf-8")
            return {"updated": updated_count, "created": created_count}, None

        return None, f"Unknown data type: {data_type}"

    @app.post("/api/sync/compare")
    async def sync_compare(req: Request):
        """
        增量同步对比：客户端发送本地数据，服务器返回差异
        body: { "tasks": [...], "bookmarks": [...] }
        返回: { "tasks": { "pull": [...], "push": [...], "conflicts": [...] }, "bookmarks": {...} }
        """
        body = await req.json()
        result = {}

        workspace_dir = app.state.workspace_dir

        for data_type in ["tasks", "bookmarks", "projects"]:
            local_items = body.get(data_type, [])
            if data_type == "tasks":
                server_items = _read_tasks_data(workspace_dir) or []
            elif data_type == "bookmarks":
                server_items = _read_bookmarks_data(workspace_dir) or []
            elif data_type == "projects":
                server_items = _read_projects_json(workspace_dir) or []
            else:
                continue

            result[data_type] = _sync_items(local_items, server_items)

        return ok(data=result)

    @app.post("/api/sync/push")
    async def sync_push(req: Request):
        """
        推送客户端变更到服务器
        body: { "tasks": [...], "bookmarks": [...], "projects": [...] }
        """
        body = await req.json()
        result = {}

        for data_type in ["tasks", "bookmarks", "projects"]:
            items = body.get(data_type, [])
            if not items:
                result[data_type] = {"updated": 0, "created": 0}
                continue
            applied, error = await _run_data(_apply_sync_push, app.state.workspace_dir, data_type, items)
            if error:
                result[data_type] = {"error": error}
            else:
                result[data_type] = applied

        return ok(data=result)

    @app.post("/api/sync/full")
    async def sync_full(req: Request):
        """
        完整同步：对比 + 自动推送 + 返回最终数据
        1. 对比差异
        2. 自动推送客户端更新的项
        3. 返回服务器最终数据（包含推送后的结果）
        """
        body = await req.json()
        workspace_dir = app.state.workspace_dir
        result = {}

        for data_type in ["tasks", "bookmarks", "projects"]:
            local_items = body.get(data_type, [])
            if data_type == "tasks":
                server_items = _read_tasks_data(workspace_dir) or []
            elif data_type == "bookmarks":
                server_items = _read_bookmarks_data(workspace_dir) or []
            elif data_type == "projects":
                server_items = _read_projects_json(workspace_dir) or []
            else:
                continue

            diff = _sync_items(local_items, server_items)

            # 自动推送客户端更新的项（push 列表）
            if diff["push"]:
                applied, error = _apply_sync_push(workspace_dir, data_type, diff["push"])
                if not error:
                    # 重新读取服务器数据
                    if data_type == "tasks":
                        server_items = _read_tasks_data(workspace_dir) or []
                    elif data_type == "bookmarks":
                        server_items = _read_bookmarks_data(workspace_dir) or []
                    elif data_type == "projects":
                        server_items = _read_projects_json(workspace_dir) or []

            result[data_type] = {
                "pull": diff["pull"],
                "conflicts": diff["conflicts"],
                "server_data": server_items,
                "pushed": len(diff["push"]),
            }

        return ok(data=result)

    # ─── 课程系统增强 API（移植自 CourseSystem/ResourceMgr/WorkflowLogger）──

    def _get_review_due(workspace_dir: str, course_id: str):
        """同步获取待复习课时列表"""
        progress_path = Path(workspace_dir) / "data" / "progress" / f"{course_id}.json"
        if not progress_path.exists():
            return []
        try:
            progress_data = _json_loads(progress_path.read_text(encoding="utf-8"))
            rs = progress_data.get("review_schedule", {})
            from datetime import datetime, timedelta
            now = datetime.now()
            due = []
            for ln_str, info in rs.items():
                try:
                    next_r = datetime.fromisoformat(info.get("next_review", ""))
                    if next_r <= now:
                        due.append({
                            "lesson_number": int(ln_str),
                            "workload": info.get("workload", 0),
                            "interval_days": info.get("interval_days", 7),
                            "review_count": info.get("review_count", 0),
                            "overdue_days": (now - next_r).days,
                        })
                except Exception:
                    pass
            due.sort(key=lambda x: -x["workload"])
            return due
        except Exception:
            return []

    def _mark_review_done(workspace_dir: str, course_id: str, lesson_number: int, status: str):
        """同步标记复习完成"""
        from datetime import datetime, timedelta
        progress_dir = Path(workspace_dir) / "data" / "progress"
        progress_path = progress_dir / f"{course_id}.json"
        progress_data = {}
        if progress_path.exists():
            try:
                progress_data = _json_loads(progress_path.read_text(encoding="utf-8"))
            except Exception:
                progress_data = {}

        rs = progress_data.setdefault("review_schedule", {})
        ln_str = str(lesson_number)
        if ln_str in rs:
            info = rs[ln_str]
            workload = info.get("workload", 5)
            review_count = info.get("review_count", 0) + 1
            info["review_count"] = review_count
            info["last_reviewed"] = datetime.now().isoformat()
            if workload <= 5: base_days = 7
            elif workload <= 15: base_days = 4
            elif workload <= 30: base_days = 2
            else: base_days = 1
            interval_days = max(1, int(base_days * (1.5 ** min(review_count, 5))))
            info["interval_days"] = interval_days
            info["next_review"] = (datetime.now() + timedelta(days=interval_days)).isoformat()
        else:
            rs[ln_str] = {
                "workload": 5, "interval_days": 7,
                "next_review": (datetime.now() + timedelta(days=7)).isoformat(),
                "last_reviewed": datetime.now().isoformat(), "review_count": 1,
            }

        try:
            progress_dir.mkdir(parents=True, exist_ok=True)
            progress_path.write_text(_json_dumps(progress_data, ensure_ascii=False, indent=2), encoding="utf-8")
        except Exception as e:
            return None, str(e)
        return progress_data, None

    def _update_review_schedule(workspace_dir: str, course_id: str, lesson_number: int, workload: int = 5):
        """同步更新复习调度（参考 course_tracker.py 的 update_review_schedule）"""
        from datetime import datetime, timedelta
        progress_dir = Path(workspace_dir) / "data" / "progress"
        progress_path = progress_dir / f"{course_id}.json"
        progress_data = {}
        if progress_path.exists():
            try:
                progress_data = _json_loads(progress_path.read_text(encoding="utf-8"))
            except Exception:
                progress_data = {}

        rs = progress_data.setdefault("review_schedule", {})
        ln_str = str(lesson_number)
        now = datetime.now()
        if workload <= 0: interval_days = 14
        elif workload <= 5: interval_days = 7
        elif workload <= 15: interval_days = 4
        elif workload <= 30: interval_days = 2
        else: interval_days = 1
        review_count = rs.get(ln_str, {}).get("review_count", 0)
        interval_days = max(1, int(interval_days * (1.5 ** min(review_count, 5))))
        rs[ln_str] = {
            "workload": workload, "interval_days": interval_days,
            "next_review": (now + timedelta(days=interval_days)).isoformat(),
            "last_reviewed": now.isoformat(), "review_count": review_count,
        }
        try:
            progress_dir.mkdir(parents=True, exist_ok=True)
            progress_path.write_text(_json_dumps(progress_data, ensure_ascii=False, indent=2), encoding="utf-8")
        except Exception as e:
            return None, str(e)
        return progress_data, None

    def _get_courses_stats(workspace_dir: str):
        """同步获取课程系统全局统计"""
        cs_path = Path(workspace_dir) / "courses_structured.json"
        if not cs_path.exists():
            return {}
        try:
            data = _json_loads(cs_path.read_text(encoding="utf-8"))
            courses = data.get("courses", [])
            from collections import Counter
            total_hours = sum(c.get("total_hours", 0) or 0 for c in courses)
            total_lessons = sum(len(c.get("lessons", [])) for c in courses)
            domain_dist = dict(Counter(c.get("domain", "UNKNOWN") for c in courses))
            progress_dir = Path(workspace_dir) / "data" / "progress"
            completed_lessons = 0
            for c in courses:
                cid = c.get("note_id", c.get("course_title", ""))
                pp = progress_dir / f"{cid}.json"
                if pp.exists():
                    try:
                        pd = _json_loads(pp.read_text(encoding="utf-8"))
                        lessons = pd.get("lessons", {})
                        completed_lessons += sum(1 for v in lessons.values() if v)
                    except Exception:
                        pass
            return {
                "total_courses": len(courses),
                "total_hours": total_hours,
                "total_lessons": total_lessons,
                "completed_lessons": completed_lessons,
                "domain_distribution": domain_dist,
                "completion_rate": round(completed_lessons / total_lessons * 100, 1) if total_lessons > 0 else 0,
            }
        except Exception:
            return {}

    def _create_course_data(workspace_dir: str, title: str, domain: str):
        """同步创建新课程"""
        import hashlib
        h = hashlib.md5(title.encode("utf-8")).hexdigest()[:16]
        new_course = {
            "note_id": f"new_{h}",
            "course_title": title,
            "total_hours": None,
            "domain": domain,
            "prerequisites": [],
            "positioning": "",
            "target_audience": "",
            "assessment": "",
            "sections": [],
            "lessons": [],
            "references": [],
        }
        cs_path = Path(workspace_dir) / "courses_structured.json"
        data = {"metadata": {}, "courses": []}
        if cs_path.exists():
            try:
                data = _json_loads(cs_path.read_text(encoding="utf-8"))
            except Exception:
                pass
        for c in data.get("courses", []):
            if c.get("note_id") == new_course["note_id"] or c.get("course_title") == title:
                return None, "课程已存在"
        data.setdefault("courses", []).append(new_course)
        cs_path.write_text(_json_dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
        return new_course, None

    def _delete_course_data(workspace_dir: str, course_id: str):
        """同步删除课程"""
        cs_path = Path(workspace_dir) / "courses_structured.json"
        if not cs_path.exists():
            return None, "课程文件不存在"
        try:
            data = _json_loads(cs_path.read_text(encoding="utf-8"))
            before = len(data.get("courses", []))
            data["courses"] = [c for c in data.get("courses", [])
                               if c.get("note_id") != course_id and c.get("course_title") != course_id]
            if len(data["courses"]) == before:
                return None, "课程未找到"
            cs_path.write_text(_json_dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
            pp = Path(workspace_dir) / "data" / "progress" / f"{course_id}.json"
            if pp.exists():
                pp.unlink()
            return True, None
        except Exception as e:
            return None, str(e)

    @app.post("/api/data/courses/review/due")
    async def courses_review_due(req: CourseProgressRequest):
        """获取待复习课时列表（间隔重复算法）"""
        data = await _run_data(_get_review_due, app.state.workspace_dir, req.course_id)
        return ok(data=data)

    @app.post("/api/data/courses/review/done")
    async def courses_review_done(req: LessonStatusRequest):
        """标记复习完成，计算下次复习时间（间隔重复）"""
        result, error = await _run_data(_mark_review_done, app.state.workspace_dir, req.course_id, req.lesson_number, req.status)
        if error:
            return err(msg=error)
        return ok(data=result)

    @app.post("/api/data/courses/updateReview")
    async def courses_update_review(req: LessonReviewRequest):
        """根据工作量更新复习调度"""
        result, error = await _run_data(_update_review_schedule, app.state.workspace_dir, req.course_id, req.lesson_number, req.workload)
        if error:
            return err(msg=error)
        return ok(data=result)

    @app.get("/api/data/courses/stats")
    async def courses_stats():
        """课程系统全局统计"""
        data = await _run_data(_get_courses_stats, app.state.workspace_dir)
        return ok(data=data)

    @app.post("/api/data/courses/create")
    async def courses_create(req: dict):
        """创建新课程"""
        title = req.get("title", "").strip()
        domain = req.get("domain", "UNKNOWN")
        if not title:
            return err(msg="课程标题不能为空")
        result, error = await _run_data(_create_course_data, app.state.workspace_dir, title, domain)
        if error:
            return err(msg=error)
        return ok(data=result)

    @app.post("/api/data/courses/delete")
    async def courses_delete(req: dict):
        """删除课程"""
        course_id = req.get("course_id", "")
        if not course_id:
            return err(msg="course_id 不能为空")
        result, error = await _run_data(_delete_course_data, app.state.workspace_dir, course_id)
        if error:
            return err(msg=error)
        return ok()

    # ─── 资源索引 API ──────────────────────────────────────

    def _get_resources(workspace_dir: str, course_id: str):
        """同步获取课程资源列表"""
        res_path = Path(workspace_dir) / "data" / "resource_index.json"
        # 回退到根目录
        if not res_path.exists():
            res_path = Path(workspace_dir) / "resource_index.json"
        if not res_path.exists():
            return []
        try:
            data = _json_loads(res_path.read_text(encoding="utf-8"))
            return data.get(course_id, [])
        except Exception:
            return []

    def _add_resource(workspace_dir: str, course_key: str, entry: dict):
        """同步添加课程资源"""
        res_path = Path(workspace_dir) / "data" / "resource_index.json"
        if not res_path.exists():
            res_path = Path(workspace_dir) / "resource_index.json"
        data = {}
        if res_path.exists():
            try:
                data = _json_loads(res_path.read_text(encoding="utf-8"))
            except Exception:
                pass
        data.setdefault(course_key, [])
        t = entry.get("type", "")
        dk = f"{t}:{entry.get('path', entry.get('url', ''))}:{entry.get('lesson_number', '')}"
        for r in data[course_key]:
            rdk = f"{r.get('type','')}:{r.get('path', r.get('url', ''))}:{r.get('lesson_number', '')}"
            if rdk == dk:
                return {"added": False}, None
        data[course_key].append(entry)
        res_path.parent.mkdir(parents=True, exist_ok=True)
        res_path.write_text(_json_dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
        return {"added": True}, None

    def _remove_resource(workspace_dir: str, course_key: str, entry: dict):
        """同步删除课程资源"""
        res_path = Path(workspace_dir) / "data" / "resource_index.json"
        if not res_path.exists():
            res_path = Path(workspace_dir) / "resource_index.json"
        if not res_path.exists():
            return {"removed": False}, None
        try:
            data = _json_loads(res_path.read_text(encoding="utf-8"))
            entries = data.get(course_key, [])
            t = entry.get("type", "")
            dk = f"{t}:{entry.get('path', entry.get('url', ''))}:{entry.get('lesson_number', '')}"
            for i, r in enumerate(entries):
                rdk = f"{r.get('type','')}:{r.get('path', r.get('url', ''))}:{r.get('lesson_number', '')}"
                if rdk == dk:
                    entries.pop(i)
                    res_path.write_text(_json_dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
                    return {"removed": True}, None
            return {"removed": False}, None
        except Exception as e:
            return None, str(e)

    @app.get("/api/data/resources/{course_id}")
    async def resources_get(course_id: str):
        """获取课程资源列表"""
        data = await _run_data(_get_resources, app.state.workspace_dir, course_id)
        return ok(data=data)

    @app.post("/api/data/resources/add")
    async def resources_add(req: dict):
        """添加课程资源"""
        course_key = req.get("course_key", "")
        entry = req.get("entry", {})
        if not course_key or not entry:
            return err(msg="参数不完整")
        result, error = await _run_data(_add_resource, app.state.workspace_dir, course_key, entry)
        if error:
            return err(msg=error)
        return ok(data=result)

    @app.post("/api/data/resources/remove")
    async def resources_remove(req: dict):
        """删除课程资源"""
        course_key = req.get("course_key", "")
        entry = req.get("entry", {})
        if not course_key:
            return err(msg="参数不完整")
        result, error = await _run_data(_remove_resource, app.state.workspace_dir, course_key, entry)
        if error:
            return err(msg=error)
        return ok(data=result)

    def _get_all_resources(workspace_dir: str, query: str = ""):
        """同步获取所有课程资源（可选搜索过滤）"""
        res_path = Path(workspace_dir) / "data" / "resource_index.json"
        if not res_path.exists():
            res_path = Path(workspace_dir) / "resource_index.json"
        if not res_path.exists():
            return []
        try:
            data = _json_loads(res_path.read_text(encoding="utf-8"))
        except Exception:
            return []
        results = []
        q = query.lower() if query else ""
        for course_id, resources in data.items():
            for r in resources:
                if q:
                    label = r.get("label", "").lower()
                    path = r.get("path", "").lower()
                    url = r.get("url", "").lower()
                    if q not in label and q not in path and q not in url and q not in course_id.lower():
                        continue
                results.append({**r, "course_id": course_id})
        # 按添加时间倒序
        results.sort(key=lambda x: x.get("added_at", x.get("lesson_number", 0) or 0), reverse=True)
        return results

    @app.get("/api/data/resources")
    async def resources_all(query: str = ""):
        """获取所有课程资源（可选搜索过滤）"""
        data = await _run_data(_get_all_resources, app.state.workspace_dir, query)
        return ok(data=data)

    # ─── 数据枢纽 API ──────────────────────────────────────

    @app.post("/api/hub/rss/list")
    async def hub_rss_list(req: Request):
        """列出 RSS 订阅源"""
        hub = _get_hub()
        if not hub:
            return err(msg="数据枢纽未初始化")
        body = await _hub_body(req)
        subs = await _run_data(hub.get_rss_subscriptions, bool(body.get("active_only", False)))
        return ok(data={
            "count": len(subs),
            "subscriptions": [
                {"id": s.id, "title": s.title, "url": s.url, "category": s.category,
                 "active": s.active, "last_polled": s.last_polled,
                 "poll_interval_minutes": s.poll_interval_minutes}
                for s in subs
            ],
        })

    @app.post("/api/hub/rss/add")
    async def hub_rss_add(req: Request):
        """添加 RSS 订阅"""
        hub = _get_hub()
        if not hub:
            return err(msg="数据枢纽未初始化")
        body = await _hub_body(req)
        url = (body.get("url") or "").strip()
        if not url:
            return err(msg="URL 不能为空")
        from ws2_data_hub import RSSSubscription
        sub = RSSSubscription(
            url=url,
            title=(body.get("title") or "").strip(),
            category=(body.get("category") or "").strip(),
            poll_interval_minutes=int(body.get("poll_interval_minutes") or 60),
        )
        result = await _run_data(hub.add_rss_subscription, sub)
        return ok(data={"id": result.id, "url": result.url, "title": result.title})

    @app.post("/api/hub/rss/remove")
    async def hub_rss_remove(req: Request):
        """移除 RSS 订阅"""
        hub = _get_hub()
        if not hub:
            return err(msg="数据枢纽未初始化")
        body = await _hub_body(req)
        sub_id = (body.get("sub_id") or "").strip()
        if not sub_id:
            return err(msg="sub_id 不能为空")
        success = await _run_data(hub.remove_rss_subscription, sub_id)
        if not success:
            return err(msg=f"未找到订阅: {sub_id}")
        return ok(data={"removed_id": sub_id})

    @app.post("/api/hub/rss/update")
    async def hub_rss_update(req: Request):
        """更新 RSS 订阅（active/poll_interval_minutes/category 等）"""
        hub = _get_hub()
        if not hub:
            return err(msg="数据枢纽未初始化")
        body = await _hub_body(req)
        sub_id = (body.get("sub_id") or "").strip()
        if not sub_id:
            return err(msg="sub_id 不能为空")
        updates = {}
        for k in ("title", "category", "poll_interval_minutes", "active"):
            if k in body and body[k] is not None:
                updates[k] = body[k]
        if not updates:
            return err(msg="未提供任何更新字段")
        success = await _run_data(lambda: hub.update_rss_subscription(sub_id, **updates))
        if not success:
            return err(msg=f"未找到订阅: {sub_id}")
        return ok(data={"updated": list(updates.keys())})

    @app.post("/api/hub/rss/poll")
    async def hub_rss_poll(req: Request):
        """轮询 RSS：sub_id 为空则轮询全部活跃订阅"""
        hub = _get_hub()
        if not hub:
            return err(msg="数据枢纽未初始化")
        body = await _hub_body(req)
        sub_id = (body.get("sub_id") or "").strip()
        if sub_id:
            new_items = await _run_data(hub.poll_rss_feed, sub_id)
            return ok(data={
                "sub_id": sub_id, "total_new": len(new_items),
                "new_items": [{"id": i.id, "title": i.title, "url": i.url} for i in new_items[:20]],
            })
        results = await _run_data(hub.poll_all_rss_feeds)
        total = sum(v for v in results.values() if v > 0)
        return ok(data={"results": results, "total_new": total})

    @app.post("/api/hub/items")
    async def hub_items_query(req: Request):
        """查询数据项（支持过滤）"""
        hub = _get_hub()
        if not hub:
            return err(msg="数据枢纽未初始化")
        body = await _hub_body(req)
        items = await _run_data(lambda: hub.query_items(
            source_type=body.get("source_type") or None,
            item_type=body.get("item_type") or None,
            tag=body.get("tag") or None,
            starred_only=bool(body.get("starred_only")),
            unread_only=bool(body.get("unread_only")),
            search=body.get("search") or None,
            limit=int(body.get("limit") or 100),
            offset=int(body.get("offset") or 0),
        ))
        return ok(data={
            "count": len(items),
            "items": [{
                "id": i.id, "title": i.title, "url": i.url, "summary": i.summary,
                "content": i.content,
                "source_type": i.source_type, "item_type": i.item_type,
                "tags": i.tags, "is_read": i.is_read, "is_starred": i.is_starred,
                "metadata": i.metadata, "created_at": i.created_at, "updated_at": i.updated_at,
            } for i in items],
        })

    @app.post("/api/hub/items/add")
    async def hub_items_add(req: Request):
        """添加数据项"""
        hub = _get_hub()
        if not hub:
            return err(msg="数据枢纽未初始化")
        body = await _hub_body(req)
        title = (body.get("title") or "").strip()
        if not title:
            return err(msg="标题不能为空")
        from ws2_data_hub import HubItem
        item = HubItem(
            title=title, url=body.get("url") or "", content=body.get("content") or "",
            summary=body.get("summary") or "",
            source_type=body.get("source_type") or "manual",
            item_type=body.get("item_type") or "webpage",
            tags=[t.strip() for t in (body.get("tags") or "").split(",") if t.strip()],
        )
        result = await _run_data(hub.add_item, item)
        return ok(data={"id": result.id, "title": result.title})

    @app.post("/api/hub/items/update")
    async def hub_items_update(req: Request):
        """更新数据项（is_read/is_starred 等）"""
        hub = _get_hub()
        if not hub:
            return err(msg="数据枢纽未初始化")
        body = await _hub_body(req)
        item_id = (body.get("item_id") or "").strip()
        if not item_id:
            return err(msg="item_id 不能为空")
        updates = {}
        for k in ("title", "content", "summary", "is_read", "is_starred", "is_archived"):
            if k in body and body[k] is not None:
                updates[k] = body[k]
        if not updates:
            return err(msg="未提供任何更新字段")
        success = await _run_data(lambda: hub.update_item(item_id, **updates))
        if not success:
            return err(msg=f"未找到数据项: {item_id}")
        return ok(data={"updated": list(updates.keys())})

    @app.post("/api/hub/items/delete")
    async def hub_items_delete(req: Request):
        """删除数据项"""
        hub = _get_hub()
        if not hub:
            return err(msg="数据枢纽未初始化")
        body = await _hub_body(req)
        item_id = (body.get("item_id") or "").strip()
        if not item_id:
            return err(msg="item_id 不能为空")
        success = await _run_data(hub.delete_item, item_id)
        if not success:
            return err(msg=f"未找到数据项: {item_id}")
        return ok(data={"deleted_id": item_id})

    @app.post("/api/hub/collections")
    async def hub_collections():
        """列出数据集合"""
        hub = _get_hub()
        if not hub:
            return err(msg="数据枢纽未初始化")
        colls = await _run_data(hub.get_collections)
        return ok(data={
            "count": len(colls),
            "collections": [
                {"id": c.id, "title": c.title, "description": c.description,
                 "item_count": len(c.item_ids), "tags": c.tags}
                for c in colls
            ],
        })

    @app.post("/api/hub/collections/create")
    async def hub_collections_create(req: Request):
        """创建数据集合"""
        hub = _get_hub()
        if not hub:
            return err(msg="数据枢纽未初始化")
        body = await _hub_body(req)
        title = (body.get("title") or "").strip()
        if not title:
            return err(msg="标题不能为空")
        from ws2_data_hub import DataCollection
        coll = DataCollection(title=title, description=body.get("description") or "")
        result = await _run_data(hub.create_collection, coll)
        return ok(data={"id": result.id, "title": result.title})

    @app.post("/api/hub/collections/addItem")
    async def hub_collections_add_item(req: Request):
        """向集合添加数据项"""
        hub = _get_hub()
        if not hub:
            return err(msg="数据枢纽未初始化")
        body = await _hub_body(req)
        ok_flag = await _run_data(
            hub.add_to_collection, body.get("collection_id", ""), body.get("item_id", ""))
        if not ok_flag:
            return err(msg="集合或数据项不存在")
        return ok(data={})

    @app.post("/api/hub/collections/delete")
    async def hub_collections_delete(req: Request):
        """删除数据集合"""
        hub = _get_hub()
        if not hub:
            return err(msg="数据枢纽未初始化")
        body = await _hub_body(req)
        await _run_data(hub.delete_collection, body.get("collection_id", ""))
        return ok(data={})

    @app.post("/api/hub/stats")
    async def hub_stats():
        """数据枢纽统计"""
        hub = _get_hub()
        if not hub:
            return err(msg="数据枢纽未初始化")
        stats = await _run_data(hub.get_statistics)
        return ok(data=stats)

    @app.post("/api/hub/pipeline/run")
    async def hub_pipeline_run(req: Request):
        """运行数据管道（空 stage=完整管道）"""
        hub = _get_hub()
        if not hub:
            return err(msg="数据枢纽未初始化")
        body = await _hub_body(req)
        engine = hub.get_pipeline_engine()
        stage = (body.get("stage") or "").strip()
        if stage:
            stage_map = {
                "scan": engine._stage_scan,
                "enrich": lambda: engine._stage_enrich(limit=int(body.get("enrich_limit") or 20)),
                "filter": engine._stage_filter,
                "update": lambda: engine._stage_update(max_age_hours=int(body.get("update_max_age_hours") or 24)),
                "syncback": engine._stage_syncback,
            }
            if stage not in stage_map:
                return err(msg=f"未知阶段: {stage}")
            result = await _run_data(stage_map[stage])
            return ok(data={"stage": stage, "result": result})
        result = await _run_data(engine.run_full_pipeline)
        return ok(data=result)

    @app.post("/api/hub/pipeline/status")
    async def hub_pipeline_status():
        """管道运行状态"""
        hub = _get_hub()
        if not hub:
            return err(msg="数据枢纽未初始化")
        status = await _run_data(hub.get_pipeline_status)
        stage_stats = await _run_data(hub.get_pipeline_stage_stats)
        return ok(data={**status, "stage_stats": stage_stats})

    @app.post("/api/hub/logs")
    async def hub_logs(req: Request):
        """管道日志"""
        hub = _get_hub()
        if not hub:
            return err(msg="数据枢纽未初始化")
        body = await _hub_body(req)
        logs = await _run_data(hub.get_pipeline_logs, int(body.get("limit") or 50))
        return ok(data={"logs": logs})

    @app.post("/api/hub/discover")
    async def hub_discover(req: Request):
        """轻度爬取探测 URL，可选自动导入发现的订阅源"""
        hub = _get_hub()
        if not hub:
            return err(msg="数据枢纽未初始化")
        body = await _hub_body(req)
        url = (body.get("url") or "").strip()
        if not url:
            return err(msg="URL 不能为空")
        try:
            info = await _run_data(hub.lightweight_crawl, url)
        except Exception as _e:
            logger.warning("hub discover failed: %s", _e)
            return ok(data={"url": url, "error": str(_e), "feeds": []})
        feeds = info.get("feeds", [])
        imported = 0
        if body.get("discover_feeds") and feeds:
            imported = await _run_data(
                hub.import_discovered_subscriptions,
                [{"url": f["url"], "type": f.get("type", "rss"), "title": f.get("title", "")} for f in feeds],
            )
        return ok(data={
            "url": url, "title": info.get("title", ""), "description": info.get("description", ""),
            "feeds": feeds, "imported": imported,
        })

    # ─── 内容推送 API ──────────────────────────────────────

    @app.get("/api/push/dashboard")
    async def push_dashboard():
        """内容推送：返回用户需要关注的摘要数据"""
        data = await _run_push(_get_push_dashboard, app.state.workspace_dir)
        return ok(data=data)

    # ─── 工作流日志 API ──────────────────────────────────────

    def _get_workflow_log(workspace_dir: str):
        """同步获取工作流日志"""
        log_path = Path(workspace_dir) / "data" / "workflow_log.json"
        if not log_path.exists():
            return []
        try:
            entries = _json_loads(log_path.read_text(encoding="utf-8"))
            entries.sort(key=lambda x: x.get("timestamp", ""))
            return entries[-200:]
        except Exception:
            return []

    def _add_workflow_log(workspace_dir: str, entry: dict):
        """同步添加工作流日志条目"""
        log_path = Path(workspace_dir) / "data" / "workflow_log.json"
        entries = []
        if log_path.exists():
            try:
                entries = _json_loads(log_path.read_text(encoding="utf-8"))
            except Exception:
                pass
        entries.append(entry)
        entries.sort(key=lambda x: x.get("timestamp", ""))
        log_path.parent.mkdir(parents=True, exist_ok=True)
        log_path.write_text(_json_dumps(entries, ensure_ascii=False, indent=2), encoding="utf-8")
        return entry

    def _get_workflow_stats(workspace_dir: str):
        """同步获取工作流统计"""
        log_path = Path(workspace_dir) / "data" / "workflow_log.json"
        if not log_path.exists():
            return {"total_entries": 0, "total_focus_time": 0, "daily": {}}
        try:
            entries = _json_loads(log_path.read_text(encoding="utf-8"))
            total_focus_time = 0.0
            note_count = 0
            complete_count = 0
            daily = {}
            course_times = {}
            for e in entries:
                ts = e.get("timestamp", "")
                etype = e.get("type", "")
                cid = e.get("course_id", "")
                if ts:
                    day = ts[:10]
                    daily[day] = daily.get(day, 0) + 1
                if etype == "note":
                    note_count += 1
                elif etype == "lesson_complete":
                    complete_count += 1
                elif etype == "timer_stop":
                    detail = e.get("detail", "")
                    try:
                        parts = detail.replace("耗时 ", "").split()
                        if parts:
                            t_str = parts[0]
                            t_parts = t_str.split(":")
                            if len(t_parts) == 3:
                                t = int(t_parts[0]) * 3600 + int(t_parts[1]) * 60 + int(t_parts[2])
                            elif len(t_parts) == 2:
                                t = int(t_parts[0]) * 60 + int(t_parts[1])
                            else:
                                t = 0
                            total_focus_time += t
                            course_times[cid] = course_times.get(cid, 0) + t
                    except Exception:
                        pass
            return {
                "total_entries": len(entries),
                "total_focus_time": round(total_focus_time, 1),
                "total_focus_hours": round(total_focus_time / 3600, 1),
                "note_count": note_count,
                "complete_count": complete_count,
                "daily": daily,
                "course_times": course_times,
            }
        except Exception:
            return {}

    def _get_push_dashboard(workspace_dir: str):
        """聚合推送数据：待办任务、待复习课程、课程资源更新、近期截止任务"""
        result = {
            "due_tasks": [],        # 近期截止的任务（今日或3天内）
            "overdue_tasks": [],    # 已超期的任务
            "in_progress_tasks": [],# 进行中的任务（截止日较远或无截止日）
            "pending_tasks": [],    # 其他待办任务
            "due_reviews": [],      # 待复习课时
            "recent_resources": [], # 最近添加的课程资源
            "today_stats": {},      # 今日统计
        }

        workspace = Path(workspace_dir)
        today = time.strftime("%Y-%m-%d")

        # 1. 任务推送
        tb_path = workspace / "task_board.json"
        if tb_path.exists():
            try:
                tasks = _json_loads(tb_path.read_text(encoding="utf-8"))
                if isinstance(tasks, list):
                    _changed = False
                    for task in tasks:
                        status = task.get("status", "")
                        if status in ("已完成", "done", "completed"):
                            continue
                        # 已过开始时间但状态仍为"待办"的，自动改为"进行中"
                        start_time = task.get("start_time", "")
                        if status == "待办" and start_time and start_time <= today:
                            status = "进行中"
                            task["status"] = "进行中"
                            _changed = True
                        due_date = task.get("due_date", "")
                        if due_date:
                            _task_base = {
                                "id": task.get("id", ""),
                                "title": task.get("title", ""),
                                "due_date": due_date,
                                "priority": task.get("priority", "中"),
                                "start_time": task.get("start_time", ""),
                                "duration": task.get("duration", 0),
                            }
                            if due_date < today:
                                _task_base["overdue_days"] = (datetime.date.today() - datetime.date.fromisoformat(due_date)).days if len(due_date) == 10 else 0
                                result["overdue_tasks"].append(_task_base)
                            elif due_date <= today:
                                result["due_tasks"].append(_task_base)
                            elif due_date <= (datetime.date.today() + datetime.timedelta(days=3)).isoformat():
                                result["due_tasks"].append(_task_base)
                            else:
                                target = "in_progress_tasks" if status == "进行中" else "pending_tasks"
                                result[target].append(_task_base)
                        else:
                            target = "in_progress_tasks" if status == "进行中" else "pending_tasks"
                            result[target].append({
                                "id": task.get("id", ""),
                                "title": task.get("title", ""),
                                "due_date": "",
                                "priority": task.get("priority", "中"),
                                "start_time": task.get("start_time", ""),
                                "duration": task.get("duration", 0),
                            })
                    if _changed:
                        try:
                            tb_path.write_text(_json_dumps(tasks), encoding="utf-8")
                        except Exception as e:
                            logger.warning(f"Push dashboard: write tasks failed: {e}")
            except Exception as e:
                logger.warning(f"Push dashboard: read tasks failed: {e}")

        # 2. 待复习课时推送
        progress_dir = workspace / "data" / "progress"
        cs_path = workspace / "courses_structured.json"
        if cs_path.exists() and progress_dir.exists():
            try:
                courses_data = _json_loads(cs_path.read_text(encoding="utf-8"))
                courses = courses_data.get("courses", []) if isinstance(courses_data, dict) else courses_data
                for course in courses:
                    course_id = course.get("note_id", course.get("id", ""))
                    course_title = course.get("course_title", course.get("title", ""))
                    pp = progress_dir / f"{course_id}.json"
                    if not pp.exists():
                        continue
                    try:
                        progress = _json_loads(pp.read_text(encoding="utf-8"))
                        review_schedule = progress.get("review_schedule", {})
                        for lesson_num, schedule in review_schedule.items():
                            if isinstance(schedule, dict):
                                next_review = schedule.get("next_review", "")
                                if next_review and next_review <= today:
                                    lesson = None
                                    for l in course.get("lessons", []):
                                        if str(l.get("lesson_number", l.get("number", ""))) == str(lesson_num):
                                            lesson = l
                                            break
                                    result["due_reviews"].append({
                                        "course_id": course_id,
                                        "course_title": course_title,
                                        "lesson_number": int(lesson_num),
                                        "lesson_title": lesson.get("lesson_title", lesson.get("title", f"课时{lesson_num}")) if lesson else f"课时{lesson_num}",
                                        "next_review": next_review,
                                    })
                    except Exception:
                        continue
            except Exception as e:
                logger.warning(f"Push dashboard: read reviews failed: {e}")

        # 3. 最近课程资源
        res_path = workspace / "data" / "resource_index.json"
        if res_path.exists():
            try:
                res_data = _json_loads(res_path.read_text(encoding="utf-8"))
                all_resources = []
                for course_id, resources in res_data.items():
                    for r in resources:
                        r["_course_id"] = course_id
                        all_resources.append(r)
                # 按修改时间排序，取最近10个
                all_resources.sort(key=lambda x: x.get("added_at", x.get("modified", "")), reverse=True)
                result["recent_resources"] = all_resources[:10]
            except Exception:
                pass

        # 3.5 RSS 新条目推送（数据枢纽未初始化时跳过）
        try:
            if _hub_ready:
                rss_items = _hub.query_items(source_type="rss", unread_only=True, limit=10)
                result["rss_new_entries"] = [{
                    "id": i.id, "title": i.title, "url": i.url, "summary": i.summary,
                    "source_type": i.source_type,
                    "sub_title": i.metadata.get("rss_sub_title", ""),
                    "published": i.metadata.get("published", ""),
                    "created_at": i.created_at,
                } for i in rss_items]
                result["rss_new_count"] = len(rss_items)
            else:
                result["rss_new_entries"] = []
                result["rss_new_count"] = 0
        except Exception as _e:
            logger.warning(f"Push dashboard: read rss failed: {_e}")
            result["rss_new_entries"] = []
            result["rss_new_count"] = 0

        # 4. 今日统计
        result["today_stats"] = {
            "overdue_tasks_count": len(result["overdue_tasks"]),
            "due_tasks_count": len(result["due_tasks"]),
            "in_progress_tasks_count": len(result["in_progress_tasks"]),
            "pending_tasks_count": len(result["pending_tasks"]),
            "due_reviews_count": len(result["due_reviews"]),
        }

        return result

    @app.get("/api/data/workflow/log")
    async def workflow_log_get():
        """获取工作流日志"""
        data = await _run_data(_get_workflow_log, app.state.workspace_dir)
        return ok(data=data)

    @app.post("/api/data/workflow/log")
    async def workflow_log_add(req: dict):
        """添加工作流日志条目"""
        from datetime import datetime
        entry = {
            "type": req.get("type", "action"),
            "action": req.get("action", ""),
            "course_id": req.get("course_id", ""),
            "lesson_number": req.get("lesson_number"),
            "detail": req.get("detail", ""),
            "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        }
        result = await _run_data(_add_workflow_log, app.state.workspace_dir, entry)
        return ok(data=result)

    @app.get("/api/data/workflow/stats")
    async def workflow_stats_get():
        """获取工作流统计"""
        data = await _run_data(_get_workflow_stats, app.state.workspace_dir)
        return ok(data=data)

    # ─── Agent 聊天 API（连接 TS2 MCP 系统）────────────────

    # 共享基础组件（LLM, config, context_injector）— 所有会话共用
    _agent_base = None
    _agent_base_lock = threading.Lock()

    # Agent 实例池 — 每个会话独立 Agent 实例
    _agent_pool: dict = {}
    _agent_pool_lock = threading.Lock()
    
    # Agent 池状态推送函数引用（供回调使用）— 使用列表作为可变容器
    _push_agent_pool_status_ref = [None]

    # 会话存储 — 独立于 checkpoint，用于会话列表/切换/恢复
    _session_store = None
    _session_store_lock = threading.Lock()

    def _get_session_store():
        nonlocal _session_store, _session_store_lock
        if _session_store is not None:
            return _session_store
        with _session_store_lock:
            if _session_store is not None:
                return _session_store
            try:
                from ..harness.session_store import SessionStore
                _session_store = SessionStore()
                logger.info(f"SessionStore initialized at {_session_store.store_dir}")
            except Exception as e:
                logger.error(f"Failed to init SessionStore: {e}")
                _session_store = None
            return _session_store

    def _extract_session_name(messages: list, checkpoints: list, fallback_id: str) -> str:
        """从对话内容中提取会话名称
        
        优先从用户消息中提取首条有意义的内容作为名称，
        回退到检查点中的文件路径或工具名称。
        """
        # 1. 尝试从用户消息中提取名称（跳过 system-reminder 伪消息，如 <current_date> 包装）
        for msg in messages:
            if msg.get("role") == "user":
                content = msg.get("content", "")
                if not content or not isinstance(content, str):
                    continue
                stripped = content.strip()
                if stripped.startswith("<system-reminder>") and stripped.endswith("</system-reminder>"):
                    continue
                # 去除空白和换行
                clean = content.strip().replace("\n", " ")
                # 提取前 30 个字符作为名称
                if len(clean) > 5:
                    name = clean[:30]
                    if len(clean) > 30:
                        name += "..."
                    return name
        
        # 2. 从检查点中的文件路径提取（如果有）
        if checkpoints:
            try:
                fdb_inst = _get_fdb_for_workspace(app.state.workspace_dir)
                if fdb_inst:
                    conn = fdb_inst._connect()
                    for cp in checkpoints[:3]:  # 只看前 3 个检查点
                        cp_id = cp.get("id", 0)
                        if cp_id:
                            row = conn.execute(
                                "SELECT path FROM file_versions WHERE checkpoint_id = ? LIMIT 1",
                                (cp_id,)
                            ).fetchone()
                            if row:
                                path = row["path"]
                                import os
                                filename = os.path.basename(path)
                                if filename and not filename.endswith("/"):
                                    return f"关于 {filename}"
            except Exception:
                pass
        
        # 3. 从检查点的工具名称提取
        if checkpoints:
            tool_counter = {}
            for cp in checkpoints:
                tool = cp.get("tool", "")
                if tool:
                    tool_counter[tool] = tool_counter.get(tool, 0) + 1
            if tool_counter:
                top_tool = max(tool_counter, key=tool_counter.get)
                tool_desc = {
                    "edit_file": "编辑文件", "read_file": "阅读文件",
                    "write_file": "写入文件", "list_directory": "浏览目录",
                    "cli_execute": "执行命令", "grep": "搜索代码",
                    "glob": "查找文件", "run_command": "运行命令",
                }
                return tool_desc.get(top_tool, top_tool)
        
        # 4. 回退到 ID
        return f"会话 {fallback_id[:8]}"

    def _get_agent_base(workspace_dir: str):
        """创建或获取共享的 Agent 基础组件（LLM, config, context_injector）"""
        nonlocal _agent_base, _agent_base_lock
        if _agent_base is not None:
            return _agent_base
        with _agent_base_lock:
            if _agent_base is not None:
                return _agent_base
            try:
                from ..config import get_config_manager
                from ..llm import MultiProviderManager, SimulatorLLM
                from ..agent import Agent, AgentConfig

                config_mgr = get_config_manager()
                provider_configs = config_mgr.get_provider_configs_for_manager()
                enabled_configs = [
                    cfg for cfg in provider_configs
                    if cfg.enabled and cfg.provider.value != 'simulator'
                ]

                if enabled_configs:
                    raw_llm = MultiProviderManager(enabled_configs)
                    class _AdapterLLM:
                        total_prompt_tokens = 0
                        total_completion_tokens = 0
                        def __init__(self, mgr):
                            self._mgr = mgr
                        def chat(self, messages, tools=None, on_token=None):
                            resp = self._mgr.chat_with_fallback(messages, tools, on_token)
                            self.total_prompt_tokens += resp.prompt_tokens
                            self.total_completion_tokens += resp.completion_tokens
                            return resp
                        def is_available(self):
                            return self._mgr.get_provider() is not None
                    llm = _AdapterLLM(raw_llm)
                    logger.info(f"Web Agent: 使用真实 LLM ({len(enabled_configs)} 个提供商)")
                else:
                    import os
                    api_key = os.environ.get("OPENAI_API_KEY", "")
                    base_url = os.environ.get("OPENAI_BASE_URL", os.environ.get("OPENAI_API_BASE", ""))
                    model_id = os.environ.get("TS2_MODEL_ID", "gpt-4o-mini")

                    if api_key:
                        from ..llm import ProviderConfig, ProviderType
                        env_config = ProviderConfig(
                            provider=ProviderType.OPENAI_COMPATIBLE,
                            name="env-openai",
                            api_key=api_key,
                            base_url=base_url or None,
                            model_id=model_id,
                            enabled=True,
                        )
                        raw_llm = MultiProviderManager([env_config])
                        class _AdapterLLM:
                            total_prompt_tokens = 0
                            total_completion_tokens = 0
                            def __init__(self, mgr):
                                self._mgr = mgr
                            def chat(self, messages, tools=None, on_token=None):
                                resp = self._mgr.chat_with_fallback(messages, tools, on_token)
                                self.total_prompt_tokens += resp.prompt_tokens
                                self.total_completion_tokens += resp.completion_tokens
                                return resp
                            def is_available(self):
                                return self._mgr.get_provider() is not None
                        llm = _AdapterLLM(raw_llm)
                        logger.info(f"Web Agent: 使用环境变量 LLM (model={model_id})")
                    else:
                        llm = SimulatorLLM()
                        logger.warning("Web Agent: 未找到 LLM 配置，使用模拟模式")

                _web_model_id = ""
                if enabled_configs:
                    _web_model_id = enabled_configs[0].model or ""
                elif api_key:
                    _web_model_id = model_id

                _web_ws2_system = None
                try:
                    import sys as _sys
                    _project_root = str(Path(__file__).resolve().parent.parent.parent)
                    if _project_root not in _sys.path:
                        _sys.path.insert(0, _project_root)
                    from course_tracker import CourseSystem
                    _web_ws2_system = CourseSystem()
                    logger.info(f"Web Agent: WS2 系统已注入 (课程数={len(_web_ws2_system.courses)})")
                except Exception as _e:
                    logger.warning(f"Web Agent: WS2 系统初始化失败: {_e}")

                from ..agent import AgentConfig
                agent_config = AgentConfig(
                    name="TS2 Web Agent",
                    base_dir=Path(workspace_dir),
                    workspace_root=workspace_dir,
                    mode="act",  # 共享 base 不带模式；per-session 模式在实例创建时按会话恢复
                    model_id=_web_model_id,
                    ws2_system=_web_ws2_system,
                )

                def _web_context_injector(a, u):
                    parts = [f"[当前工作目录: {workspace_dir}]"]
                    try:
                        ws_mgr = getattr(app.state, 'ws_manager', None)
                        if ws_mgr:
                            sessions = ws_mgr.get_sessions_info()
                            if sessions:
                                parts.append(f"[当前 Web 客户端连接数: {len(sessions)}]")
                                for s in sessions[:5]:
                                    parts.append(
                                        f"  - app_id={s.get('app_id','?')}, "
                                        f"session={s.get('session_id','?')[:8]}, "
                                        f"type={s.get('client_type','?')}"
                                    )
                            else:
                                parts.append("[当前无 Web 客户端连接]")
                    except Exception:
                        pass
                    return "\n".join(parts)

                _agent_base = {
                    'llm': llm,
                    'config': agent_config,
                    'context_injector': _web_context_injector,
                    'workspace_dir': workspace_dir,
                    'is_simulator': isinstance(llm, SimulatorLLM),
                }
                logger.info(f"Agent base 已初始化 (LLM: {'已配置' if not _agent_base['is_simulator'] else '未配置(模拟模式)'})")
            except Exception as e:
                logger.warning(f"Agent base 初始化失败: {e}")
                import traceback
                traceback.print_exc()
                _agent_base = None
        return _agent_base

    def _agent_state_of(agent):
        """统一读取 Agent 状态机的运行状态三元组（口径收口）。

        `_chat_active` 语义反转（is_set()=True 空闲 / False 活跃），
        `_is_streaming` / `_cancelled` 为布尔字段；三者本应同步，
        但历史上多处消费点各自 getattr 翻译双标志，口径漂移导致
        前端 agentStreaming 被覆写（如 cancel 残留 _cancelled 被当成流式中）。
        此处统一为一处判定，各端点不得再自行 getattr 翻译。
        返回 (is_streaming, is_active, is_cancelled)
        """
        is_streaming = bool(getattr(agent, '_is_streaming', False))
        is_active = False
        chat_active = getattr(agent, '_chat_active', None)
        if chat_active is not None:
            is_active = not chat_active.is_set()  # clear() = chat 正在运行
        is_cancelled = bool(getattr(agent, '_cancelled', False))
        return is_streaming, is_active, is_cancelled

    def _agent_status_name(agent):
        """会话状态枚举（对齐 kimi-code）：idle / running / streaming / awaiting_approval / aborted。

        由 _agent_state_of 三元组派生（口径收口，不自行翻译标志）：
          - _cancelled          → aborted
          - _is_streaming       → streaming
          - _awaiting_approval  → awaiting_approval（预留：审批机制置位）
          - chat 正在运行       → running（非流式运行中，如工具执行/思考间隙）
          - 其余                → idle
        """
        is_streaming, is_active, is_cancelled = _agent_state_of(agent)
        if is_cancelled:
            return "aborted"
        if is_streaming:
            return "streaming"
        if bool(getattr(agent, "_awaiting_approval", False)):
            return "awaiting_approval"
        if is_active:
            return "running"
        return "idle"

    def _sync_agent_from_store(agent, session_id: str, store) -> bool:
        """将 Agent 实例与 SessionStore 双向同步 — 自主状态机的核心

        重要：`_chat_active` 的语义是反的：
        - `_chat_active.set()` = "无 chat 运行"，`is_set()` = True
        - `_chat_active.clear()` = "chat 正在运行"，`is_set()` = False
        - 所以 `is_set()` 返回 True 表示 "chat 处于空闲状态"

        规则（双向合并，避免单向覆盖导致消息回滚丢失）：
        - chat 正在运行时（_is_streaming=True 或 _chat_active.is_set()=False）不进行同步
        - Store 消息数 > Agent → 加载到 Agent（store 更新）
        - Agent 消息数 > Store → 保存到 Store（agent 更新，如 chat 完成未保存）
        - 数量一致但 Agent 为空 → 从 Store 加载（确保 switch 后内存有历史）
        """
        is_streaming, is_active, _ = _agent_state_of(agent)
        # chat_idle = True 表示"无 chat 在运行"，可以安全同步
        # （_chat_active.is_set() = True 意味着空闲；is_active 为其取反）
        chat_idle = not is_active

        if is_streaming or not chat_idle:
            logger.debug(f"[state-machine] Skipping sync for '{session_id[:12]}': streaming={is_streaming}, chat_idle={chat_idle}")
            return False
        
        # 仲裁只看「真实对话消息」（剔除 system 模板消息）。
        # system 每次由 _init_messages + context 注入重新生成，内容可能不同；
        # 若把 [system] 计入仲裁，新 agent（仅 system）会与 store 完整历史前缀失配
        # → 落入分叉分支 → 新 agent 覆盖 store（只剩 system）→ 切换后历史丢失、空覆盖。
        def _strip_system(msgs) -> list:
            return [m for m in (msgs or []) if m.get("role") != "system"]

        agent_real = _strip_system(agent.messages)
        agent_msg_count = len(agent_real)
        record = store.get(session_id)
        
        if record and record.messages:
            store_msgs = record.messages
            store_real = _strip_system(store_msgs)
            store_msg_count = len(store_real)
            # 仲裁方向不能只看数量：数量相等但内容分叉（checkpoint 回滚 / 手工清理重生成）
            # 时新旧无法识别，旧 store 会反向覆盖新 agent（加载旧会话的根因）。
            # 用「内容前缀一致性」仲裁：谁包含对方完整前缀并更多 → 谁更新。
            def _is_prefix(prefix, full) -> bool:
                return len(prefix) <= len(full) and full[:len(prefix)] == prefix
            if agent_msg_count == 0:
                # Agent 无真实对话（仅有 system 或为空）→ 从 store 加载历史
                agent.restore_messages(store_msgs)
                logger.info(f"[state-machine] Loaded '{session_id[:12]}': {store_msg_count} msgs (empty agent)")
            elif store_msg_count == 0:
                # Store 无真实对话 → agent 胜出写回（新对话内容以 agent 为准）
                store.update(session_id, messages=agent.snapshot_messages())
                logger.info(f"[state-machine] Saved '{session_id[:12]}': pushed {agent_msg_count} msgs (empty store)")
            elif _is_prefix(store_real, agent_real) and agent_msg_count > store_msg_count:
                # Agent 是 store 的追加扩展（如 chat 完成但尚未保存）→ 防止旧 store 覆盖新消息
                store.update(session_id, messages=agent.snapshot_messages())
                logger.info(f"[state-machine] Saved '{session_id[:12]}': pushed {agent_msg_count} msgs to store (store had {store_msg_count})")
            elif _is_prefix(agent_real, store_real) and store_msg_count > agent_msg_count:
                # Store 是 agent 的追加扩展 → Store 更新，加载到 Agent
                agent.restore_messages(store_msgs)
                logger.info(f"[state-machine] Activated '{session_id[:12]}': loaded {store_msg_count} msgs from store (agent had {agent_msg_count})")
            elif agent_real == store_real:
                pass  # 真实对话内容完全一致，无需同步
            else:
                # 真实对话内容分叉：以 agent（活跃端）为准写回 store，避免旧 store 反向覆盖新 agent
                store.update(session_id, messages=agent.snapshot_messages())
                logger.info(f"[state-machine] Reconciled '{session_id[:12]}': content diverged, agent wins ({agent_msg_count} vs store {store_msg_count})")
            return True
        elif agent_msg_count > 0:
            # Agent 有真实对话但 Store 没有 → 创建 Store 记录（标题自动从首条非系统消息生成）
            store.create_with_id(session_id=session_id,
                                 name=_extract_session_name(agent.snapshot_messages(), [], session_id))
            store.update(session_id, messages=agent.snapshot_messages())
            logger.info(f"[state-machine] Created store record for '{session_id[:12]}' with {agent_msg_count} msgs")
            return True
        return False

    def _force_snapshot_to_store(agent, session_id: str, store) -> bool:
        """流式/运行中 Agent 的强制快照落盘（switch 离开时兜底）

        `_sync_agent_from_store` 在 chat 运行中（streaming / _chat_active 未空闲）
        会跳过同步，导致「流式中切换会话」时消息不落盘 —— 若该会话之后被
        回收（agent 池清理 / 服务重启）则内容丢失。
        此函数绕过 streaming 判定，直接把当前内存快照条件化写入 SessionStore：
        - store 已有更多消息时不覆盖（避免流式快照丢失 store 独有数据）；
        - 无记录时自动创建（标题沿用 _extract_session_name）。
        """
        try:
            msgs = agent.snapshot_messages()
            if not msgs:
                return False
            record = store.get(session_id)
            if record and record.messages and len(record.messages) > len(msgs):
                return False  # store 已有更多消息，不覆盖
            if not record:
                store.create_with_id(session_id=session_id,
                                     name=_extract_session_name(msgs, [], session_id))
            store.update(session_id, messages=msgs)
            logger.info(f"[force-snapshot] Saved '{session_id[:12]}' {len(msgs)} msgs (streaming fallback)")
            return True
        except Exception as e:
            logger.debug(f"[force-snapshot] error: {e}")
            return False

    # ─── Web 审批联动 ─────────────────────────────────────────
    # request_id → (ApprovalRequest, agent)，供前端 decide API 决策
    _web_approval_requests = {}
    # session_id → 活跃 SSE 流的 _safe_put（审批回调转发到当前前端连接）
    _web_approval_dispatch = {}

    def _get_agent_for_session(workspace_dir: str, session_id: str = ""):
        """获取或创建指定会话的 Agent 实例 — 多实例池 + 自主状态机"""
        base = _get_agent_base(workspace_dir)
        if not base:
            return None

        sid = session_id or "default"
        store = _get_session_store()

        with _agent_pool_lock:
            if sid in _agent_pool:
                agent = _agent_pool[sid]
                # 校正实例↔会话绑定：防止实例被其他路径（chat(session_id=...)）覆写 _active_session_id 后脱节
                if getattr(agent, '_active_session_id', '') != sid:
                    old_sid = agent._active_session_id or ''
                    agent._active_session_id = sid
                    logger.debug(f"[state-machine] Rebound agent '{sid[:12]}' (was '{old_sid[:12]}')")
                # 已有实例：同步 SessionStore 状态
                if store and sid:
                    _sync_agent_from_store(agent, sid, store)
                return agent
            
            # 创建新 Agent 实例
            from ..agent import Agent
            agent = Agent(llm=base['llm'], config=base['config'])
            agent.register_context_injector(base['context_injector'])
            agent._active_session_id = sid
            # per-session 模式恢复：从会话记录 metadata 读取，多会话互不串扰
            try:
                if store:
                    record = store.get(sid)
                    if record and record.metadata:
                        _m = record.metadata.get("agent_mode")
                        if _m in ("act", "plan"):
                            agent.config.mode = _m
            except Exception:
                pass

            # 将 Agent + 工具注册表注入共享 WorkflowEngine 单例
            # （单例已在模块加载时固定；注入后 workflow 的 AGENT/TOOL 步骤
            #   通过 WorkflowTool / /api/workflow/* 触发时获得真实执行能力）
            try:
                from ..workflow_engine import get_workflow_engine as _gwf
                _wf_engine = _gwf()  # 返回已固定的单例
                _wf_engine.set_agent(agent, {t.name: t for t in agent.tools})
            except Exception as e:
                logger.debug(f"WorkflowEngine set_agent 注入失败: {e}")
            
            # 注册状态回调
            def _agent_status_callback(agent_instance, event_type, **kwargs):
                if _push_agent_pool_status_ref[0]:
                    try:
                        _push_agent_pool_status_ref[0]()
                    except Exception:
                        pass
            
            agent.register_status_callback(_agent_status_callback)
            _agent_pool[sid] = agent
            
            # 新创建的 Agent：立即从 SessionStore 加载历史
            if store and sid:
                _sync_agent_from_store(agent, sid, store)
            
            agent_msg_count = len(agent.messages) if agent.messages else 0
            logger.info(f"Agent pool: created instance for session '{sid[:12]}' (pool size={len(_agent_pool)}, agent_msgs={agent_msg_count})")
            return agent

    def _get_web_agent(workspace_dir: str):
        """向后兼容 — 返回默认 Agent 实例"""
        return _get_agent_for_session(workspace_dir, "default")

    def _peek_agent_for_session(session_id: str = ""):
        """只读获取实例：存在则返回，不存在返回 None。

        与 _get_agent_for_session 不同：不创建实例、不触发状态机同步。
        用于只读接口（getAgentSession 30s 轮询 / checkpoint 查询），
        避免为不存在的会话反复创建空实例导致 pool 膨胀、实例与会话 id 脱节。
        """
        sid = session_id or "default"
        with _agent_pool_lock:
            return _agent_pool.get(sid)

    def _get_session_preview(messages_snapshot: list) -> str:
        """从消息快照中提取适合预览的文本（优先取 user/assistant 消息）"""
        if not messages_snapshot:
            return ""
        # 从后往前找第一条 user 或 assistant 消息
        for msg in reversed(messages_snapshot):
            if not isinstance(msg, dict):
                continue
            role = msg.get("role", "")
            content = msg.get("content", "")
            if role in ("user", "assistant") and content:
                text = content if isinstance(content, str) else str(content)
                return text[:80]
        return ""

    @app.post("/api/agent/chat")
    async def agent_chat(req: AgentChatRequest):
        """Agent 聊天接口（连接 TS2 MCP 系统）"""
        agent = _get_agent_for_session(app.state.workspace_dir, req.session_id)
        if agent is None:
            return ok(data={
                "reply": "Agent 初始化失败，请检查配置文件 (config_dir: ~/.ts2/agent_config)",
                "source": "error",
            })

        try:
            # 构建多模态消息
            message = _build_multimodal_message(req.message, req.attachments)
            reply = await asyncio.wait_for(
                _run_agent(agent.chat, message, session_id=req.session_id),
                timeout=300.0
            )
            # 对话完成后写回 SessionStore（非流式路径此前缺失，导致会话记录滞后）
            try:
                store = _get_session_store()
                if store and req.session_id:
                    _sync_agent_from_store(agent, req.session_id, store)
            except Exception as e:
                logger.warning(f"Auto-save session (non-stream) error: {e}")
            # 自动保存检查点（每5轮对话保存一次）
            try:
                msg_count = len(agent.messages)
                if msg_count > 0 and msg_count % 10 == 0:  # 每10条消息（约5轮）
                    summary = req.message[:80]
                    agent.create_checkpoint(summary=summary)
            except Exception:
                pass
            return ok(data={"reply": reply, "source": "ts2_mcp"})
        except asyncio.TimeoutError:
            logger.error("Agent chat timed out after 300s")
            return ok(data={
                "reply": "处理超时，AI 助手未在 300 秒内响应，请检查 LLM 配置或网络连接，或简化问题描述",
                "source": "timeout",
            })
        except Exception as e:
            logger.error(f"Agent chat error: {e}")
            return ok(data={
                "reply": f"处理出错：{str(e)}",
                "source": "error",
            })

    @app.post("/api/agent/inject-skill")
    async def agent_inject_skill(req: AgentInjectSkillRequest):
        """真实注入技能：读取 SKILL.md 全文，作为 user 指令插入 agent 会话上下文。

        与图片附件不同——图片只随单条消息传递；技能注入是持久指令，
        直接写入 agent.messages（role=user），后续所有对话均感知该技能。
        """
        skill_name = (req.skill_name or "").strip()
        if not skill_name:
            return err(msg="skill_name 不能为空")

        # 0. direct_text 通用注入（tool/mcp/workflow/plugin 等非技能类型）
        if (req.direct_text or "").strip():
            agent = _get_agent_for_session(app.state.workspace_dir, req.session_id)
            if agent is None:
                return err(msg="Agent 未初始化，无法注入")
            try:
                inject_msg = req.direct_text.strip()
                with getattr(agent, "_messages_lock", __import__("contextlib").nullcontext()):
                    if getattr(agent, "messages", None) is None:
                        agent.messages = []
                    agent.messages.append({"role": "user", "content": inject_msg})
                try:
                    store = _get_session_store()
                    if store and req.session_id:
                        _sync_agent_from_store(agent, req.session_id, store)
                except Exception as e:
                    logger.warning(f"inject-skill(direct) session save error: {e}")
                return ok(data={
                    "injected": True,
                    "skill_name": skill_name,
                    "content_length": len(inject_msg),
                    "preview": inject_msg[:200],
                    "message": f"指令已注入当前会话（{len(inject_msg)} 字符）",
                })
            except Exception as e:
                return err(msg=f"注入失败: {e}")

        # 1. 定位 SKILL.md（优先激活 skills/，其次市场 skills_market/）
        skill_md_path = None
        project_root = Path(__file__).resolve().parent.parent.parent
        for base in ("skills", "skills_market"):
            cand = project_root / base / skill_name / "SKILL.md"
            if cand.exists():
                skill_md_path = cand
                break
        if skill_md_path is None:
            # 兼容 discover 返回的 name 与目录名不一致（如 admin→Administrator）
            for base in ("skills", "skills_market"):
                base_dir = project_root / base
                if not base_dir.exists():
                    continue
                for sub in base_dir.iterdir():
                    if sub.is_dir() and (sub / "SKILL.md").exists():
                        try:
                            from ..skill_system import Skill as _SK
                            s = _SK.from_skill_md(sub)
                            if s and s.name == skill_name:
                                skill_md_path = sub / "SKILL.md"
                                break
                        except Exception:
                            continue
                if skill_md_path:
                    break
        if skill_md_path is None:
            return err(msg=f"未找到技能 {skill_name} 的 SKILL.md（skills/ 与 skills_market/ 均已检查）")

        # 2. 读取全文
        try:
            content = skill_md_path.read_text(encoding="utf-8")
        except Exception as e:
            return err(msg=f"读取技能文件失败: {e}")

        # 3. 注入 agent 会话
        agent = _get_agent_for_session(app.state.workspace_dir, req.session_id)
        if agent is None:
            return err(msg="Agent 未初始化，无法注入")

        try:
            inject_msg = (
                f"[系统注入技能指令] 用户启用了技能「{skill_name}」，以下为技能完整定义，"
                f"请遵守其触发条件与执行步骤（如需可调用其依赖工具）：\n\n{content}"
            )
            with getattr(agent, "_messages_lock", __import__("contextlib").nullcontext()):
                if hasattr(agent, "messages"):
                    if getattr(agent, "messages", None) is None:
                        agent.messages = []
                    agent.messages.append({"role": "user", "content": inject_msg})
            # 写回 SessionStore 持久化
            try:
                store = _get_session_store()
                if store and req.session_id:
                    _sync_agent_from_store(agent, req.session_id, store)
            except Exception as e:
                logger.warning(f"inject-skill session save error: {e}")
        except Exception as e:
            return err(msg=f"注入失败: {e}")

        return ok(data={
            "injected": True,
            "skill_name": skill_name,
            "content_length": len(content),
            "preview": content[:200],
            "message": f"技能「{skill_name}」已注入当前会话（{len(content)} 字符）",
        })

    class AgentSessionRenameRequest(BaseModel):
        """会话重命名请求（前端自主命名：根据对话内容提取）"""
        session_id: str
        name: str = ""  # 可选：显式名称；留空则由后端从对话提取

    @app.post("/api/agent/session/rename")
    async def agent_session_rename(req: AgentSessionRenameRequest):
        """自主重命名会话：显式名称或根据最新对话内容提取"""
        try:
            store = _get_session_store()
            if store is None or not req.session_id:
                return err(msg="会话存储不可用")
            record = store.get(req.session_id)
            if record is None:
                return err(msg=f"会话不存在: {req.session_id}")
            # 提取名称：显式优先；否则从会话消息提取
            name = (req.name or "").strip()
            if not name:
                agent = _peek_agent_for_session(req.session_id)
                msgs = agent.messages if (agent and getattr(agent, "messages", None)) else record.messages
                name = _extract_session_name(msgs or [], [], req.session_id)
                if not name or len(name) < 2:
                    return ok(data={"renamed": False, "name": record.name, "reason": "内容过短"})
            store.update(req.session_id, name=name)
            return ok(data={"renamed": True, "session_id": req.session_id, "name": name})
        except Exception as e:
            logger.error(f"session rename error: {e}")
            return err(msg=f"重命名失败: {e}")

    class AgentPluginActivateRequest(BaseModel):
        """插件激活请求：真正加载插件并注册其工具到 agent"""
        plugin_name: str
        session_id: str = ""

    @app.post("/api/agent/plugin/activate")
    async def agent_plugin_activate(req: AgentPluginActivateRequest):
        """激活插件：从 PluginManager 真正 load_plugin，把插件工具注册进 agent.tools。
        与 skill 注入不同——插件是"工具集合"，需加载才有实际能力。
        """
        name = (req.plugin_name or "").strip()
        if not name:
            return err(msg="plugin_name 不能为空")
        try:
            agent = _get_agent_for_session(app.state.workspace_dir, req.session_id)
            if agent is None:
                return err(msg="Agent 未初始化，无法激活插件")
            pm = getattr(agent, "_plugin_mgr", None)
            if pm is None:
                from ..plugins import PluginManager
                pm = PluginManager(plugins_dirs=[Path(app.state.workspace_dir or ".") / "plugins"])
                pm.discover_plugins()
                agent._plugin_mgr = pm
            # 检查插件存在性与环境依赖
            entry = None
            for e in pm.list_plugins():
                if e.name == name:
                    entry = e
                    break
            if entry is None:
                return err(msg=f"插件不存在: {name}（已发现: {[e.name for e in pm.list_plugins()]}）")
            missing_env = [v for v in (entry.requires_env or []) if not os.environ.get(v)]
            if missing_env:
                return err(msg=f"插件 {name} 缺少环境变量: {', '.join(missing_env)}，无法激活")
            # 加载插件（注册工具）
            pm.enable_plugin(name)
            ctx = pm.load_plugin(name)
            if ctx is None:
                return err(msg=f"插件 {name} 加载失败（register 未执行或出错）")
            # 把新工具注册进 agent.tools
            from ..tools import Tool
            registered = []
            for tool_reg in ctx.get_registered_tools():
                if any(getattr(t, 'name', '') == tool_reg.name for t in agent.tools):
                    continue
                class _PluginTool(Tool):
                    name = tool_reg.name
                    description = tool_reg.schema.get("description", "")
                    parameters = tool_reg.schema.get("parameters", {})
                    def execute(self, **kwargs):
                        result = tool_reg.handler(**kwargs)
                        return str(result) if result is not None else ""
                agent.tools.append(_PluginTool())
                registered.append(tool_reg.name)
            agent._instance_tool_schemas = None
            return ok(data={
                "activated": True,
                "plugin_name": name,
                "registered_tools": registered,
                "message": f"插件「{name}」已激活，注册工具: {', '.join(registered) or '(无新工具)'}",
            })
        except Exception as e:
            logger.error(f"plugin activate error: {e}")
            return err(msg=f"插件激活失败: {e}")

    class AgentWorkflowActivateRequest(BaseModel):
        """工作流激活请求：校验工作流存在并注入其定义"""
        workflow_id: str
        session_id: str = ""

    @app.post("/api/agent/workflow/activate")
    async def agent_workflow_activate(req: AgentWorkflowActivateRequest):
        """激活工作流：校验定义存在，将其描述作为指令注入会话（供 agent 按需启动）。"""
        wid = (req.workflow_id or "").strip()
        if not wid:
            return err(msg="workflow_id 不能为空")
        try:
            engine = _get_workflow_api_engine()
            wf_def = None
            for d in engine.persistence.list_definitions():
                if d.get("workflow_id") == wid or d.get("name") == wid:
                    wf_def = d
                    break
            if wf_def is None:
                return err(msg=f"工作流不存在: {wid}")
            name = wf_def.get("name") or wid
            desc = wf_def.get("description") or ""
            text = (f"[系统注入工作流指令] 用户引用了工作流「{name}」（workflow_id={wid}）。"
                    f"描述：{desc}。需要时请使用 workflow 工具启动它（workflow_id={wid}）。")
            agent = _get_agent_for_session(app.state.workspace_dir, req.session_id)
            if agent is None:
                return err(msg="Agent 未初始化，无法注入")
            with getattr(agent, "_messages_lock", __import__("contextlib").nullcontext()):
                if getattr(agent, "messages", None) is None:
                    agent.messages = []
                agent.messages.append({"role": "user", "content": text})
            try:
                store = _get_session_store()
                if store and req.session_id:
                    _sync_agent_from_store(agent, req.session_id, store)
            except Exception as e:
                logger.warning(f"workflow activate save error: {e}")
            return ok(data={
                "activated": True,
                "workflow_id": wid,
                "name": name,
                "message": f"工作流「{name}」已注入，可按需启动",
            })
        except Exception as e:
            logger.error(f"workflow activate error: {e}")
            return err(msg=f"工作流激活失败: {e}")


    # ─── Agent SSE 流式聊天 API ──────────────────────────────

    @app.post("/api/agent/chat/stream")
    async def agent_chat_stream(req: AgentChatRequest):
        """Agent 流式聊天接口（SSE）— 状态机自动同步"""
        logger.info(f"[stream] ENTRY session='{(req.session_id or 'default')[:12]}' msg='{(req.message or '')[:40]}'")
        agent = _get_agent_for_session(app.state.workspace_dir, req.session_id)
        if agent is None:
            async def _err():
                yield f"data: {_json_dumps_compact({'type': 'error', 'content': 'Agent 初始化失败'})}\n\n"
                yield "data: [DONE]\n\n"
            return StreamingResponse(_err(), media_type="text/event-stream")

        message_queue = asyncio.Queue()
        main_loop = asyncio.get_event_loop()
        _loop_alive = True

        def _is_loop_alive():
            """检查事件循环是否仍然存活"""
            return _loop_alive and not main_loop.is_closed()

        def _safe_put(message: dict):
            """安全地将消息放入队列，处理事件循环关闭的情况"""
            if not _is_loop_alive():
                return
            try:
                main_loop.call_soon_threadsafe(message_queue.put_nowait, message)
            except RuntimeError:
                pass  # 事件循环已关闭，静默忽略
            except Exception as e:
                logger.warning(f"_safe_put error: {e}")

        # ── Web 审批联动：工具请求授权 → SSE 推送前端弹窗 ──
        # 同一 agent 只注册一次回调（agent 为共享实例）；回调按 session_id
        # 转发到当前活跃 SSE 连接（_web_approval_dispatch 由本流设置）。
        try:
            _approval = getattr(agent, 'approval_manager', None)
            if _approval is not None:
                _reg_key = id(agent)
                if getattr(agent, '_web_approval_cb_registered', 0) != _reg_key:
                    agent._web_approval_cb_registered = _reg_key

                    def _on_web_approval(request):
                        try:
                            _web_approval_requests[request.id] = (request, agent, req.session_id or 'default')
                            # 后台线程等待请求结束（decide 立即返回 / 300s 超时返回），
                            # 无论结果如何都清理 _web_approval_requests，防止超时泄漏/幽灵条目
                            def _cleanup_when_done(r):
                                try:
                                    r.wait(timeout=300)
                                finally:
                                    _web_approval_requests.pop(r.id, None)
                            threading.Thread(target=_cleanup_when_done, args=(request,), daemon=True).start()
                            try:
                                agent._awaiting_approval = True
                            except Exception:
                                pass
                            put = _web_approval_dispatch.get(req.session_id or 'default')
                            if put is None:
                                logger.warning(f"[web审批] 无活跃连接 (sid={req.session_id[:12]})，等待超时默认拒绝")
                                return
                            tool_input = getattr(request, 'tool_input', {}) or {}
                            try:
                                ti_str = json.dumps(tool_input, ensure_ascii=False)[:800]
                            except Exception:
                                ti_str = str(tool_input)[:800]
                            put({
                                "type": "approval_request",
                                "request_id": getattr(request, 'id', ''),
                                "tool_name": getattr(request, 'tool_name', ''),
                                "reason": getattr(request, 'reason', ''),
                                "risk_level": getattr(request, 'risk_level', 'medium'),
                                "tool_input_preview": ti_str,
                            })
                        except Exception as e:
                            logger.warning(f"Web 审批回调异常: {e}")

                    _approval.on_request(_on_web_approval)
                    logger.info(f"[web审批] 回调已注册 (agent id={_reg_key})")
        except Exception as e:
            logger.warning(f"Web 审批回调注册失败: {e}")
        _web_approval_dispatch[req.session_id or 'default'] = _safe_put

        def _on_token(token: str):
            """LLM 流式 token 回调"""
            _safe_put({"type": "token", "content": token})
            # 流式 token 期间也定时持久化真实快照（agent.messages 在 LLM 返回后
            # 才有完整正文，store 稍滞后无妨：读路径走 agent_pool 内存，流式结束后
            # 强制落盘补齐）。不再把 token 累积 buffer 追加进 store——那会产生
            # 半截正文/重复 assistant 消息（2026-08-11 移除）。
            _persist_snapshot()

        def _persist_snapshot():
            """节流把 agent 内存消息快照实时写回 SessionStore。

            流式对话进行中（_sync_agent_from_store 会跳过 streaming 状态），
            这里主动持久化，避免会话管理/刷新读到 store 旧数据导致前后端偏移。
            """
            try:
                now = time.time()
                if now - _last_store_write[0] < _STORE_WRITE_INTERVAL:
                    return
                _last_store_write[0] = now
                store = _get_session_store()
                if not store or not req.session_id:
                    return
                msgs = agent.snapshot_messages()
                if not msgs:
                    return
                record = store.get(req.session_id)
                # 条件化覆写：store 已有比 agent 更多的消息（如 restore/他页操作落库）
                # 时不覆盖，避免流式快照丢失 store 独有数据（与 _sync_agent_from_store 仲裁精神一致）
                if record and record.messages and len(record.messages) > len(msgs):
                    return
                if not record:
                    store.create_with_id(session_id=req.session_id,
                                         name=_extract_session_name(msgs, [], req.session_id))
                store.update(req.session_id, messages=msgs)
            except Exception as e:
                logger.debug(f"[persist] snapshot error: {e}")

        _last_store_write = [0.0]
        _STORE_WRITE_INTERVAL = 2.0

        def _on_tool(name: str, args: dict):
            """工具调用开始回调"""
            _safe_put({"type": "tool_call", "name": name, "args": args})

        def _on_tool_result(name: str, result: str):
            """工具调用结果回调"""
            if not _is_loop_alive():
                return
            try:
                # 获取检查点（优先整数 id 与 /diff 接口 int() 主路径匹配，无则回退 git hash）
                cp_hash = ""
                try:
                    if hasattr(agent, '_middleware_chain') and agent._middleware_chain:
                        for mw in agent._middleware_chain._middlewares:
                            from ..middleware.shadow_checkpoint import CheckpointMiddleware
                            if isinstance(mw, CheckpointMiddleware):
                                cp_id = getattr(mw, '_last_checkpoint_id', 0)
                                if cp_id and cp_id > 0:
                                    cp_hash = str(cp_id)
                                else:
                                    cp_hash = getattr(mw, '_last_hash', '') or ''
                                break
                except Exception:
                    pass

                _safe_put({
                    "type": "tool_result",
                    "name": name,
                    "result": result,
                    "checkpoint_hash": cp_hash,
                })
                # 工具结果落盘后，实时持久化最新消息（含工具调用与结果）
                _persist_snapshot()
            except Exception as e:
                logger.warning(f"_on_tool_result error: {e}")

        def _run_agent_stream():
            """在线程池中运行 Agent（流式）"""
            try:
                # ask/answer 追问闭环：把 SSE 推送函数注入 agent，
                # 供 ask_followup_question 挂起时推送 ask 事件到前端。
                agent._current_ask_sender = _safe_put
                message = _build_multimodal_message(req.message, req.attachments)
                logger.info(f"Agent stream: starting chat for message: {req.message[:50]}...")
                reply = agent.chat(
                    message,
                    on_token=_on_token,
                    on_tool=_on_tool,
                    on_tool_result=_on_tool_result,
                    session_id=req.session_id,
                    # 子代理执行进度事件 → SSE：前端实时显示工具请求与进程
                    on_sub_agent_event=lambda e: _safe_put({"type": "sub_agent", **e}),
                )
                logger.info(f"Agent stream: chat completed, reply length={len(reply) if reply else 0}")
                # 对话完成后，强制持久化一次（覆盖节流，确保最终内容写回）
                _last_store_write[0] = 0.0
                _persist_snapshot()
                # 对话完成后，通过状态机自动保存消息到 SessionStore
                try:
                    store = _get_session_store()
                    if store and req.session_id:
                        _sync_agent_from_store(agent, req.session_id, store)
                    # 同时保存到 checkpoint 系统
                    try:
                        cp = agent.create_checkpoint(
                            summary=req.message[:80] if req.message else req.session_id[:12]
                        )
                        if cp:
                            logger.info(f"Auto-checkpoint saved for session '{req.session_id[:12]}'")
                    except Exception as cpe:
                        logger.warning(f"Auto-checkpoint error: {cpe}")
                except Exception as e:
                    logger.warning(f"Auto-save session error: {e}")
                try:
                    agent._current_ask_sender = None
                except Exception:
                    pass
                _safe_put({"type": "done", "content": reply})
            except Exception as e:
                try:
                    agent._current_ask_sender = None
                except Exception:
                    pass
                logger.error(f"Agent stream error: {e}")
                import traceback
                traceback.print_exc()
                _safe_put({"type": "error", "content": str(e)})

        async def _stream():
            """SSE 事件生成器"""
            nonlocal _loop_alive
            agent_future = main_loop.run_in_executor(_agent_executor, _run_agent_stream)

            try:
                while True:
                    try:
                        msg = await asyncio.wait_for(message_queue.get(), timeout=90.0)
                    except asyncio.TimeoutError:
                        yield f"data: {_json_dumps_compact({'type': 'error', 'content': '超时'})}\n\n"
                        break

                    if msg.get("type") == "done":
                        yield f"data: {_json_dumps_compact(msg)}\n\n"
                        yield "data: [DONE]\n\n"
                        break
                    elif msg.get("type") == "error":
                        yield f"data: {_json_dumps_compact(msg)}\n\n"
                        yield "data: [DONE]\n\n"
                        break
                    else:
                        yield f"data: {_json_dumps_compact(msg)}\n\n"
            except asyncio.CancelledError:
                # 客户端断开（如前端刷新/关闭标签页）：【不中断】Agent 生成。
                # 生成任务运行在 ThreadPoolExecutor，取消 async 响应不会停掉线程；
                # 让其继续跑完 → _run_agent_stream 完成后自动落盘 SessionStore，
                # 前端刷新回来即可通过恢复路径看到完整结果。
                # （若这里调用 agent.cancel() 会立即打断 LLM 流式 → 刷新即中断）
                logger.info(f"[stream] client disconnected, agent keeps running (session='{(req.session_id or '')[:12]}')")
            finally:
                _loop_alive = False  # 标记事件循环不再接受新消息（done/error 不再入队，跑完落盘即可）
                # 不取消 agent_future：任务可能仍在线程池排队/运行，cancel() 会阻止其执行或中断，
                # 导致生成不完整；让其自然完成并持久化。
                # 清理审批联动的本流转发与状态
                try:
                    _sid = req.session_id or 'default'
                    if _web_approval_dispatch.get(_sid) is _safe_put:
                        _web_approval_dispatch.pop(_sid, None)
                    try:
                        agent._awaiting_approval = False
                    except Exception:
                        pass
                except Exception:
                    pass

        return StreamingResponse(
            _stream(),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "Connection": "keep-alive",
                "X-Accel-Buffering": "no",
                "Content-Encoding": "identity",  # 阻止 GZip 中间件压缩 SSE
            }
        )

    # ─── Agent 审批决策 ──────────────────────────────────────

    @app.post("/api/agent/approval/decide")
    async def agent_approval_decide(req: Request):
        """前端审批决策：approve / deny / always_approve

        由 SSE approval_request 事件携带 request_id，前端弹窗后调用本接口，
        解除 harness ApprovalManager.request_approval 的阻塞等待。
        """
        body = {}
        try:
            body = await req.json()
        except Exception:
            pass
        request_id = body.get("request_id", "")
        decision_str = body.get("decision", "deny")
        if not request_id:
            return err(msg="缺少 request_id")
        item = _web_approval_requests.get(request_id)
        if not item:
            return err(msg=f"审批请求不存在或已过期: {request_id}")
        request, agent, _sid = item
        from ..harness import ApprovalDecision
        decision_map = {
            "approve": ApprovalDecision.APPROVE,
            "deny": ApprovalDecision.DENY,
            "always_approve": ApprovalDecision.ALWAYS_APPROVE,
        }
        decision = decision_map.get(decision_str)
        if decision is None:
            return err(msg="无效的决策类型")
        try:
            request.decide(decision)
        except Exception as e:
            return err(msg=f"决策失败: {e}")
        _web_approval_requests.pop(request_id, None)
        try:
            agent._awaiting_approval = False
        except Exception:
            pass
        return ok(data={"decided": True, "request_id": request_id, "decision": decision_str})

    @app.get("/api/agent/approval/pending")
    async def agent_approval_pending(req: Request):
        """刷新/断线恢复：返回当前 session 的待审批请求列表（前端据此重新弹窗）"""
        sid = (req.query_params.get("session_id") or "").strip()
        items = []
        for rid, (rq, _ag, rq_sid) in _web_approval_requests.items():
            if sid and rq_sid != sid:
                continue
            tool_input = getattr(rq, 'tool_input', {}) or {}
            try:
                ti_str = json.dumps(tool_input, ensure_ascii=False)[:800]
            except Exception:
                ti_str = str(tool_input)[:800]
            items.append({
                "request_id": rid,
                "tool_name": getattr(rq, 'tool_name', ''),
                "reason": getattr(rq, 'reason', ''),
                "risk_level": getattr(rq, 'risk_level', 'medium'),
                "tool_input_preview": ti_str,
            })
        return ok(data={"pending": items})

    # ─── Agent 控制接口 ──────────────────────────────────────

    @app.post("/api/agent/cancel")
    async def agent_cancel(req: Request):
        """取消指定会话的 Agent 对话"""
        session_id = ""
        try:
            body = await req.json()
            session_id = body.get("session_id", "")
        except Exception:
            pass
        agent = _get_agent_for_session(app.state.workspace_dir, session_id)
        if agent and hasattr(agent, 'cancel'):
            agent.cancel()
            return ok(data={"cancelled": True})
        return ok(data={"cancelled": False})

    @app.post("/api/agent/reset")
    async def agent_reset(req: Request):
        """重置指定会话的 Agent 对话历史"""
        session_id = ""
        try:
            body = await req.json()
            session_id = body.get("session_id", "")
        except Exception:
            pass
        agent = _get_agent_for_session(app.state.workspace_dir, session_id)
        if agent:
            await asyncio.get_event_loop().run_in_executor(None, agent.reset)
            return ok(data={"reset": True})
        return ok(data={"reset": False})

    @app.get("/api/agent/status")
    async def agent_status(session_id: str = ""):
        """获取 Agent 状态（只读：不创建实例，避免 pool 膨胀）"""
        agent = _peek_agent_for_session(session_id)
        if agent is None:
            return ok(data={"available": False, "tools": 0})
        # 模式 / 审批模式 / 工具组 / 上下文使用量
        from ..harness import get_global_approval_manager
        am = get_global_approval_manager()
        try:
            from ..prompt.context_window import estimate_messages_tokens
            est = estimate_messages_tokens(agent.messages)
        except Exception:
            est = 0
        ctx_window = getattr(agent, "model_context_window", 0) or 0
        mgr = getattr(agent, "_tool_group_mgr", None)
        try:
            active_groups = sorted(list(getattr(mgr, "_activated_groups", set()))) if mgr else []
        except Exception:
            active_groups = []
        return ok(data={
            "available": True,
            "tools": len(agent.tools),
            "model": agent.config.model_id or "unknown",
            "messages": len(agent.messages),
            "current_session": {
                "message_count": len(agent.messages),
                "has_context": bool(agent.messages and len(agent.messages) > 1),
            },
            "mode": getattr(agent, "get_mode", lambda: "act")(),
            "approval_mode": (am.get_mode().value if am else "suggest"),
            "active_tool_groups": active_groups,
            "context_usage": {
                "estimated_tokens": est,
                "context_window": ctx_window,
                "ratio": round(est / ctx_window, 3) if ctx_window else 0.0,
            },
        })

    def _iter_all_agents():
        with _agent_pool_lock:
            return list(_agent_pool.values())

    @app.get("/api/agent/mode")
    async def agent_get_mode(session_id: str = ""):
        # per-session：优先返回当前会话实例的模式，无实例则从会话记录读取
        sid = session_id or "default"
        agent = _peek_agent_for_session(sid)
        if agent is not None:
            return ok(data={"mode": getattr(agent, "get_mode", lambda: "act")()})
        mode = "act"
        try:
            store = _get_session_store()
            record = store.get(sid) if store else None
            if record and record.metadata:
                m = record.metadata.get("agent_mode")
                if m in ("act", "plan"):
                    mode = m
        except Exception:
            pass
        return ok(data={"mode": mode})

    @app.post("/api/agent/mode")
    async def agent_set_mode(req: Request):
        try:
            body = await req.json()
        except Exception:
            body = {}
        new_mode = str(body.get("mode", "act")).lower()
        if new_mode not in ("act", "plan"):
            return ok(data={"error": "invalid mode"})
        # 只作用当前会话实例，不串扰其它会话
        sid = str(body.get("session_id", "")).strip() or "default"
        agent = _peek_agent_for_session(sid)
        if agent is not None:
            try:
                agent.set_mode(new_mode)
            except Exception:
                pass
        # per-session 持久化（会话记录 metadata），刷新/重建后恢复
        try:
            store = _get_session_store()
            if store:
                record = store.get(sid)
                if record is None:
                    # 记录不存在则先创建（避免 update 对缺失记录静默失败）
                    store.create_with_id(session_id=sid, name=sid[:12])
                    record = store.get(sid)
                meta = dict(record.metadata) if record else {}
                meta["agent_mode"] = new_mode
                store.update(sid, metadata=meta)
        except Exception:
            pass
        return ok(data={"mode": new_mode})

    @app.get("/api/agent/approval/mode")
    async def agent_get_approval_mode():
        from ..harness import get_global_approval_manager
        am = get_global_approval_manager()
        return ok(data={"mode": am.get_mode().value})

    @app.post("/api/agent/approval/mode")
    async def agent_set_approval_mode(req: Request):
        try:
            body = await req.json()
        except Exception:
            body = {}
        m = str(body.get("mode", "suggest")).lower()
        from ..harness import get_global_approval_manager, ApprovalMode
        from ..config import get_config_manager
        try:
            mode_enum = ApprovalMode(m)
        except Exception:
            return ok(data={"error": "invalid mode"})
        am = get_global_approval_manager()
        am.set_mode(mode_enum)
        get_config_manager().set_approval_mode(m)
        return ok(data={"mode": m})

    @app.get("/api/agent/approval/tools")
    async def agent_list_approval_tools():
        from ..harness import get_global_approval_manager
        am = get_global_approval_manager()
        return ok(data={"tools": sorted(am.get_approved_tools())})

    @app.post("/api/agent/approval/tools")
    async def agent_add_approval_tool(req: Request):
        try:
            body = await req.json()
        except Exception:
            body = {}
        tool = str(body.get("tool", "")).strip()
        if not tool:
            return ok(data={"error": "tool required"})
        from ..harness import get_global_approval_manager
        from ..config import get_config_manager
        am = get_global_approval_manager()
        am.add_approved_tool(tool)
        try:
            get_config_manager().add_always_approved_tool(tool)
        except Exception:
            pass
        return ok(data={"tools": sorted(am.get_approved_tools())})

    @app.delete("/api/agent/approval/tools")
    async def agent_remove_approval_tool(req: Request):
        try:
            body = await req.json()
        except Exception:
            body = {}
        tool = str(body.get("tool", "")).strip()
        if not tool:
            return ok(data={"error": "tool required"})
        from ..harness import get_global_approval_manager
        from ..config import get_config_manager
        am = get_global_approval_manager()
        am.remove_approved_tool(tool)
        try:
            get_config_manager().remove_always_approved_tool(tool)
        except Exception:
            pass
        return ok(data={"tools": sorted(am.get_approved_tools())})

    @app.post("/api/agent/approval/reset")
    async def agent_reset_approval():
        from ..harness import get_global_approval_manager
        from ..config import get_config_manager
        am = get_global_approval_manager()
        am.reset_permissions()
        try:
            get_config_manager().clear_always_approved_tools()
        except Exception:
            pass
        return ok(data={"tools": []})

    @app.get("/api/agent/toolgroups")
    async def agent_list_toolgroups(sid: str = ""):
        # 优先返回当前会话对应 agent 的工具组状态，避免切换会话后显示错位
        if sid and sid in _agent_pool:
            targets = [_agent_pool[sid]]
        else:
            targets = _iter_all_agents()
        for ag in targets:
            mgr = getattr(ag, "_tool_group_mgr", None)
            if mgr is not None:
                try:
                    return ok(data=mgr.get_group_status())
                except Exception:
                    pass
        return ok(data={})

    @app.post("/api/agent/toolgroups/activate")
    async def agent_activate_toolgroup(req: Request):
        try:
            body = await req.json()
        except Exception:
            body = {}
        group = str(body.get("group", "")).strip()
        if not group:
            return ok(data={"error": "group required"})
        sid = str(body.get("session_id", "")).strip()
        if sid and sid in _agent_pool:
            targets = [_agent_pool[sid]]
        else:
            targets = _iter_all_agents()
        for ag in targets:
            mgr = getattr(ag, "_tool_group_mgr", None)
            if mgr is not None:
                try:
                    res = mgr.activate_group(group)
                    return ok(data=res)
                except Exception as e:
                    return ok(data={"error": str(e)})
        return ok(data={"error": "no active agent"})

    @app.post("/api/agent/ask/answer")
    async def agent_ask_answer(req: Request):
        """ask 追问闭环：用户答案注入 agent，恢复执行。"""
        try:
            body = await req.json()
        except Exception:
            body = {}
        request_id = str(body.get("request_id", "")).strip()
        answer = str(body.get("answer", ""))
        if not request_id:
            return ok(data={"error": "request_id required"})
        for ag in _iter_all_agents():
            if getattr(ag, "has_pending_ask", lambda r: False)(request_id):
                # 并发多挂起：答案按 rid 写入字典，唤醒对应挂起线程
                answers = getattr(ag, "_last_ask_answers", None)
                if answers is None:
                    ag._last_ask_answers = answers = {}
                answers[request_id] = answer or "（用户未输入）"
                ag.clear_pending_ask(request_id)
                print(f"[ASK] answer {request_id} 主路径匹配 写入 (实例id={id(ag)})", flush=True)
                return ok(data={"accepted": True})
        # 兜底：从全局 pending ask 池按 request_id 找 agent（WS2/隔离环境实例可能不在 pool）
        try:
            from ..agent import _GLOBAL_PENDING_ASKS
            info = _GLOBAL_PENDING_ASKS.get(request_id)
            ag = info.get("agent") if info else None
            print(f"[ASK] answer {request_id} 主路径未匹配，全局池兜底 agent={'有' if ag is not None else '无'}", flush=True)
            if ag is not None and getattr(ag, "has_pending_ask", lambda r: False)(request_id):
                answers = getattr(ag, "_last_ask_answers", None)
                if answers is None:
                    ag._last_ask_answers = answers = {}
                answers[request_id] = answer or "（用户未输入）"
                ag.clear_pending_ask(request_id)
                print(f"[ASK] answer {request_id} 全局池兜底 写入 (实例id={id(ag)})", flush=True)
                return ok(data={"accepted": True})
        except Exception:
            pass
        print(f"[ASK] answer {request_id} 匹配失败 → no matching (accepted=false)", flush=True)
        return ok(data={"accepted": False, "error": "no matching pending ask"})

    @app.get("/api/agent/ask/pending")
    async def agent_ask_pending():
        """非 SSE 环境（WS2 Agent）挂起的 ask 问题池：前端轮询拉取并弹窗"""
        try:
            from ..agent import list_global_pending_asks
            return ok(data={"pending": list_global_pending_asks()})
        except Exception as e:
            return ok(data={"pending": [], "error": str(e)})

    @app.get("/api/agent/model-info")
    async def agent_model_info():
        """获取当前模型信息（含多模态能力标记）"""
        agent = _get_web_agent(app.state.workspace_dir)
        if agent is None:
            return ok(data={"model": "unknown", "capabilities": {}})
        model_id = agent.config.model_id or "unknown"
        try:
            from ..llm import DEFAULT_MODEL_INFOS, ModelInfo
            info = DEFAULT_MODEL_INFOS.get(model_id)
            if info is None:
                if hasattr(agent, 'llm') and hasattr(agent.llm, 'config'):
                    info = agent.llm.config.model_info
                else:
                    info = ModelInfo(name=model_id, provider=None)
            return ok(data={
                "model": info.name,
                "provider": info.provider.value if info.provider else "unknown",
                "context_window": info.context_window,
                "max_tokens": info.max_tokens,
                "is_reasoning_model": info.is_reasoning_model,
                "capabilities": {
                    "supports_image_input": info.supports_image_input,
                    "supports_video_input": info.supports_video_input,
                    "supports_tools": info.supports_tools,
                    "supports_streaming": info.supports_streaming,
                },
                "pricing": {
                    "input_per_million": info.pricing_input,
                    "output_per_million": info.pricing_output,
                },
            })
        except Exception as e:
            return ok(data={"model": model_id, "capabilities": {}, "error": str(e)})

    @app.get("/api/models")
    async def list_models():
        """列出所有已知模型及其能力"""
        try:
            from ..llm import DEFAULT_MODEL_INFOS, PROVIDER_DEFAULT_MODELS, PROVIDER_DISPLAY_NAMES
            models = []
            for name, info in DEFAULT_MODEL_INFOS.items():
                models.append({
                    "name": info.name,
                    "provider": info.provider.value if info.provider else "unknown",
                    "provider_display": PROVIDER_DISPLAY_NAMES.get(info.provider, info.provider.value if info.provider else "unknown"),
                    "context_window": info.context_window,
                    "max_tokens": info.max_tokens,
                    "is_reasoning_model": info.is_reasoning_model,
                    "supports_image_input": info.supports_image_input,
                    "supports_video_input": info.supports_video_input,
                    "supports_tools": info.supports_tools,
                    "pricing_input": info.pricing_input,
                    "pricing_output": info.pricing_output,
                })
            providers = {}
            for pt, model_names in PROVIDER_DEFAULT_MODELS.items():
                providers[pt.value] = {
                    "display": PROVIDER_DISPLAY_NAMES.get(pt, pt.value),
                    "models": model_names,
                }
            return ok(data={"models": models, "providers": providers})
        except Exception as e:
            return ok(data={"models": [], "providers": {}, "error": str(e)})

    @app.get("/api/models/status")
    async def model_status(session_id: str = ""):
        """查询当前模型选择状态（mode / 全局默认 / 会话覆盖）"""
        try:
            sel = _get_model_selector()
            return ok(data=sel.get_status(session_id=session_id or None))
        except Exception as e:
            return err(msg=f"获取状态失败: {e}")

    class ModelSelectRequest(BaseModel):
        mode: Optional[str] = None
        default_model: Optional[str] = None

    @app.post("/api/models/select")
    async def model_select(req: ModelSelectRequest):
        """设置全局默认模型与 mode（route/fixed）"""
        try:
            sel = _get_model_selector()
            if req.mode:
                sel.set_mode(req.mode)
            if req.default_model:
                sel.set_default_model(req.default_model)
            return ok(data=sel.get_status())
        except Exception as e:
            return err(msg=f"设置失败: {e}")

    class SessionSelectRequest(BaseModel):
        session_id: str
        model_key: Optional[str] = None

    @app.post("/api/models/session-select")
    async def model_session_select(req: SessionSelectRequest):
        """设置会话级模型覆盖"""
        try:
            sel = _get_model_selector()
            if req.model_key:
                sel.set_session_model(req.session_id, req.model_key)
            else:
                sel.clear_session_model(req.session_id)
            return ok(data=sel.get_status(session_id=req.session_id))
        except Exception as e:
            return err(msg=f"设置会话模型失败: {e}")

    class ModelRefreshRequest(BaseModel):
        provider: Optional[str] = None

    @app.post("/api/models/refresh")
    async def model_refresh(req: ModelRefreshRequest):
        """从 provider /v1/models 刷新模型目录"""
        try:
            sel = _get_model_selector()
            results = sel.refresh_catalog(provider_value=req.provider)
            return ok(data=results)
        except Exception as e:
            return err(msg=f"刷新失败: {e}")

    @app.post("/api/models/import-opencode")
    async def model_import_opencode():
        """单向导入 opencode 配置到 providers.json"""
        try:
            from ..opencode_adapter import import_opencode_config
            return ok(data=import_opencode_config())
        except Exception as e:
            return err(msg=f"导入失败: {e}")

    class RegisterModelRequest(BaseModel):
        name: str
        provider: str = "custom"
        max_tokens: int = 4096
        context_window: int = 8192
        supports_image_input: bool = False
        supports_video_input: bool = False
        supports_tools: bool = True
        is_reasoning_model: bool = False
        pricing_input: float = 0.0
        pricing_output: float = 0.0

    @app.post("/api/models/register")
    async def register_model(req: RegisterModelRequest):
        """注册自定义模型（运行时动态添加到模型信息表）"""
        try:
            from ..llm import register_custom_model, ProviderType
            provider = ProviderType(req.provider) if req.provider in [p.value for p in ProviderType] else ProviderType.CUSTOM
            info = register_custom_model(
                name=req.name, provider=provider, max_tokens=req.max_tokens,
                context_window=req.context_window, supports_image_input=req.supports_image_input,
                supports_video_input=req.supports_video_input, supports_tools=req.supports_tools,
                is_reasoning_model=req.is_reasoning_model, pricing_input=req.pricing_input,
                pricing_output=req.pricing_output,
            )
            return ok(data={"registered": True, "model": info.name, "capabilities": {
                "supports_image_input": info.supports_image_input, "supports_video_input": info.supports_video_input,
            }})
        except Exception as e:
            return err(msg=f"注册失败: {e}")

    # ─── Agent 会话管理 API ──────────────────────────────────

    def _build_agent_pool_snapshot() -> list:
        """Agent 池实例快照统一构建（三处共用：sessions 合并 / pool-status 端点 / WS 推送）。

        统一口径：_agent_state_of（is_streaming 为流式权威，is_active 兼容 _chat_active）；
        字段含 last_active_time，消除原三处同构代码的字段漂移（pool-status 端点曾有额外字段）。
        """
        snapshot = []
        with _agent_pool_lock:
            for sid, agent in list(_agent_pool.items()):
                msg_count = len(agent.messages) if agent.messages else 0
                user_count = sum(1 for m in agent.messages if m.get('role') == 'user') if agent.messages else 0
                assistant_count = sum(1 for m in agent.messages if m.get('role') == 'assistant') if agent.messages else 0
                tool_count = sum(1 for m in agent.messages if m.get('role') == 'tool') if agent.messages else 0

                # 获取最后活跃时间（优先显式标记，回退最近消息时间戳）
                last_active = None
                if hasattr(agent, '_last_active_time') and agent._last_active_time:
                    last_active = agent._last_active_time
                elif agent.messages:
                    for m in reversed(agent.messages):
                        if m.get('role') in ('user', 'assistant'):
                            last_active = m.get('_timestamp', time.time())
                            break
                if last_active is None:
                    last_active = time.time()

                # 统一口径：_agent_state_of（is_streaming 为流式权威）
                is_streaming, is_active, _ = _agent_state_of(agent)

                snapshot.append({
                    "session_id": sid,
                    "status": _agent_status_name(agent),
                    "is_active": is_active,
                    "is_streaming": is_streaming,
                    "message_count": msg_count,
                    "user_message_count": user_count,
                    "assistant_message_count": assistant_count,
                    "tool_message_count": tool_count,
                    "last_active": last_active,
                    "last_active_time": agent._last_active_time if hasattr(agent, '_last_active_time') else None,
                })
        return snapshot

    @app.get("/api/agent/sessions")
    async def agent_sessions(include_checkpoint: bool = False):
        """列出所有会话 — 统一状态容器，每个ID只出现一次
        
        Args:
            include_checkpoint: 是否包含检查点会话（默认 False）
        """
        # 使用字典按 session_id 存储，确保唯一性
        session_map = {}
        agent_pool_status = _build_agent_pool_snapshot()

        # 1. 从 SessionStore 获取已保存的会话
        store = _get_session_store()
        if store:
            try:
                # 清理孤立会话（无用户消息的空会话）
                store.cleanup_orphaned_sessions()
                
                # 列出会话（include_checkpoint=True 时包含检查点会话）
                records = store.list_sessions(limit=100, include_checkpoint=include_checkpoint)
                for r in records:
                    user_msgs = [m for m in r.messages if m.get('role') == 'user']
                    if len(user_msgs) == 0:
                        continue
                        
                    session_map[r.id] = {
                        "id": r.id,
                        "timestamp": r.updated_at,
                        "last_accessed": r.last_accessed_at,
                        "message_count": len(r.messages),
                        "user_message_count": len(user_msgs),
                        "token_count": r.total_tokens,
                        "turn_count": r.turn_count,
                        "summary": r.name or "",
                        "preview": _get_session_preview(r.messages),
                        "session_type": r.session_type,
                        "source": "session_store",
                        "is_active": False,
                        "is_streaming": False,
                    }
            except Exception as e:
                logger.warning(f"SessionStore list error: {e}")

        # 2. Agent 池活跃实例快照 → 合并/更新会话数据（快照字段来自 _build_agent_pool_snapshot）
        pool_map = {it["session_id"]: it for it in agent_pool_status}
        with _agent_pool_lock:
            for sid, agent in list(_agent_pool.items()):
                item = pool_map.get(sid)
                if item is None:
                    continue
                msg_count = item["message_count"]
                user_count = item["user_message_count"]
                assistant_count = item["assistant_message_count"]
                tool_count = item["tool_message_count"]
                last_active = item["last_active"]
                is_active = item["is_active"]
                is_streaming = item["is_streaming"]

                # 合并逻辑：如果 SessionStore 已有此会话，更新它；否则添加新的
                if msg_count > 0 or is_active or is_streaming:
                    if sid in session_map:
                        # 更新已存在的会话（来自 Agent Pool 的数据更新）
                        existing = session_map[sid]
                        existing.update({
                            "message_count": max(existing["message_count"], msg_count),
                            "user_message_count": max(existing["user_message_count"], user_count),
                            "timestamp": max(existing["timestamp"], last_active),
                            "last_accessed": last_active,
                            "is_active": is_active,
                            "is_streaming": is_streaming,
                            "source": "agent_pool",  # 标记为活跃实例
                        })
                        # 如果 Agent Pool 有消息，更新预览
                        if msg_count > 0:
                            existing["preview"] = _get_session_preview(agent.messages) if agent.messages else existing.get("preview", "")
                    else:
                        # 添加新的活跃会话
                        summary = ""
                        if agent.messages:
                            for m in reversed(agent.messages):
                                if m.get("role") == "user":
                                    content = m.get("content", "")
                                    summary = content[:100] if isinstance(content, str) else str(content)[:100]
                                    break
                        session_map[sid] = {
                            "id": sid,
                            "timestamp": last_active,
                            "last_accessed": last_active,
                            "message_count": msg_count,
                            "user_message_count": user_count,
                            "assistant_message_count": assistant_count,
                            "tool_message_count": tool_count,
                            "token_count": 0,
                            "turn_count": user_count + assistant_count,
                            "summary": summary or f"会话 {sid[:8]}",
                            "preview": _get_session_preview(agent.messages) if agent.messages else "",
                            "session_type": "chat",
                            "source": "agent_pool",
                            "is_active": is_active,
                            "is_streaming": is_streaming,
                        }

        # 3. 转换为列表并排序
        sessions = list(session_map.values())
        sessions.sort(key=lambda s: s.get("last_accessed") or s.get("timestamp", 0), reverse=True)
        
        return ok(data={
            "sessions": sessions[:50],
            "total": len(sessions),
            "agent_pool": agent_pool_status,
            "pool_size": len(agent_pool_status),
        })

    @app.get("/api/agent/pool/status")
    async def agent_pool_status():
        """获取 Agent 池状态 — 所有活跃实例（快照字段统一来自 _build_agent_pool_snapshot）"""
        status = _build_agent_pool_snapshot()
        return ok(data={
            "pool_size": len(status),
            "instances": status,
        })

    # Agent 池状态推送 — 供后台线程调用
    def _push_agent_pool_status():
        """推送 Agent 池状态到 WebSocket — 线程安全"""
        try:
            ws_mgr = app.state.ws_manager
            if not ws_mgr:
                return
            
            # 构建状态数据（快照字段与 pool-status 端点/会话列表同源，无字段漂移）
            status_data = {
                "pool_size": len(_agent_pool),
                "timestamp": time.time(),
                "instances": _build_agent_pool_snapshot(),
            }
            
            # 在事件循环中发送
            loop = asyncio.get_event_loop()
            if loop.is_running():
                loop.create_task(ws_mgr.broadcast("agent_pool_status", 0, "", status_data))
        except Exception as e:
            logger.debug(f"Push agent pool status error: {e}")

    # 设置全局引用，供 Agent 状态回调使用
    _push_agent_pool_status_ref[0] = _push_agent_pool_status

    @app.get("/api/agent/sessions/{session_id}")
    async def agent_session_get(session_id: str):
        """获取单个会话的完整状态 — 状态机驱动"""
        store = _get_session_store()
        try:
            def _get_session():
                messages_source = None
                source_type = "session_store"
                is_streaming = False
                streaming_status = "idle"
                agent_present = False  # 实例在池中 = 活跃会话（流式/刚发送过）
                
                # 1. 只读获取已有实例（不创建，避免只读轮询导致 pool 膨胀 / 实例脱节）
                try:
                    sid = session_id or "default"
                    # 【诊断】打印池中所有实例的 key 与消息数，确认请求的 sid 是否命中活跃实例
                    try:
                        _pool_keys = [f"{k[:12]}({len(a.messages) if a.messages else 0})" for k, a in list(_agent_pool.items())]
                        logger.info(f"[session-get] req sid='{sid[:12]}' pool={_pool_keys}")
                    except Exception:
                        pass
                    agent = _peek_agent_for_session(sid)
                    if agent:
                        agent_present = True
                        messages_source = agent.snapshot_messages()
                        logger.info(f"[session-get] HIT agent '{sid[:12]}': snapshot={len(messages_source) if messages_source else 0} msgs")
                        is_streaming, _is_active, is_cancelled = _agent_state_of(agent)
                        # 统一口径：_agent_status_name（idle/running/streaming/awaiting_approval/aborted），
                        # 替代原局部 idle/streaming/error 三分支，避免与快照状态漂移
                        streaming_status = _agent_status_name(agent)
                        if messages_source:
                            # agent_live 只表示「此刻真正在流式」；_cancelled 残留（cancel 后
                            # 到下次 chat/reset 之间）不得再映射 agent_live，否则前端轮询会把
                            # 已停止的会话反向打回"流式中"（source='agent_live' → agentStreaming=true）
                            source_type = "agent_live" if is_streaming else "agent_pool"
                except Exception as e:
                    logger.warning(f"Error getting agent for session '{session_id[:12]}': {e}")
                
                # 2. 如果 Agent 没有消息，从 SessionStore 读取作为 fallback
                if not messages_source and store:
                    target = store.get(session_id)
                    if target and target.messages:
                        messages_source = target.messages
                        source_type = "session_store"
                        logger.info(f"[session-get] MISS agent, fallback session_store: {len(messages_source)} msgs (sid='{session_id[:12]}')")
                
                # 3. 检查点兜底仅用于「实例不存在且 store 无」的冷启动/旧数据迁移场景。
                #    活跃会话（agent_present=True，实例在池中）即使消息暂为空也绝不回退检查点——
                #    否则流式进行中 token 未落盘时会读到检查点里的冻结快照，前端渲染旧历史、
                #    跟不上流式而卡死（2026-08-10 修复）。
                if not messages_source and not agent_present:
                    try:
                        fdb = _get_fdb_for_workspace(app.state.workspace_dir)
                        if fdb:
                            checkpoints = fdb.get_checkpoints(session_id, count=100)
                            if checkpoints:
                                messages_source = _build_messages_from_checkpoints(checkpoints)
                                source_type = "checkpoint"
                    except Exception as ce:
                        logger.debug(f"Checkpoint restore for {session_id}: {ce}")
                
                if not messages_source:
                    return None
                
                # 【压缩视图】若当前消息含 [对话历史摘要] system 消息，计算展开历史总数，
                # 随每条摘要 system 消息附带 expanded_total（仅总数，不附内容）。
                # 完整展开历史改由前端经 /api/agent/messages/expand 分页懒加载，
                # 避免一次会话轮询就把数千条历史全量推给前端（加载性能优化）。
                compact_expanded_total = 0
                if _has_compact_summary(messages_source):
                    try:
                        from ..cache.context_reloader import get_context_reloader
                        _rel = get_context_reloader()
                        _cex = _expand_for_snapshot(_rel, messages_source)
                        compact_expanded_total = len(_cex) if _cex else 0
                        _cex = None
                    except Exception as _ce:
                        logger.debug(f"[session-get] compact expand 失败: {_ce}")
                        compact_expanded_total = 0
                
                # 构造 UI 消息列表
                ui_messages = []
                for msg in messages_source:
                    role = msg.get("role", "")
                    content = msg.get("content", "")
                    # 【诚实渲染】系统消息不再剔除，原样返回给前端展示。
                    # 此前直接 continue 跳过 system，前端永远看不到 system，
                    # 会话 UI 与后端 agent.messages 不一致，掩盖真实状态。
                    if role == "system":
                        entry = {
                            "role": "system",
                            "content": content if isinstance(content, str) else str(content),
                        }
                        # 压缩摘要：附带可找回总数（内容由 expand 端点分页懒加载）
                        if isinstance(content, str) and content.startswith("[对话历史摘要") and compact_expanded_total > 0:
                            entry["expanded_total"] = compact_expanded_total
                        ui_messages.append(entry)
                    elif role == "tool":
                        tool_call_id = msg.get("tool_call_id", "")
                        tool_content = content if isinstance(content, str) else str(content)
                        tool_name = ""
                        checkpoint_hash = msg.get("checkpoint_hash", "")
                        try:
                            for prev in reversed(ui_messages):
                                if prev.get("role") == "assistant" and prev.get("tool_calls"):
                                    for tc in prev["tool_calls"]:
                                        tc_dict = tc if isinstance(tc, dict) else {}
                                        if tc_dict.get("id") == tool_call_id:
                                            tool_name = tc_dict.get("function", {}).get("name", "")
                                            if not checkpoint_hash:
                                                checkpoint_hash = tc_dict.get("checkpoint_hash", "")
                                            break
                                break
                        except Exception:
                            pass
                        ui_messages.append({
                            "role": "tool",
                            "tool_call_id": tool_call_id,
                            "tool_name": tool_name,
                            "content": tool_content,
                            "checkpoint_hash": checkpoint_hash,
                        })
                    elif role == "assistant":
                        entry = {
                            "role": "assistant",
                            "content": content if isinstance(content, str) else str(content),
                        }
                        # 【诚实渲染】思考过程与内容分开透传，前端在同一气泡内按
                        # "思考 → 内容 → 工具卡" 顺序渲染，保持与 agent.messages 一致。
                        if msg.get("reasoning_content"):
                            entry["reasoning_content"] = msg["reasoning_content"]
                        if msg.get("tool_calls"):
                            entry["tool_calls"] = msg["tool_calls"]
                        ui_messages.append(entry)
                    elif role == "user" and content:
                        ui_messages.append({
                            "role": "user",
                            "content": content if isinstance(content, str) else str(content),
                        })
                
                return {
                    "messages": ui_messages,
                    "status": streaming_status,
                    "is_streaming": is_streaming,
                    "source": source_type,
                    "message_count": len(ui_messages),
                }
            
            result = await asyncio.get_event_loop().run_in_executor(None, _get_session)
            if result is None:
                logger.info(f"[session-get] RETURN empty (sid='{session_id[:12]}')")
                return ok(data={
                    "messages": [],
                    "status": "not_found",
                    "is_streaming": False,
                    "source": "none",
                    "message_count": 0,
                })
            logger.info(f"[session-get] RETURN {len(result['messages'])} ui msgs (source={result.get('source')}) sid='{session_id[:12]}'")
            return ok(data=result)
        except Exception as e:
            logger.error(f"Get session error: {e}")
            return ok(data={
                "messages": [],
                "status": "error",
                "is_streaming": False,
                "error": str(e),
            })

    @app.post("/api/agent/sessions/create")
    async def agent_session_create():
        """保存当前会话并创建新会话 — 避免创建空会话"""
        store = _get_session_store()
        try:
            def _create_session():
                summary = ""
                saved_any = False
                
                # 处理 default agent（当前会话）
                current_agent = _agent_pool.get("default")
                if current_agent:
                    messages = current_agent.snapshot_messages()
                    # 只有当有用户消息时才保存
                    user_msgs = [m for m in messages if m.get('role') == 'user']
                    if len(user_msgs) >= 1:
                        # 提取摘要（跳过 system-reminder 伪消息，如 <current_date> 包装）
                        summary = None
                        for msg in reversed(messages):
                            if msg.get("role") == "user":
                                content = msg.get("content", "")
                                if not isinstance(content, str):
                                    continue
                                stripped = content.strip()
                                if stripped.startswith("<system-reminder>") and stripped.endswith("</system-reminder>"):
                                    continue
                                summary = content[:100]
                                break
                        # 保存到 SessionStore
                        if store:
                            current_sid = current_agent._active_session_id or "default"
                            existing = store.get(current_sid)
                            if existing:
                                store.update(current_sid, messages=messages)
                            else:
                                store.create_with_id(session_id=current_sid,
                                                     name=_extract_session_name(messages, [], current_sid) or summary or "会话")
                                store.update(current_sid, messages=messages)
                        # 同时保存到 checkpoint
                        try:
                            cp = current_agent.create_checkpoint(summary=summary or "会话")
                            if cp:
                                logger.info(f"Session create: checkpoint saved for '{current_sid[:12]}'")
                                saved_any = True
                        except Exception as ce:
                            logger.warning(f"Checkpoint save error: {current_sid}: {ce}")
                
                # 重置当前 agent（清空消息，准备新对话）
                if current_agent:
                    current_agent.reset()
                
                # 生成新会话 ID（标准 sess_* 格式）
                new_id = f"sess_{uuid.uuid4().hex[:12]}"
                # 用新 ID 创建新的 agent 实例
                new_agent = _get_agent_for_session(app.state.workspace_dir, new_id)
                if new_agent:
                    new_agent.reset()  # 确保干净状态
                
                return new_id, summary

            session_id, summary = await asyncio.get_event_loop().run_in_executor(None, _create_session)
            return ok(data={
                "created": True,
                "session_id": session_id,
                "summary": summary,
                "saved_previous": True,
            })
        except Exception as e:
            logger.error(f"Create session error: {e}")
            return ok(data={"created": False, "error": str(e)})

    @app.post("/api/agent/sessions/switch")
    async def agent_session_switch(req: AgentSessionSwitchRequest):
        """切换到指定会话 — 主动激活状态机"""
        store = _get_session_store()
        try:
            def _switch_session():
                # 0. 记录目标会话是否原本就在池中（活跃会话）：是则后续绝不回退检查点，
                #    否则流式进行中切换会读到检查点冻结快照，污染活跃实例（2026-08-10 修复）
                was_in_pool = False
                with _agent_pool_lock:
                    was_in_pool = req.session_id in _agent_pool

                # 1. 保存所有非目标会话的 Agent 状态到 SessionStore
                #    _sync_agent_from_store 在 chat 运行中(streaming)会跳过；
                #    跳过时用快照强制落盘，避免流式中切换导致消息不保存、回收后丢失
                if store:
                    with _agent_pool_lock:
                        for sid, ag in _agent_pool.items():
                            if sid != req.session_id and ag.messages:
                                if not _sync_agent_from_store(ag, sid, store):
                                    _force_snapshot_to_store(ag, sid, store)

                # 2. 主动激活目标 Agent — 状态机会自动从 Store 加载
                logger.info(f"[switch] Activating session '{req.session_id[:12]}'")
                target_agent = _get_agent_for_session(app.state.workspace_dir, req.session_id)
                if not target_agent:
                    logger.warning(f"[switch] Failed to get agent for '{req.session_id[:12]}'")
                    return None
                
                # 3. 再次强制同步（确保从 Store 加载到内存）
                if store:
                    _sync_agent_from_store(target_agent, req.session_id, store)
                
                # 4. 检查目标会话状态（统一口径：_agent_state_of）
                target_is_streaming, target_is_active, _ = _agent_state_of(target_agent)
                
                # 5. 从目标 Agent 获取消息
                messages_source = target_agent.snapshot_messages()
                source_type = "agent_live" if (target_is_streaming or target_is_active) else "agent_pool"
                
                # 6. 如果 Agent 仍然没有消息，尝试从 SessionStore 直接读取（最后兜底）
                if not messages_source and store:
                    target = store.get(req.session_id)
                    if target and target.messages:
                        messages_source = target.messages
                        source_type = "session_store"
                        # 同步到 Agent
                        target_agent.restore_messages(messages_source)
                        logger.info(f"[switch] Loaded {len(messages_source)} msgs from SessionStore (fallback)")
                
                # 7. 最后兜底：从 FileVersionDB 检查点恢复（旧数据迁移兜底）。
                #    仅当目标会话原本不在池中（was_in_pool=False，冷启动/迁移场景）才允许；
                #    活跃会话即使消息为空也不回退检查点，避免读到冻结快照卡死流式（2026-08-10 修复）
                if not messages_source and not was_in_pool:
                    try:
                        fdb = _get_fdb_for_workspace(app.state.workspace_dir)
                        if fdb:
                            checkpoints = fdb.get_checkpoints(req.session_id, count=100)
                            if checkpoints:
                                messages_source = _build_messages_from_checkpoints(checkpoints)
                                source_type = "checkpoint"
                                target_agent.restore_messages(messages_source)
                                logger.info(f"[switch] Loaded {len(messages_source)} msgs from FileVersionDB checkpoints (fallback)")
                    except Exception as ce:
                        logger.warning(f"Checkpoint restore error: {ce}")
                
                if not messages_source:
                    logger.warning(f"[switch] No messages found for session '{req.session_id[:12]}'")
                    return None
                
                logger.info(f"[switch] Session '{req.session_id[:12]}' activated: {len(messages_source)} msgs (source={source_type})")
                
                # 构造 UI 消息列表
                ui_messages = []
                for msg in messages_source:
                    role = msg.get("role", "")
                    content = msg.get("content", "")
                    if role == "system":
                        continue
                    if role == "tool":
                        tool_call_id = msg.get("tool_call_id", "")
                        tool_content = content if isinstance(content, str) else str(content)
                        tool_name = ""
                        checkpoint_hash = msg.get("checkpoint_hash", "")
                        try:
                            for prev in reversed(ui_messages):
                                if prev.get("role") == "assistant" and prev.get("tool_calls"):
                                    for tc in prev["tool_calls"]:
                                        tc_dict = tc if isinstance(tc, dict) else {}
                                        if tc_dict.get("id") == tool_call_id:
                                            tool_name = tc_dict.get("function", {}).get("name", "")
                                            if not checkpoint_hash:
                                                checkpoint_hash = tc_dict.get("checkpoint_hash", "")
                                            break
                                    break
                        except Exception:
                            pass
                        ui_messages.append({
                            "role": "tool",
                            "tool_call_id": tool_call_id,
                            "tool_name": tool_name,
                            "content": tool_content,
                            "checkpoint_hash": checkpoint_hash,
                        })
                    elif role == "assistant":
                        entry = {
                            "role": "assistant",
                            "content": content if isinstance(content, str) else str(content),
                        }
                        if msg.get("tool_calls"):
                            entry["tool_calls"] = msg["tool_calls"]
                        ui_messages.append(entry)
                    elif role == "user" and content:
                        ui_messages.append({
                            "role": "user",
                            "content": content if isinstance(content, str) else str(content),
                        })
                return ui_messages

            # 获取会话状态信息（在 _switch_session 外部获取，统一口径：_agent_state_of）
            session_is_streaming = False
            session_is_active = False
            with _agent_pool_lock:
                agent = _agent_pool.get(req.session_id)
                if agent:
                    session_is_streaming, session_is_active, _ = _agent_state_of(agent)
            
            ui_messages = await asyncio.get_event_loop().run_in_executor(None, _switch_session)
            if ui_messages is not None:
                return ok(data={
                    "switched": True,
                    "session_id": req.session_id,
                    "messages": ui_messages,
                    "is_streaming": session_is_streaming,
                    "is_active": session_is_active,
                    "source": "agent_live" if (session_is_streaming or session_is_active) else "session_store",
                })
            else:
                return ok(data={"switched": False, "error": "会话不存在"})
        except Exception as e:
            logger.error(f"Switch session error: {e}")
            return ok(data={"switched": False, "error": str(e)})

    @app.post("/api/agent/sessions/delete")
    async def agent_session_delete(req: AgentSessionSwitchRequest):
        """删除指定会话 — 同时清理 SessionStore、checkpoint 和 Agent 实例池"""
        deleted = False
        # 1. 从 SessionStore 删除
        store = _get_session_store()
        if store:
            try:
                deleted = store.delete(req.session_id) or deleted
            except Exception as e:
                logger.warning(f"SessionStore delete error: {e}")

        # 2. 清理 FileVersionDB 中该会话的检查点（检查点按 session_id 字段关联）
        try:
            fdb = _get_fdb_for_workspace(app.state.workspace_dir)
            if fdb:
                n = fdb.delete_checkpoints_by_session(req.session_id)
                if n:
                    deleted = True
                    logger.info(f"Session delete: removed {n} checkpoints for '{req.session_id[:12]}'")
        except Exception as e:
            logger.warning(f"Checkpoint delete error: {e}")

        # 3. 清理 Agent 实例池
        with _agent_pool_lock:
            agent = _agent_pool.pop(req.session_id, None)
        if agent:
            try:
                agent.cancel()
            except Exception:
                pass
            logger.info(f"Agent pool: removed session '{req.session_id[:12]}' (pool size={len(_agent_pool)})")

        return ok(data={"deleted": deleted})

    # ─── Swarm 子 Agent API ──────────────────────────────────

    @app.get("/api/swarm/agents")
    async def swarm_list_agents():
        """列出所有已注册的子 Agent 及其状态"""
        agent = _get_web_agent(app.state.workspace_dir)
        if not agent or not hasattr(agent, '_coordinator') or not agent._coordinator:
            return ok(data={"agents": [], "available": False, "swarm_enabled": False})
        try:
            coordinator = agent._coordinator
            agent_names = coordinator.list_agents()
            agents_info = []
            for name in agent_names:
                spec = coordinator._specs.get(name)
                sa = coordinator._agents.get(name)
                running_tasks = []
                for tid, holder in coordinator._pending_tasks.items():
                    if holder.result and holder.result.agent_name == name:
                        running_tasks.append({"task_id": tid, "status": holder.result.status.value, "started_at": holder.result.started_at})
                info = {
                    "name": name,
                    "role": spec.role.value if spec else "custom",
                    "system_prompt": (spec.system_prompt[:80] + "...") if spec and spec.system_prompt else "",
                    "model": spec.model if spec else "",
                    "max_turns": spec.max_turns if spec else 0,
                    "allowed_tools": spec.allowed_tools if spec else None,
                    "is_busy": sa.is_busy if sa else False,
                    "status": sa._result.status.value if sa else "idle",
                    "running_tasks": running_tasks,
                }
                agents_info.append(info)
            return ok(data={"agents": agents_info, "available": True, "swarm_enabled": coordinator.swarm_enabled})
        except Exception as e:
            logger.warning(f"swarm_list_agents error: {e}")
            return ok(data={"agents": [], "available": False, "error": str(e)})

    @app.get("/api/swarm/agents/{agent_name}")
    async def swarm_get_agent(agent_name: str):
        """获取指定子 Agent 的详情"""
        agent = _get_web_agent(app.state.workspace_dir)
        if not agent or not hasattr(agent, '_coordinator') or not agent._coordinator:
            return err(msg="Swarm 系统未初始化")
        coordinator = agent._coordinator
        spec = coordinator._specs.get(agent_name)
        if not spec:
            return err(msg=f"子 Agent '{agent_name}' 不存在")
        sa = coordinator._agents.get(agent_name)
        result = sa._result if sa else None
        return ok(data={
            "name": spec.name, "role": spec.role.value, "system_prompt": spec.system_prompt,
            "model": spec.model, "max_turns": spec.max_turns, "allowed_tools": spec.allowed_tools,
            "denied_tools": spec.denied_tools, "is_busy": sa.is_busy if sa else False,
            "status": result.status.value if result else "idle",
            "last_result": {
                "content": result.content if result and result.content else "",
                "error": result.error if result else None,
                "duration_ms": result.duration_ms if result else 0,
                "tool_calls_count": result.tool_calls_count if result else 0,
                "prompt_tokens": result.prompt_tokens if result else 0,
                "completion_tokens": result.completion_tokens if result else 0,
                "completed_at": result.completed_at if result else 0,
            } if result and result.status.value in ("completed", "failed") else None,
        })

    class SwarmRunRequest(BaseModel):
        agent_name: str
        prompt: str
        context: Optional[Dict[str, Any]] = None
        background: bool = False

    @app.post("/api/swarm/run")
    async def swarm_run(req: SwarmRunRequest):
        """运行子 Agent（同步或后台）"""
        agent = _get_web_agent(app.state.workspace_dir)
        if not agent or not hasattr(agent, '_coordinator') or not agent._coordinator:
            return err(msg="Swarm 系统未初始化")
        coordinator = agent._coordinator
        if req.agent_name not in coordinator._specs:
            return err(msg=f"子 Agent '{req.agent_name}' 不存在")
        if req.background:
            task_id = coordinator.run_async(req.agent_name, req.prompt, req.context)
            ws_mgr = get_ws_manager()
            await ws_mgr.broadcast("swarm_task_started", 0, "", {
                "task_id": task_id, "agent_name": req.agent_name, "prompt": req.prompt[:100],
            })
            return ok(data={"task_id": task_id, "background": True})
        else:
            def _run():
                return coordinator.run(req.agent_name, req.prompt, req.context)
            result = await asyncio.get_event_loop().run_in_executor(None, _run)
            return ok(data={
                "status": result.status.value, "content": result.content,
                "reasoning_content": result.reasoning_content, "error": result.error,
                "duration_ms": result.duration_ms, "tool_calls_count": result.tool_calls_count,
                "prompt_tokens": result.prompt_tokens, "completion_tokens": result.completion_tokens,
            })

    @app.post("/api/swarm/cancel/{agent_name}")
    async def swarm_cancel(agent_name: str):
        """取消子 Agent 的运行"""
        agent = _get_web_agent(app.state.workspace_dir)
        if not agent or not hasattr(agent, '_coordinator') or not agent._coordinator:
            return err(msg="Swarm 系统未初始化")
        agent._coordinator.cancel(agent_name)
        return ok(data={"cancelled": True, "agent_name": agent_name})

    @app.get("/api/swarm/tasks")
    async def swarm_list_tasks():
        """列出所有后台任务及其状态"""
        agent = _get_web_agent(app.state.workspace_dir)
        if not agent or not hasattr(agent, '_coordinator') or not agent._coordinator:
            return ok(data={"tasks": []})
        coordinator = agent._coordinator
        tasks = []
        for tid, holder in list(coordinator._pending_tasks.items()):
            task_info = {"task_id": tid, "completed": holder.completed.is_set()}
            if holder.result:
                task_info.update({
                    "agent_name": holder.result.agent_name, "status": holder.result.status.value,
                    "content": holder.result.content[:300] if holder.result.content else "",
                    "error": holder.result.error, "duration_ms": holder.result.duration_ms,
                    "tool_calls_count": holder.result.tool_calls_count,
                })
            tasks.append(task_info)
        return ok(data={"tasks": tasks})

    @app.get("/api/swarm/tasks/{task_id}")
    async def swarm_get_task(task_id: str):
        """获取后台任务结果"""
        agent = _get_web_agent(app.state.workspace_dir)
        if not agent or not hasattr(agent, '_coordinator') or not agent._coordinator:
            return err(msg="Swarm 系统未初始化")
        coordinator = agent._coordinator
        holder = coordinator._pending_tasks.get(task_id)
        if not holder:
            return ok(data={"task_id": task_id, "status": "unknown", "completed": True})
        result = holder.result
        if not result:
            return ok(data={"task_id": task_id, "status": "running", "completed": False})
        return ok(data={
            "task_id": task_id, "status": result.status.value, "completed": holder.completed.is_set(),
            "agent_name": result.agent_name, "content": result.content,
            "reasoning_content": result.reasoning_content, "error": result.error,
            "duration_ms": result.duration_ms, "tool_calls_count": result.tool_calls_count,
            "prompt_tokens": result.prompt_tokens, "completion_tokens": result.completion_tokens,
        })

    @app.post("/api/swarm/poll/{task_id}")
    async def swarm_poll_task(task_id: str):
        """轮询后台任务结果（等待最多 30 秒）"""
        agent = _get_web_agent(app.state.workspace_dir)
        if not agent or not hasattr(agent, '_coordinator') or not agent._coordinator:
            return err(msg="Swarm 系统未初始化")
        coordinator = agent._coordinator
        result = await asyncio.get_event_loop().run_in_executor(
            None, lambda: coordinator.poll_result(task_id, timeout=30.0)
        )
        if result is None:
            return ok(data={"task_id": task_id, "status": "running", "completed": False})
        ws_mgr = get_ws_manager()
        await ws_mgr.broadcast("swarm_task_completed", 0, "", {
            "task_id": task_id, "agent_name": result.agent_name, "status": result.status.value,
            "duration_ms": result.duration_ms, "content_preview": result.content[:200] if result.content else "",
        })
        return ok(data={
            "task_id": task_id, "status": result.status.value, "completed": True,
            "agent_name": result.agent_name, "content": result.content,
            "reasoning_content": result.reasoning_content, "error": result.error,
            "duration_ms": result.duration_ms, "tool_calls_count": result.tool_calls_count,
            "prompt_tokens": result.prompt_tokens, "completion_tokens": result.completion_tokens,
        })

    class SwarmEnableRequest(BaseModel):
        reason: str = ""

    @app.post("/api/swarm/enable")
    async def swarm_enable(req: SwarmEnableRequest):
        """启用大规模 Swarm 集群模式（>4个并行任务需启用）"""
        agent = _get_web_agent(app.state.workspace_dir)
        if not agent or not hasattr(agent, '_coordinator') or not agent._coordinator:
            return err(msg="Swarm 系统未初始化")
        coordinator = agent._coordinator
        if coordinator.swarm_enabled:
            return ok(data={"swarm_enabled": True, "message": "Swarm 集群模式已处于启用状态"})
        if not req.reason.strip():
            return err(msg="启用 Swarm 集群模式需要提供原因说明（如：大规模并行研究任务）")
        success = coordinator.enable_swarm(reason=req.reason)
        if success:
            return ok(data={"swarm_enabled": True, "message": f"Swarm 集群模式已启用 (原因: {req.reason})"})
        else:
            return err(msg="Swarm 启用失败：LLM 未配置")

    @app.post("/api/swarm/disable")
    async def swarm_disable():
        """禁用大规模 Swarm 集群模式"""
        agent = _get_web_agent(app.state.workspace_dir)
        if not agent or not hasattr(agent, '_coordinator') or not agent._coordinator:
            return err(msg="Swarm 系统未初始化")
        agent._coordinator.disable_swarm()
        return ok(data={"swarm_enabled": False, "message": "Swarm 集群模式已禁用（单次调用和≤4并行仍可用）"})

    # ─── 发现 API（Skill/Tool/MCP/Workflow 聚合）──────────

    @app.get("/api/discover")
    async def discover():
        """聚合发现：返回可用的技能/工具组/MCP工具/工作流定义，供前端输入栏菜单引用"""
        result = {"skills": [], "tools": {}, "mcp": [], "workflows": [], "plugins": []}

        # 1. Skill — 配置注册表 + skills_market/ 动态扫描合并（市场技能无需写 config 也可引用）
        try:
            from ..config import get_config_manager
            config_mgr = get_config_manager()
            seen = set()
            result["skills"] = []
            for s in config_mgr.get_enabled_skills():
                result["skills"].append({"name": s.name, "description": s.description, "type": getattr(s, "type", "")})
                seen.add(s.name)
            # 合并 skills_market/ 第三方技能仓库（未注册 config 也能被发现）
            try:
                from ..skill_system import Skill
                from pathlib import Path
                market_dir = Path(__file__).resolve().parent.parent.parent / "skills_market"
                if market_dir.exists():
                    for skill_subdir in sorted(market_dir.iterdir()):
                        if not skill_subdir.is_dir() or skill_subdir.name.startswith((".", "_")):
                            continue
                        try:
                            skill_obj = Skill.from_skill_md(skill_subdir)
                            if skill_obj and skill_obj.name and skill_obj.name not in seen:
                                result["skills"].append({
                                    "name": skill_obj.name,
                                    "description": skill_obj.description,
                                    "type": skill_obj.category.value if hasattr(skill_obj.category, "value") else str(skill_obj.category),
                                })
                                seen.add(skill_obj.name)
                        except Exception:
                            continue
            except Exception as e:
                logger.debug(f"discover market skills error: {e}")
        except Exception as e:
            logger.warning(f"discover skills error: {e}")

        # 2. Tool — 从工具组加载追踪器取实际加载的工具组
        try:
            from ..tools import get_loaded_tool_groups
            result["tools"] = get_loaded_tool_groups()
        except Exception as e:
            logger.warning(f"discover tools error: {e}")

        # 3. MCP — 标准 MCP 客户端已连接工具 + 远程 MCP 服务组
        try:
            from ..mcp_client.client import MCPClientManager
            mgr = MCPClientManager()
            for t in mgr.list_tools():
                result["mcp"].append({
                    "name": t.get("name", ""),
                    "description": t.get("description", ""),
                    "client": "mcp_client",
                })
            for name, info in mgr.list_clients().items():
                if getattr(info, "state", None) and str(info.state) not in ("ClientState.CONNECTED", "CONNECTED"):
                    continue
                result["mcp"].append({
                    "name": name,
                    "description": f"MCP 服务（{getattr(info, 'tool_count', 0)} 工具）",
                    "client": name,
                })
        except Exception as e:
            logger.warning(f"discover mcp error: {e}")

        # 4. Workflow — 从 WorkflowEngine 持久化层取定义（list_definitions 返回 dict 列表）
        try:
            engine = _get_workflow_api_engine()
            for d in engine.persistence.list_definitions():
                result["workflows"].append({
                    "workflow_id": d.get("workflow_id", ""),
                    "name": d.get("name", "") or d.get("workflow_id", ""),
                    "description": d.get("description", ""),
                    "version": d.get("version", ""),
                    "created_at": d.get("created_at", ""),
                })
        except Exception as e:
            logger.warning(f"discover workflows error: {e}")

        # 5. Plugin — 从 PluginManager 动态发现的插件（含 mcp/plugins/ 内置）
        try:
            from ..plugins import PluginManager
            pm = PluginManager()
            entries = pm.discover_plugins()
            result["plugins"] = [
                {
                    "name": e.name,
                    "description": f"{e.kind.value} 插件 · {len(e.provides_tools or [])} 工具",
                    "kind": e.kind.value,
                    "tools": e.provides_tools or [],
                }
                for e in entries
            ]
        except Exception as e:
            logger.debug(f"discover plugins error: {e}")

        return ok(data=result)

    # ─── Workflow 引擎 API ──────────────────────────────────

    def _get_workflow_api_engine():
        """获取进程内 WorkflowEngine 单例（模块加载时已固定并注入 agent/tools）"""
        engine = getattr(app.state, "workflow_engine", None)
        if engine is None:
            from ..workflow_engine import get_workflow_engine as _gwf
            engine = _gwf()
            app.state.workflow_engine = engine
        return engine

    @app.get("/api/workflow/definitions")
    async def workflow_list_definitions():
        """列出所有工作流定义"""
        try:
            engine = _get_workflow_api_engine()
            defs = engine.persistence.list_definitions()
            return ok(data={"definitions": defs, "count": len(defs)})
        except Exception as e:
            logger.warning(f"workflow_list_definitions error: {e}")
            return err(msg=f"获取工作流定义失败: {e}")

    class WorkflowStartRequest(BaseModel):
        workflow_id: str
        input_data: Optional[Dict[str, Any]] = None

    @app.post("/api/workflow/start")
    async def workflow_start(req: WorkflowStartRequest):
        """启动一个工作流"""
        try:
            engine = _get_workflow_api_engine()
            wf_def = engine.persistence.get_definition(req.workflow_id)
            if not wf_def:
                return err(msg=f"工作流不存在: {req.workflow_id}")
            instance_id = engine.start_workflow(wf_def, req.input_data or {})
            return ok(data={"instance_id": instance_id, "workflow_id": req.workflow_id})
        except Exception as e:
            logger.warning(f"workflow_start error: {e}")
            return err(msg=f"启动工作流失败: {e}")

    @app.get("/api/workflow/instances")
    async def workflow_list_instances():
        """列出运行实例（含状态/进度/当前步骤）"""
        try:
            engine = _get_workflow_api_engine()
            instances = engine.persistence.list_instances(limit=50)
            return ok(data={"instances": instances, "count": len(instances)})
        except Exception as e:
            logger.warning(f"workflow_list_instances error: {e}")
            return err(msg=f"获取实例失败: {e}")

    @app.post("/api/workflow/pause/{instance_id}")
    async def workflow_pause(instance_id: str):
        engine = _get_workflow_api_engine()
        return ok(data={"paused": engine.pause_workflow(instance_id), "instance_id": instance_id})

    @app.post("/api/workflow/resume/{instance_id}")
    async def workflow_resume(instance_id: str):
        engine = _get_workflow_api_engine()
        return ok(data={"resumed": engine.resume_workflow(instance_id), "instance_id": instance_id})

    @app.post("/api/workflow/cancel/{instance_id}")
    async def workflow_cancel(instance_id: str):
        engine = _get_workflow_api_engine()
        return ok(data={"cancelled": engine.cancel_workflow(instance_id), "instance_id": instance_id})

    @app.get("/api/workflow/status/{instance_id}")
    async def workflow_status(instance_id: str):
        engine = _get_workflow_api_engine()
        status = engine.get_status(instance_id)
        if not status:
            return err(msg=f"实例不存在: {instance_id}")
        return ok(data=status)

    @app.get("/api/workflow/step_results/{instance_id}")
    async def workflow_step_results(instance_id: str):
        engine = _get_workflow_api_engine()
        return ok(data={"instance_id": instance_id, "steps": engine.get_step_results(instance_id)})

    @app.get("/api/workflow/logs/{instance_id}")
    async def workflow_logs(instance_id: str, limit: int = 100):
        engine = _get_workflow_api_engine()
        return ok(data={"instance_id": instance_id, "logs": engine.get_logs(instance_id, limit=min(limit, 500))})

    # ─── Agent 检查点 API ────────────────────────────────────

    def _get_fdb_for_workspace(workspace_dir: str) -> Optional[Any]:
        """获取 FileVersionDB 实例（独立于 agent，可直接查询检查点）"""
        try:
            from ..middleware.file_version_db import FileVersionDB, default_db_path
            db_path = default_db_path(workspace_dir)
            if os.path.exists(db_path):
                return FileVersionDB(db_path)
        except Exception as e:
            logger.debug(f"Failed to get FileVersionDB: {e}")
        return None

    def _build_messages_from_checkpoints(checkpoints: list, fallback_content: str = "历史会话") -> list:
        """从 FileVersionDB 检查点记录构建 UI 消息列表（会话兜底 / 迁移恢复复用）

        注意：这里的 checkpoints 是「按 session_id 关联」的文件快照记录，
        与会话 ID（sess_*）绑定，而非 ContextReloader 的 cp-* 对话检查点。
        """
        messages = []
        for cp in reversed(checkpoints):
            tool_name = cp.get("tool", "")
            step = cp.get("step", 0)
            file_count = cp.get("file_count", 0)
            if tool_name:
                cp_id = cp.get("id", "")
                messages.append({
                    "role": "assistant",
                    "content": "",
                    "tool_calls": [{
                        "id": f"restore-{cp_id}",
                        "type": "function",
                        "function": {
                            "name": tool_name,
                            "arguments": "{}"
                        }
                    }]
                })
                messages.append({
                    "role": "tool",
                    "tool_call_id": f"restore-{cp_id}",
                    "content": f"检查点 #{step}: {file_count} 个文件变更",
                    "checkpoint_hash": cp.get("checkpoint_hash", ""),
                })
        if not messages:
            messages.append({
                "role": "user",
                "content": fallback_content,
            })
        return messages

    def _find_conversation_cp_by_ref(reloader, commit_hash: str) -> str:
        """反查 commit_hash 对应的 ContextReloader 对话快照 id（cp-*）

        用于 conversation_checkpoint_id 回填缺失/旧数据的兜底：
        - SQLite 数字 id：cp-* 快照 summary 带 '#<sqlite_id> ' 前缀（_auto_conversation_checkpoint 写入）
        - git hash：cp-* 快照 git_commit_hash 字段匹配
        """
        try:
            commit = str(commit_hash or "")
            if not commit:
                return ""
            for key in reloader.list_checkpoints():
                cp = reloader.restore_checkpoint(key)
                if cp is None:
                    continue
                if commit.isdigit():
                    if ((getattr(cp, "summary", "") or "").startswith(f"#{commit} ")):
                        return key
                else:
                    if getattr(cp, "git_commit_hash", "") == commit:
                        return key
        except Exception:
            pass
        return ""

    def _message_key(msg: dict) -> tuple:
        """消息同一性 key（用于判断检查点快照间的新增/保留消息）"""
        role = msg.get("role", "")
        if role == "tool":
            return ("tool", str(msg.get("tool_call_id", "")))
        tcs = msg.get("tool_calls")
        if tcs:
            ids = tuple(
                str(tc.get("id", "")) if isinstance(tc, dict) else ""
                for tc in tcs
            )
            return ("assistant", ids)
        return (role, str(msg.get("content", "")))

    def _is_compact_summary_msg(msg: dict) -> bool:
        return (
            msg.get("role") == "system"
            and str(msg.get("content", "")).startswith("[对话历史摘要")
        )

    def _has_compact_summary(messages: list) -> bool:
        return any(_is_compact_summary_msg(m) for m in messages)

    def _expand_compact_checkpoint(reloader, cp_id: str, snapshot_msgs: list) -> Optional[list]:
        """压缩视图展开（日志式）：回退到含摘要的检查点时，向上合并最近的完整基底快照，
        找回被压缩的具体历史消息。多次压缩时基底为最早的完整快照，后续新增消息全部保留，
        天然支持多次压缩还原。
        """
        try:
            cur = reloader.restore_checkpoint(cp_id)
            if cur is None:
                return None
            cur_ts = getattr(cur, "timestamp", 0) or 0
            base_msgs = None
            for key in reloader.list_checkpoints():
                cp = reloader.restore_checkpoint(key)
                if cp is None:
                    continue
                if (getattr(cp, "timestamp", 0) or 0) >= cur_ts:
                    continue  # 只取比当前检查点更早的快照
                msgs = getattr(cp, "messages_snapshot", None) or []
                if not msgs or _has_compact_summary(msgs):
                    continue  # 跳过同样被压缩的快照，找到最早的完整基底
                base_msgs = msgs
                break
            if not base_msgs:
                return None
            base_keys = set(_message_key(m) for m in base_msgs)
            expanded = list(base_msgs)
            for m in snapshot_msgs:
                if _is_compact_summary_msg(m):
                    continue  # 摘要 system 消息不并入（前端诚实渲染不再出现摘要块）
                if _message_key(m) not in base_keys:
                    expanded.append(m)  # 压缩点之后新增的具体消息
            return expanded
        except Exception:
            pass
        return None

    def _expand_for_snapshot(reloader, snapshot_msgs: list) -> Optional[list]:
        """按当前会话消息自动定位最近压缩快照并展开为完整历史。

        前端「展开找回具体历史」按需调用：当前消息含 [对话历史摘要] 时，
        找到最近创建的含摘要对话快照作为当前检查点，复用 _expand_compact_checkpoint
        合并最早完整基底 + 压缩点之后新增消息（多次压缩递归还原）。

        缓存层：基于消息指纹（数量 + 末尾消息签名）缓存展开结果——
        前端高频轮询会话时消息未变化 → 直接命中缓存，零检查点遍历/零锁/零磁盘 IO，
        不再打断正在进行的对话（restore_checkpoint 持锁 + checkpoint_store 磁盘读是干扰源）。
        """
        cache = getattr(_expand_for_snapshot, "_cache", None)
        if cache is None:
            cache = {}
            _expand_for_snapshot._cache = cache
        try:
            if not snapshot_msgs or not _has_compact_summary(snapshot_msgs):
                return None
            # 轻量指纹：消息数 + 末尾 3 条消息签名（新消息/新压缩才会变化）
            try:
                tail = snapshot_msgs[-3:]
                sig_parts = []
                for m in tail:
                    role = m.get("role", "")
                    content = m.get("content", "")
                    if not isinstance(content, str):
                        content = str(content)
                    sig_parts.append(f"{role}:{content[:60]}")
                sig = (len(snapshot_msgs), "|".join(sig_parts))
            except Exception:
                sig = (len(snapshot_msgs), "")
            hit = cache.get(sig)
            if hit is not None:
                return hit
            cur_cp_id = None
            # list_checkpoints 已按新→旧排序（reverse=True），第一个命中即最近一次压缩快照
            for key in reloader.list_checkpoints():
                cp = reloader.restore_checkpoint(key)
                if cp is None:
                    continue
                msgs = getattr(cp, "messages_snapshot", None) or []
                if msgs and _has_compact_summary(msgs):
                    cur_cp_id = key  # 第一个命中 = 最新压缩快照，立即停止，避免遍历全部检查点刷日志/耗性能
                    break
            if not cur_cp_id:
                return None
            result = _expand_compact_checkpoint(reloader, cur_cp_id, snapshot_msgs)
            if result:
                if len(cache) >= 200:
                    cache.clear()
                cache[sig] = result
            return result
        except Exception:
            return None

    @app.post("/api/agent/messages/expand")
    async def agent_messages_expand(req: Request):
        """压缩视图展开（分页懒加载）：找回被 [对话历史摘要] 压缩掉的具体历史
        （日志式，多次压缩可还原）。

        body: { session_id, offset?, limit? } →
              { expanded: [本页消息], total: 总条数, offset, limit, has_more, count }
        前端按需分页加载，避免一次性把完整历史全量推给前端（加载性能优化）。
        无压缩/无归档返回 expanded=[] / total=0。
        """
        body = await _hub_body(req)
        session_id = body.get("session_id", "")
        try:
            offset = int(body.get("offset", 0) or 0)
            limit = int(body.get("limit", 30) or 30)
        except Exception:
            offset, limit = 0, 30
        if offset < 0:
            offset = 0
        if limit <= 0:
            limit = 30
        if limit > 300:
            limit = 300

        def _do() -> dict:
            try:
                from ..cache.context_reloader import get_context_reloader
                reloader = get_context_reloader()
            except Exception:
                return {"expanded": [], "total": 0, "offset": offset, "limit": limit, "has_more": False, "count": 0}
            if reloader is None:
                return {"expanded": [], "total": 0, "offset": offset, "limit": limit, "has_more": False, "count": 0}
            # 当前会话消息（含摘要的压缩后状态）
            snapshot = None
            try:
                agent = _peek_agent_for_session(session_id)
                if agent and agent.messages:
                    snapshot = agent.snapshot_messages()
            except Exception:
                pass
            if not snapshot:
                try:
                    store = _get_session_store()
                    if store:
                        t = store.get(session_id)
                        if t and t.messages:
                            snapshot = t.messages
                except Exception:
                    pass
            if not snapshot:
                return {"expanded": [], "total": 0, "offset": offset, "limit": limit, "has_more": False, "count": 0}
            expanded = _expand_for_snapshot(reloader, snapshot) or []
            total = len(expanded)
            page = expanded[offset:offset + limit]
            return {
                "expanded": page,
                "total": total,
                "offset": offset,
                "limit": limit,
                "has_more": (offset + len(page)) < total,
                "count": len(page),
            }

        try:
            data = await asyncio.get_event_loop().run_in_executor(None, _do)
            return ok(data=data)
        except Exception as e:
            logger.warning(f"[messages/expand] 展开失败: {e}")
            return ok(data={"expanded": [], "total": 0, "offset": offset, "limit": limit, "has_more": False, "count": 0})

    def _format_checkpoint_row(c: dict) -> dict:
        """格式化检查点记录为前端需要的格式"""
        # 映射 created_at → timestamp（前端统一用 timestamp）
        if "created_at" in c and "timestamp" not in c:
            c["timestamp"] = c["created_at"]
        if "file_count" in c and "diff_count" not in c:
            c["diff_count"] = c["file_count"]
        # SQLite 的 id → hash（兼容前端）
        if "id" in c and "hash" not in c:
            c["hash"] = c.get("checkpoint_hash", "") or str(c["id"])
        if "meta" not in c:
            c["meta"] = {
                "source": c.get("source", "auto"),
                "step": c.get("step", 0),
                "tool": c.get("tool", ""),
                "instance": c.get("session_id", "")[:8],
                "duration_ms": c.get("duration_ms"),
            }
        return c

    @app.get("/api/agent/checkpoints")
    async def agent_checkpoints(session_id: str = ""):
        """列出当前会话的检查点列表（优先从 agent，回退到独立 FileVersionDB）

        Args:
            session_id: 前端持久化的会话 ID（query parameter），用于查询正确的检查点

        支持两种模式：
        1. 有活跃 agent 实例：通过 CheckpointMiddleware 获取（更快）
        2. 无活跃 agent 实例：直接从 FileVersionDB 查询（独立读取）
        """
        # 优先尝试通过 agent 实例获取
        agent = _get_agent_for_session(app.state.workspace_dir, session_id)
        if agent is not None:
            try:
                if hasattr(agent, '_middleware_chain') and agent._middleware_chain:
                    for mw in agent._middleware_chain._middlewares:
                        from ..middleware.shadow_checkpoint import CheckpointMiddleware
                        if isinstance(mw, CheckpointMiddleware):
                            # 不再改写 mw._instance_id（会污染中间件、清空 _last_hash）。
                            # 仅当实例已绑定该会话时才用中间件，否则走 FileVersionDB 独立查询。
                            bound = (not session_id) or (getattr(agent, '_active_session_id', '') == session_id)
                            if not bound:
                                logger.debug(f"[checkpoints] instance not bound to '{session_id[:12]}' (bound={getattr(agent, '_active_session_id', '')[:12] or 'none'}), using FileVersionDB")
                                break

                            # 优先从 SQLite 获取
                            commits = mw.get_checkpoints(count=50)
                            if commits:
                                for c in commits:
                                    _format_checkpoint_row(c)
                                return ok(data={
                                    "checkpoints": commits,
                                    "available": True,
                                    "instance_id": getattr(mw, 'instance_id', ''),
                                    "version": mw.global_version,
                                })
                            # 回退到 Shadow Git
                            if mw.checkpointer:
                                cp = mw.checkpointer
                                git_commits = cp.get_commits(count=50)
                                for i, c in enumerate(git_commits):
                                    if i < len(git_commits) - 1:
                                        c["diff_count"] = cp.get_diff_count(git_commits[i + 1]["hash"], c["hash"])
                                    else:
                                        c["diff_count"] = 0
                                return ok(data={
                                    "checkpoints": git_commits,
                                    "available": True,
                                    "instance_id": getattr(mw, 'instance_id', ''),
                                    "version": 0,
                                })
            except Exception as e:
                logger.warning(f"通过 agent 获取检查点失败: {e}")

        # 回退：直接从 FileVersionDB 查询（独立于 agent）
        if session_id:
            try:
                fdb = _get_fdb_for_workspace(app.state.workspace_dir)
                if fdb:
                    commits = fdb.get_checkpoints(session_id, count=50)
                    if commits:
                        for c in commits:
                            _format_checkpoint_row(c)
                        version = fdb.global_version()
                        return ok(data={
                            "checkpoints": commits,
                            "available": True,
                            "instance_id": session_id,
                            "version": version,
                        })
            except Exception as e:
                logger.warning(f"独立获取检查点失败: {e}")

        return ok(data={"checkpoints": [], "available": False, "version": 0})

    @app.get("/api/agent/checkpoints/{commit_hash}/diff")
    async def agent_checkpoint_diff(commit_hash: str, session_id: str = ""):
        """获取指定检查点自身的增量差异（该工具调用实际修改的文件）

        优先从 agent 实例获取，回退到独立 FileVersionDB 查询。
        """
        # 优先尝试通过 agent 实例获取
        agent = _get_agent_for_session(app.state.workspace_dir, session_id)
        if agent is not None:
            try:
                if hasattr(agent, '_middleware_chain') and agent._middleware_chain:
                    for mw in agent._middleware_chain._middlewares:
                        from ..middleware.shadow_checkpoint import CheckpointMiddleware
                        if isinstance(mw, CheckpointMiddleware):
                            fdb = mw.fdb
                            if fdb:
                                try:
                                    cp_id = int(commit_hash)
                                    diff_files = fdb.get_checkpoint_diff_files(cp_id)
                                    if diff_files:
                                        if mw.checkpointer:
                                            cp_git = mw.checkpointer
                                            cp_row = fdb.get_checkpoint(cp_id)
                                            git_hash = cp_row.get("checkpoint_hash", "") if cp_row else ""
                                            for f in diff_files:
                                                if not f.get("diff"):
                                                    try:
                                                        if git_hash:
                                                            f["diff"] = cp_git.get_incremental_diff_content(git_hash, f["path"], max_lines=80)
                                                        else:
                                                            import difflib
                                                            full_path = os.path.join(mw._workspace_root, f["path"]) if not os.path.isabs(f["path"]) else f["path"]
                                                            if os.path.isfile(full_path):
                                                                new_lines = open(full_path, encoding="utf-8", errors="replace").readlines()
                                                                diff = difflib.unified_diff([], new_lines, fromfile="/dev/null", tofile=f["path"])
                                                                f["diff"] = "".join(diff)
                                                    except Exception:
                                                        pass
                                        total_adds = sum(f.get("additions", 0) for f in diff_files)
                                        total_dels = sum(f.get("deletions", 0) for f in diff_files)
                                        return ok(data={
                                            "diff": diff_files,
                                            "summary": {
                                                "additions": total_adds,
                                                "deletions": total_dels,
                                                "files_changed": len(diff_files),
                                            },
                                        })
                                except (ValueError, TypeError):
                                    pass

                            if mw.checkpointer:
                                cp = mw.checkpointer
                                diff_files = cp.get_incremental_diff(commit_hash)
                                total_adds = 0
                                total_dels = 0
                                for f in diff_files:
                                    content = cp.get_incremental_diff_content(commit_hash, f["path"], max_lines=80)
                                    f["diff"] = content
                                    total_adds += f.get("additions", 0)
                                    total_dels += f.get("deletions", 0)
                                return ok(data={
                                    "diff": diff_files,
                                    "summary": {
                                        "additions": total_adds,
                                        "deletions": total_dels,
                                        "files_changed": len(diff_files),
                                    },
                                })
            except Exception as e:
                logger.warning(f"通过 agent 获取 diff 失败: {e}")

        # 回退：直接从 FileVersionDB 查询
        if session_id:
            try:
                fdb = _get_fdb_for_workspace(app.state.workspace_dir)
                if fdb:
                    try:
                        cp_id = int(commit_hash)
                        diff_files = fdb.get_checkpoint_diff_files(cp_id)
                        if diff_files:
                            total_adds = sum(f.get("additions", 0) for f in diff_files)
                            total_dels = sum(f.get("deletions", 0) for f in diff_files)
                            return ok(data={
                                "diff": diff_files,
                                "summary": {
                                    "additions": total_adds,
                                    "deletions": total_dels,
                                    "files_changed": len(diff_files),
                                },
                            })
                    except (ValueError, TypeError):
                        pass
            except Exception as e:
                logger.warning(f"独立获取 diff 失败: {e}")

        return ok(data={"diff": [], "summary": {"additions": 0, "deletions": 0, "files_changed": 0}, "error": "无差异数据"})

    @app.get("/api/agent/checkpoints/{commit_hash}/diff-count")
    async def agent_checkpoint_diff_count(commit_hash: str, session_id: str = ""):
        """快速获取指定检查点的变更文件数（优先 SQLite）"""
        agent = _get_agent_for_session(app.state.workspace_dir, session_id)
        if agent is None:
            return ok(data={"count": 0})
        try:
            if hasattr(agent, '_middleware_chain') and agent._middleware_chain:
                for mw in agent._middleware_chain._middlewares:
                    from ..middleware.shadow_checkpoint import CheckpointMiddleware
                    if isinstance(mw, CheckpointMiddleware):
                        # 尝试 SQLite
                        fdb = mw.fdb
                        if fdb:
                            try:
                                cp_id = int(commit_hash)
                                count = fdb.get_diff_count(cp_id)
                                return ok(data={"count": count})
                            except (ValueError, TypeError):
                                pass
                        # 回退到 Shadow Git
                        if mw.checkpointer:
                            count = mw.checkpointer.get_diff_count(commit_hash)
                            return ok(data={"count": count})
        except Exception as e:
            logger.warning(f"获取 diff-count 失败: {e}")
        return ok(data={"count": 0})

    @app.post("/api/agent/checkpoints/{commit_hash}/restore")
    async def agent_checkpoint_restore(commit_hash: str, req: Request):
        """恢复到指定检查点（files / task / taskAndFiles）"""
        body = {}
        try:
            body = await req.json()
        except Exception:
            pass
        restore_type = body.get("restore_type", "files")
        session_id = body.get("session_id", "")
        agent = _get_agent_for_session(app.state.workspace_dir, session_id)
        if agent is None:
            return ok(data={"restored": False, "error": "Agent 未就绪"})
        try:
            if not (hasattr(agent, '_middleware_chain') and agent._middleware_chain):
                return ok(data={"restored": False, "error": "无中间件"})
            from ..middleware.shadow_checkpoint import CheckpointMiddleware
            from ..cache.context_reloader import get_context_reloader
            reloader = get_context_reloader()
            from ..agent import sanitize_messages

            for mw in agent._middleware_chain._middlewares:
                if isinstance(mw, CheckpointMiddleware) and mw.checkpointer:
                    cp = mw.checkpointer

                    # 解析前端传的 commit_hash：
                    # - SQLite 整数 id → 查该记录关联的 git hash（文件恢复）与 cp-* 对话快照（对话恢复）
                    # - git hash / cp-* 直接使用
                    fdb = getattr(mw, 'fdb', None)
                    git_hash = commit_hash
                    conversation_cp_id = ""
                    if str(commit_hash).isdigit():
                        try:
                            row = fdb.get_checkpoint(int(commit_hash)) if fdb else None
                            if row:
                                git_hash = row.get("checkpoint_hash", "") or ""
                                conversation_cp_id = row.get("conversation_checkpoint_id", "") or ""
                        except Exception:
                            pass
                    elif str(commit_hash).startswith("cp-"):
                        conversation_cp_id = commit_hash

                    # 恢复文件（用 git hash，仅当解析出有效 hash 时执行）
                    if restore_type in ("files", "taskAndFiles"):
                        try:
                            if git_hash:
                                cp.restore_files(git_hash)
                            else:
                                logger.warning(f"恢复文件快照失败: 无 git hash (commit_hash={commit_hash})")
                        except Exception as e:
                            logger.warning(f"恢复文件快照失败: {e}")

                    # 恢复对话历史（用 ContextReloader 的 cp-* 对话快照 id）
                    restored_msgs = None
                    if restore_type in ("task", "taskAndFiles"):
                        try:
                            if conversation_cp_id:
                                restored_msgs = reloader.rollback_to_checkpoint(conversation_cp_id)
                            else:
                                # 兜底反查（回填缺失 / 旧数据）：按 summary '#<sqlite_id>' 或 git_commit_hash 匹配
                                conversation_cp_id = _find_conversation_cp_by_ref(reloader, commit_hash)
                                if conversation_cp_id:
                                    restored_msgs = reloader.rollback_to_checkpoint(conversation_cp_id)
                            if restored_msgs:
                                # 压缩视图展开：找回被压缩的具体历史（日志式，多次压缩可还原）
                                if conversation_cp_id and _has_compact_summary(restored_msgs):
                                    expanded = _expand_compact_checkpoint(reloader, conversation_cp_id, restored_msgs)
                                    if expanded:
                                        restored_msgs = expanded
                                        logger.info(f"检查点对话已展开为完整历史: {len(restored_msgs)} 条 (cp={conversation_cp_id})")
                                restored_msgs = sanitize_messages(restored_msgs)
                                with agent._messages_lock:
                                    agent.messages = restored_msgs
                                mgr = getattr(agent, '_tool_group_mgr', None)
                                if mgr:
                                    mgr.reset_session()
                                agent._instance_tool_schemas = None
                                logger.info(f"检查点消息已恢复到 agent: {len(restored_msgs)} 条 (cp={conversation_cp_id or commit_hash})")

                                # 关键：同步到 SessionStore，确保状态机一致性
                                store = _get_session_store()
                                if store:
                                    store.update(session_id, messages=agent.snapshot_messages())
                                    logger.info(f"检查点消息已同步到 SessionStore: {session_id[:12]}")
                            else:
                                logger.warning(f"恢复对话失败: 无对话快照 (commit_hash={commit_hash}, cp={conversation_cp_id or '无'})")
                        except Exception as e:
                            logger.warning(f"恢复对话失败: {e}")

                    # 构建前端可用的消息列表
                    ui_messages = []
                    if restored_msgs:
                        for msg in restored_msgs:
                            role = msg.get("role", "")
                            content = msg.get("content", "")
                            if role == "system":
                                continue
                            if role == "tool":
                                tool_call_id = msg.get("tool_call_id", "")
                                tool_content = content if isinstance(content, str) else str(content)
                                tool_name = ""
                                checkpoint_hash = msg.get("checkpoint_hash", "")
                                for prev in reversed(ui_messages):
                                    if prev.get("role") == "assistant" and prev.get("tool_calls"):
                                        for tc in prev["tool_calls"]:
                                            tc_dict = tc if isinstance(tc, dict) else {}
                                            if tc_dict.get("id") == tool_call_id:
                                                tool_name = tc_dict.get("function", {}).get("name", "")
                                                if not checkpoint_hash:
                                                    checkpoint_hash = tc_dict.get("checkpoint_hash", "")
                                                break
                                        break
                                ui_messages.append({
                                    "role": "tool",
                                    "tool_call_id": tool_call_id,
                                    "tool_name": tool_name,
                                    "content": tool_content,
                                    "checkpoint_hash": checkpoint_hash,
                                })
                            elif role == "assistant":
                                entry = {
                                    "role": "assistant",
                                    "content": content if isinstance(content, str) else str(content),
                                }
                                if msg.get("tool_calls"):
                                    entry["tool_calls"] = msg["tool_calls"]
                                ui_messages.append(entry)
                            elif role == "user" and content:
                                ui_messages.append({
                                    "role": "user",
                                    "content": content if isinstance(content, str) else str(content),
                                })

                    return ok(data={
                        "restored": True,
                        "restore_type": restore_type,
                        "commit_hash": commit_hash,
                        "restored_messages": len(restored_msgs or []),
                        "ui_messages": ui_messages,
                    })
        except Exception as e:
            logger.warning(f"恢复检查点失败: {e}")
        return ok(data={"restored": False, "error": str(e)})

    @app.post("/api/agent/migrate-checkpoint-sessions")
    async def agent_migrate_checkpoint_sessions(force: bool = False, session_id: str = ""):
        """一次性迁移旧的检查点会话到标准格式
        
        迁移策略：
        1. 如果旧 cp- 会话在 SessionStore 中有完整消息 → 重命名文件为 sess_* 格式
        2. 如果只有 FileVersionDB 中的检查点 → 从检查点重建消息，创建新会话
        3. 不删除旧文件（安全起见），迁移后旧文件保留为备份
        """
        store = _get_session_store()
        if not store:
            return err(msg="SessionStore 不可用")
        
        fdb = _get_fdb_for_workspace(app.state.workspace_dir)
        if not fdb:
            return ok(data={"migrated": 0, "skipped": 0, "filtered": 0, "message": "无检查点数据"})
        
        import re
        import uuid
        import json
        import shutil
        from pathlib import Path
        
        try:
            conn = fdb._connect()
            if session_id:
                rows = conn.execute(
                    "SELECT DISTINCT session_id FROM checkpoints WHERE session_id = ?",
                    (session_id,)
                ).fetchall()
            else:
                rows = conn.execute(
                    "SELECT DISTINCT session_id FROM checkpoints WHERE session_id != ''"
                ).fetchall()
            
            migrated = 0
            skipped = 0
            filtered = 0
            migration_map = {}
            
            store_dir = Path(store.store_dir)
            
            for row in rows:
                old_sid = row["session_id"]
                
                # 过滤异常会话 ID
                if re.match(r'^[0-9a-f]{8}$', old_sid):
                    filtered += 1
                    continue
                
                # 确定新的会话 ID
                if old_sid.startswith("cp-"):
                    new_sid = f"sess_{uuid.uuid4().hex[:12]}"
                elif old_sid.startswith("sess_") or re.match(r'^[0-9a-f]{8}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{12}$', old_sid):
                    # 已经是标准格式，跳过
                    skipped += 1
                    continue
                else:
                    new_sid = old_sid
                
                # 检查旧文件是否存在
                old_file = store_dir / f"{old_sid}.json"
                new_file = store_dir / f"{new_sid}.json"
                
                if new_file.exists():
                    skipped += 1
                    continue
                
                # 如果旧文件存在，直接重命名并修改 ID
                if old_file.exists():
                    try:
                        # 读取旧文件
                        with open(old_file, 'r', encoding='utf-8') as f:
                            data = json.load(f)
                        
                        # 更新 ID
                        data["id"] = new_sid
                        # 更新 metadata
                        if "metadata" not in data:
                            data["metadata"] = {}
                        data["metadata"]["source"] = "checkpoint_migrated"
                        data["metadata"]["original_id"] = old_sid
                        data["metadata"]["migration_type"] = "one_time"
                        
                        # 写入新文件
                        with open(new_file, 'w', encoding='utf-8') as f:
                            json.dump(data, f, ensure_ascii=False, indent=2)
                        
                        # 删除旧文件
                        old_file.unlink()
                        
                        # 更新缓存
                        if old_sid in store._cache:
                            record = store._cache.pop(old_sid)
                            record.id = new_sid
                            record.metadata = data["metadata"]
                            store._cache[new_sid] = record
                        
                        migration_map[old_sid] = new_sid
                        migrated += 1
                        logger.info(f"[迁移] 重命名会话 {old_sid[:12]} → {new_sid[:12]}")
                    except Exception as e:
                        logger.warning(f"[迁移] 重命名失败 {old_sid[:12]}: {e}")
                        skipped += 1
                else:
                    # 旧文件不存在，从检查点重建
                    checkpoints = fdb.get_checkpoints(old_sid, count=100)
                    
                    # 从检查点构建消息
                    messages = []
                    for cp in reversed(checkpoints):
                        tool_name = cp.get("tool", "")
                        step = cp.get("step", 0)
                        if tool_name:
                            messages.append({
                                "role": "assistant",
                                "content": "",
                                "tool_calls": [{
                                    "id": f"cp-{cp.get('id', '')}",
                                    "type": "function",
                                    "function": {"name": tool_name, "arguments": "{}"}
                                }]
                            })
                            messages.append({
                                "role": "tool",
                                "tool_call_id": f"cp-{cp.get('id', '')}",
                                "content": f"检查点 #{step}: {cp.get('file_count', 0)} 个文件变更",
                                "checkpoint_hash": cp.get("checkpoint_hash", ""),
                            })
                    
                    if not messages:
                        messages.append({"role": "user", "content": "历史会话"})
                    
                    # 生成会话名称
                    name = _extract_session_name(messages, checkpoints, old_sid)
                    
                    # 创建新会话
                    try:
                        record = store.create_with_id(
                            session_id=new_sid,
                            name=name,
                            metadata={
                                "source": "checkpoint_migrated",
                                "original_id": old_sid,
                                "migration_type": "one_time"
                            }
                        )
                        store.update(record.id, messages=messages)
                        migration_map[old_sid] = new_sid
                        migrated += 1
                        logger.info(f"[迁移] 重建会话 {old_sid[:12]} → {new_sid[:12]}")
                    except Exception as e:
                        logger.warning(f"[迁移] 重建失败 {old_sid[:12]}: {e}")
                        skipped += 1
            
            return ok(data={
                "migrated": migrated,
                "skipped": skipped,
                "filtered": filtered,
                "total": len(rows),
                "migration_map": migration_map,
                "message": f"已迁移 {migrated} 个会话，跳过 {skipped} 个已存在，过滤 {filtered} 个异常会话。",
            })
        except Exception as e:
            logger.error(f"Migration failed: {e}")
            return err(msg=f"迁移失败: {str(e)}")

    @app.get("/api/agent/checkpoint-sessions")
    async def agent_list_checkpoint_sessions():
        """列出 FileVersionDB 中所有唯一的 session_id（用于迁移参考）"""
        fdb = _get_fdb_for_workspace(app.state.workspace_dir)
        if not fdb:
            return ok(data={"sessions": [], "total": 0})
        
        try:
            conn = fdb._connect()
            rows = conn.execute(
                "SELECT session_id, COUNT(*) as cp_count FROM checkpoints WHERE session_id != '' GROUP BY session_id ORDER BY cp_count DESC"
            ).fetchall()
            
            sessions = []
            store = _get_session_store()
            
            for row in rows:
                session_id = row["session_id"]
                cp_count = row["cp_count"]
                
                # 检查是否已在 SessionStore 中
                in_session_store = False
                if store:
                    existing = store.get(session_id)
                    in_session_store = existing is not None
                
                sessions.append({
                    "session_id": session_id,
                    "checkpoint_count": cp_count,
                    "in_session_store": in_session_store,
                })
            
            return ok(data={
                "sessions": sessions,
                "total": len(sessions),
                "migrated_count": sum(1 for s in sessions if s["in_session_store"]),
                "pending_count": sum(1 for s in sessions if not s["in_session_store"]),
            })
        except Exception as e:
            return err(msg=f"查询失败: {str(e)}")

    @app.post("/api/agent/restore-checkpoint-session/{session_id}")
    async def agent_restore_checkpoint_session(session_id: str):
        """恢复单个检查点会话到 SessionStore
        
        将旧的 cp-* 会话或其他格式的检查点会话恢复为可用会话，
        保留其检查点数据作为对话历史。
        """
        store = _get_session_store()
        if not store:
            return err(msg="SessionStore 不可用")
        
        fdb = _get_fdb_for_workspace(app.state.workspace_dir)
        if not fdb:
            return err(msg="无检查点数据")
        
        try:
            # 检查是否已在 SessionStore 中
            existing = store.get(session_id)
            if existing:
                return ok(data={"restored": False, "reason": "already_exists", "session_id": session_id})
            
            # 从 FileVersionDB 获取该会话的检查点
            checkpoints = fdb.get_checkpoints(session_id, count=100)
            if not checkpoints:
                return err(msg=f"会话 {session_id} 无检查点数据")
            
            # 构建消息列表
            messages = _build_messages_from_checkpoints(checkpoints, fallback_content=f"历史会话 {session_id[:8]}")
            
            # 使用 _extract_session_name 生成会话名称
            name = _extract_session_name(messages, checkpoints, session_id)
            
            # 创建 SessionStore 记录
            record = store.create_with_id(
                session_id=session_id,
                name=name,
                metadata={"source": "checkpoint_restored", "original_id": session_id}
            )
            store.update(record.id, messages=messages)
            
            return ok(data={
                "restored": True,
                "session_id": session_id,
                "name": name,
                "checkpoint_count": len(checkpoints),
                "message_count": len(messages),
            })
        except Exception as e:
            logger.error(f"Restore checkpoint session failed: {e}")
            return err(msg=f"恢复失败: {str(e)}")

    # ─── 多实例集群 API ──────────────────────────────────────

    def _scan_instances(port: int):
        """同步扫描本机其他 TS2 实例"""
        instances = []
        for p in range(6906, 6916):
            if p == port:
                continue
            try:
                import urllib.request
                req = urllib.request.urlopen(
                    f"http://127.0.0.1:{p}/api/system/version",
                    data=b'{}',
                    timeout=2,
                )
                if req.status == 200:
                    data = _json_loads(req.read().decode('utf-8'))
                    instances.append({
                        "port": p,
                        "url": f"http://127.0.0.1:{p}",
                        "version": data.get("data", {}).get("version", "unknown"),
                        "local_ip": data.get("data", {}).get("local_ip", ""),
                    })
            except Exception:
                continue
        return instances

    def _fetch_remote_file(remote_url: str, file_path: str):
        """同步从远端实例获取文件内容"""
        import urllib.request
        # 先获取文件元数据
        meta_req = urllib.request.Request(
            f"{remote_url}/api/file/getFile",
            data=_json_dumps_compact({"path": file_path}).encode('utf-8'),
            headers={"Content-Type": "application/json"},
            method="POST",
        )
        with urllib.request.urlopen(meta_req, timeout=30) as resp:
            result = _json_loads(resp.read().decode('utf-8'))
        if result.get("code") != 0:
            return None, result.get("msg", "远端获取失败")
        file_data = result.get("data", {})
        content = file_data.get("content", "")
        entry = file_data.get("entry", {})
        return {"content": content, "entry": entry}, None

    def _fetch_remote_dir(remote_url: str, dir_path: str):
        """同步从远端实例获取目录列表"""
        import urllib.request
        req = urllib.request.Request(
            f"{remote_url}/api/file/readDir",
            data=_json_dumps_compact({"path": dir_path}).encode('utf-8'),
            headers={"Content-Type": "application/json"},
            method="POST",
        )
        with urllib.request.urlopen(req, timeout=10) as resp:
            result = _json_loads(resp.read().decode('utf-8'))
        if result.get("code") != 0:
            return None, result.get("msg", "远端获取失败")
        return result.get("data", []), None

    def _fetch_remote_search(remote_url: str, query: str, subdir: str = ""):
        """同步从远端实例搜索文件"""
        import urllib.request
        req = urllib.request.Request(
            f"{remote_url}/api/file/search",
            data=_json_dumps_compact({"query": query, "subdir": subdir}).encode('utf-8'),
            headers={"Content-Type": "application/json"},
            method="POST",
        )
        with urllib.request.urlopen(req, timeout=15) as resp:
            result = _json_loads(resp.read().decode('utf-8'))
        if result.get("code") != 0:
            return None, result.get("msg", "远端搜索失败")
        return result.get("data", []), None

    def _transfer_file_from_remote(remote_url: str, remote_path: str, local_path: str):
        """同步从远端实例传输文件到本地"""
        file_data, err = _fetch_remote_file(remote_url, remote_path)
        if err:
            return None, err
        content = file_data["content"]
        # 写入本地
        engine: FileSyncEngine = app.state.sync_engine
        entry = engine.put_file(local_path, content)
        if entry is None:
            return None, "本地写入失败"
        return entry.to_dict(), None

    def _batch_transfer_from_remote(remote_url: str, file_pairs: list):
        """同步批量从远端传输文件到本地

        file_pairs: [{"remote_path": "...", "local_path": "..."}, ...]
        """
        results = []
        ok_count = 0
        fail_count = 0
        for pair in file_pairs:
            remote_path = pair.get("remote_path", "")
            local_path = pair.get("local_path", remote_path)  # 默认同路径
            entry, err = _transfer_file_from_remote(remote_url, remote_path, local_path)
            if err:
                fail_count += 1
                results.append({"path": remote_path, "success": False, "error": err})
            else:
                ok_count += 1
                results.append({"path": remote_path, "success": True, "entry": entry})
        return {"ok": ok_count, "fail": fail_count, "results": results}, None

    @app.get("/api/cluster/instances")
    async def cluster_instances():
        """发现本机其他 TS2 实例"""
        instances = await _run_data(_scan_instances, app.state.port)
        # 加入自身信息
        self_info = {
            "port": app.state.port,
            "url": f"http://127.0.0.1:{app.state.port}",
            "version": "1.0.0",
            "local_ip": app.state.local_ip,
            "self": True,
            "workspace": str(app.state.workspace_dir),
        }
        return ok(data={"self": self_info, "peers": instances})

    @app.post("/api/cluster/remote/readDir")
    async def cluster_remote_read_dir(req: Request):
        """从远端实例读取目录"""
        body = await req.json()
        remote_url = body.get("remote_url", "")
        dir_path = body.get("path", "")
        if not remote_url:
            return err(msg="缺少 remote_url")
        result, error = await _run_data(_fetch_remote_dir, remote_url, dir_path)
        if error:
            return err(msg=error)
        return ok(data=result)

    @app.post("/api/cluster/remote/search")
    async def cluster_remote_search(req: Request):
        """从远端实例搜索文件"""
        body = await req.json()
        remote_url = body.get("remote_url", "")
        query = body.get("query", "")
        subdir = body.get("subdir", "")
        if not remote_url or not query:
            return err(msg="缺少 remote_url 或 query")
        result, error = await _run_data(_fetch_remote_search, remote_url, query, subdir)
        if error:
            return err(msg=error)
        return ok(data=result)

    @app.post("/api/cluster/transfer")
    async def cluster_transfer(req: Request):
        """从远端实例传输文件到本地（单个文件）"""
        body = await req.json()
        remote_url = body.get("remote_url", "")
        remote_path = body.get("remote_path", "")
        local_path = body.get("local_path", remote_path)
        if not remote_url or not remote_path:
            return err(msg="缺少 remote_url 或 remote_path")
        result, error = await _run_data(_transfer_file_from_remote, remote_url, remote_path, local_path)
        if error:
            return err(msg=error)
        return ok(data=result)

    @app.post("/api/cluster/transfer/batch")
    async def cluster_transfer_batch(req: Request):
        """从远端实例批量传输文件到本地"""
        body = await req.json()
        remote_url = body.get("remote_url", "")
        file_pairs = body.get("files", [])
        if not remote_url or not file_pairs:
            return err(msg="缺少 remote_url 或 files")
        result, error = await _run_data(_batch_transfer_from_remote, remote_url, file_pairs)
        if error:
            return err(msg=error)
        return ok(data=result)

    # ─── Health Check ──────────────────────────────────────

    @app.get("/api/ping")
    @app.head("/api/ping")
    async def ping():
        return ok(data={"status": "ok"})

    # ─── Mobile Bootstrap API ──────────────────────────────

    @app.get("/api/mobile/bootstrap")
    async def mobile_bootstrap():
        """移动端启动引导：一次请求获取所有模块数据"""
        workspace_dir = app.state.workspace_dir

        # 并行获取各模块数据
        tasks_data, courses_data, bookmarks_data, projects_data, agent_data, push_data = await asyncio.gather(
            _run_data(_read_tasks_data, workspace_dir),
            _run_data(_read_courses_data, workspace_dir),
            _run_data(_read_bookmarks_data, workspace_dir),
            _run_data(_read_projects_data, workspace_dir),
            _run_data(_read_agent_status, workspace_dir),
            _run_push(_get_push_dashboard, workspace_dir),
            return_exceptions=True,
        )

        return ok(data={
            "tasks": tasks_data if not isinstance(tasks_data, Exception) else [],
            "courses": courses_data if not isinstance(courses_data, Exception) else [],
            "bookmarks": bookmarks_data if not isinstance(bookmarks_data, Exception) else [],
            "projects": projects_data if not isinstance(projects_data, Exception) else [],
            "agent": agent_data if not isinstance(agent_data, Exception) else {"available": False},
            "push": push_data if not isinstance(push_data, Exception) else {},
            "server": {
                "version": "1.0.0",
                "local_ip": app.state.local_ip,
                "port": app.state.port,
                "uptime": time.time() - app.state.start_time,
            },
        })

    # ─── 启动/关闭事件 ──────────────────────────────────────

    @app.on_event("startup")
    async def on_startup():
        engine: FileSyncEngine = app.state.sync_engine
        # 捕获主事件循环引用（用于线程池安全的事件通知）
        try:
            engine._main_loop = asyncio.get_running_loop()
        except RuntimeError:
            engine._main_loop = asyncio.get_event_loop()
        # 设置 server_tools 的主事件循环引用
        try:
            from .server_tools import set_main_loop
            set_main_loop(asyncio.get_running_loop())
        except Exception:
            pass
        # Run scan in background to avoid blocking startup
        import threading
        def _bg_scan():
            try:
                engine.scan_file_tree()
            except Exception as e:
                logger.warning(f"Background scan error: {e}")
        threading.Thread(target=_bg_scan, daemon=True).start()
        local_ip = app.state.local_ip
        logger.info(f"TS2 Server started: {host}:{port}, workspace={workspace_dir}")
        logger.info(f"LAN access: http://{local_ip}:{port}")

        # 启动定时推送任务
        async def _periodic_push():
            """每5分钟推送一次 dashboard 数据到所有 WebSocket 客户端"""
            while True:
                await asyncio.sleep(300)  # 5分钟
                try:
                    ws_mgr = app.state.ws_manager
                    if ws_mgr and ws_mgr.get_session_count() > 0:
                        push_data = await _run_push(_get_push_dashboard, app.state.workspace_dir)
                        await ws_mgr.broadcast("pushDashboard", data=push_data)
                except Exception as e:
                    logger.warning(f"Periodic push error: {e}")

        asyncio.create_task(_periodic_push())

    @app.on_event("shutdown")
    async def on_shutdown():
        engine: FileSyncEngine = app.state.sync_engine
        await engine.stop_watching()
        # 清理线程池，防止僵尸线程阻止进程退出
        _file_executor.shutdown(wait=False)
        _data_executor.shutdown(wait=False)
        _agent_executor.shutdown(wait=False)
        _push_executor.shutdown(wait=False)
        # 清理 Agent 实例池（释放 LLM 连接等资源）
        with _agent_pool_lock:
            for sid, agent in list(_agent_pool.items()):
                try:
                    agent.cancel()
                except Exception:
                    pass
            _agent_pool.clear()
        logger.info("TS2 Server shutdown complete")

    # ─── PDF 智能阅读 API ──────────────────────────────────────

    from .pdf_service import extract_pdf_text, index_pdf, query_pdf

    class PdfExtractRequest(BaseModel):
        file_path: str

    class PdfIndexRequest(BaseModel):
        file_path: str

    class PdfQueryRequest(BaseModel):
        query: str
        top_k: int = 4

    @app.post("/api/pdf/extract")
    async def pdf_extract(req: PdfExtractRequest, request: Request = None):
        """提取 PDF 文本内容"""
        try:
            abs_path = sync_engine._absolute_path(req.file_path)
        except ValueError:
            if not _is_resource_path(req.file_path, app.state.workspace_dir):
                return err(msg="路径不在允许的读取目录中")
            if not await check_auth(request, load_api_config()):
                return err(msg="未授权，请提供有效 token 或授权码")
            abs_path = Path(req.file_path).resolve()
        if not os.path.isfile(abs_path):
            return err(msg=f"文件不存在: {req.file_path}")
        try:
            result = await asyncio.get_event_loop().run_in_executor(
                None, extract_pdf_text, str(abs_path)
            )
            return ok(data=result)
        except Exception as e:
            return err(msg=f"PDF 提取失败: {e}")

    @app.post("/api/pdf/index")
    async def pdf_index(req: PdfIndexRequest, request: Request = None):
        """将 PDF 索引到 RAG 向量库"""
        try:
            abs_path = sync_engine._absolute_path(req.file_path)
        except ValueError:
            if not _is_resource_path(req.file_path, app.state.workspace_dir):
                return err(msg="路径不在允许的读取目录中")
            if not await check_auth(request, load_api_config()):
                return err(msg="未授权，请提供有效 token 或授权码")
            abs_path = Path(req.file_path).resolve()
        if not os.path.isfile(abs_path):
            return err(msg=f"文件不存在: {req.file_path}")
        try:
            result = await asyncio.get_event_loop().run_in_executor(
                None, index_pdf, str(abs_path), app.state.workspace_dir
            )
            return ok(data=result)
        except Exception as e:
            return err(msg=f"PDF 索引失败: {e}")

    @app.post("/api/pdf/query")
    async def pdf_query(req: PdfQueryRequest):
        """对已索引 PDF 进行 RAG 查询"""
        try:
            result = await asyncio.get_event_loop().run_in_executor(
                None, query_pdf, req.query, app.state.workspace_dir, req.top_k
            )
            return ok(data=result)
        except Exception as e:
            return err(msg=f"PDF 查询失败: {e}")

    @app.post("/api/pdf/chat")
    async def pdf_chat(req: AgentChatRequest):
        """PDF AI 对话 — 先 RAG 检索，再注入 Agent 上下文"""
        agent = _get_web_agent(app.state.workspace_dir)
        if agent is None:
            return ok(data={"reply": "Agent 未初始化", "source": "error"})

        # RAG 检索相关内容
        contexts = []
        try:
            rag_result = await asyncio.get_event_loop().run_in_executor(
                None, query_pdf, req.message, app.state.workspace_dir, 4
            )
            contexts = rag_result.get("results", [])
        except Exception:
            pass

        # 构造带 PDF 上下文的消息
        if contexts:
            context_text = "\n\n".join(
                f"[{c['file_name']} 第{c['page']}页]\n{c['content']}"
                for c in contexts
            )
            enhanced_message = (
                f"以下是与用户问题相关的 PDF 内容：\n\n{context_text}\n\n"
                f"用户问题：{req.message}\n\n请基于以上 PDF 内容回答问题。"
            )
        else:
            enhanced_message = req.message

        try:
            reply = await asyncio.wait_for(
                _run_agent(agent.chat, enhanced_message),
                timeout=300.0
            )
            return ok(data={"reply": reply, "source": "pdf_rag", "contexts": contexts})
        except asyncio.TimeoutError:
            return ok(data={"reply": "处理超时", "source": "timeout"})
        except Exception as e:
            return ok(data={"reply": f"处理出错：{e}", "source": "error"})

    # ─── 主题清单（动态扫描 static/theme 目录，替代前端硬编码）──────────
    def _scan_theme_manifest() -> list:
        """扫描 static/theme 下每个子目录，读取 theme.json + 入口 CSS，返回主题清单"""
        from urllib.parse import quote
        theme_root = static_dir / "theme"
        themes = []
        if not theme_root.exists():
            return themes
        for d in sorted(theme_root.iterdir()):
            if not d.is_dir() or d.name.startswith('.'):
                continue
            meta = {
                "id": d.name,
                "dir": d.name,
                "name": d.name,
                "displayName": d.name,
                "author": "",
                "version": "",
                "modes": [],
                "description": "",
                "css": "",
                "cssUrl": "",
                "iconUrl": "",
                "previewUrl": "",
            }
            tj = d / "theme.json"
            if tj.exists():
                try:
                    tj_data = json.loads(tj.read_text(encoding="utf-8"))
                    meta["name"] = tj_data.get("name") or d.name
                    disp = tj_data.get("displayName") or {}
                    meta["displayName"] = (disp.get("zh_CN") or disp.get("default") or d.name)
                    meta["author"] = tj_data.get("author") or ""
                    meta["version"] = tj_data.get("version") or ""
                    meta["modes"] = tj_data.get("modes") or []
                    desc = tj_data.get("description") or {}
                    meta["description"] = (desc.get("zh_CN") or desc.get("default") or "")
                except Exception:
                    pass
            # 入口 CSS：tsw2.css（TWS2 适配层）> theme.css > {目录名}.css > 目录内第一个 .css
            css_rel = None
            for cand in (d / "tsw2.css", d / "theme.css", d / (d.name + ".css")):
                if cand.exists() and cand.is_file():
                    css_rel = cand.name
                    break
            if css_rel is None:
                for f in sorted(d.iterdir()):
                    if f.is_file() and f.suffix.lower() == ".css":
                        css_rel = f.name
                        break
            if css_rel:
                meta["css"] = css_rel
                rel_url = "static/theme/" + d.name + "/" + css_rel
                meta["cssUrl"] = "/" + quote(rel_url.replace("\\", "/"), safe="/")
            for asset, key in (("icon.png", "iconUrl"), ("preview.png", "previewUrl")):
                if (d / asset).exists():
                    rel_url = "static/theme/" + d.name + "/" + asset
                    meta[key] = "/" + quote(rel_url.replace("\\", "/"), safe="/")
            themes.append(meta)
        return themes

    # ─── /api/themes 主题清单 API（异步刷新，保证目录最新）────
    @app.get("/api/themes")
    async def list_themes():
        return ok(data={"themes": _scan_theme_manifest()})

    # ─── /static/theme/manifest.js 主题清单注入（同步 <script> 加载，首帧兜底）────
    # 必须注册在 /static mount 之前：显式路由优先于 mount 前缀匹配。
    @app.get("/static/theme/manifest.js")
    async def theme_manifest_js():
        payload = json.dumps({"themes": _scan_theme_manifest()}, ensure_ascii=False)
        body = "window.__TS2_THEME_MANIFEST = " + payload + ";\n"
        return Response(content=body, media_type="application/javascript")

    # ─── /vditor/ 路由（APP 端离线时 cdn 为 /vditor，映射到 static/vditor/）────
    @app.get("/vditor/{file_path:path}")
    async def vditor_fallback(file_path: str):
        vditor_file = static_dir / "vditor" / file_path
        if vditor_file.exists() and vditor_file.is_file():
            return FileResponse(str(vditor_file))
        return HTMLResponse(content="", status_code=404)

    # ─── 静态文件挂载（必须在 catch-all 路由之前）────────────────
    if static_dir.exists():
        app.mount("/static", StaticFiles(directory=str(static_dir)), name="static")

    # ─── 静态资源禁用浏览器强缓存（开发期避免加载到旧的 app.js / index.html）────
    @app.middleware("http")
    async def _no_cache_static(request: Request, call_next):
        response = await call_next(request)
        if request.url.path.startswith("/static"):
            response.headers["Cache-Control"] = "no-cache"
            # 修正 Windows mimetypes 对模块脚本的 MIME 误判（collab loro_wasm 等）
            _p = request.url.path
            if _p.endswith(".wasm"):
                response.headers["Content-Type"] = "application/wasm"
            elif _p.endswith(".js") and (response.headers.get("content-type") or "").startswith("application/json"):
                response.headers["Content-Type"] = "text/javascript; charset=utf-8"
        return response

    # ─── SaberSystem 路由挂载（必须在 SPA fallback 之前）──────────────
    try:
        from mcp.server.saber.api import create_saber_router
        app.include_router(create_saber_router())
    except Exception as _e:
        # SaberSystem 模块加载失败不应阻断主应用启动
        pass

    # ─── 笔记分析 & 间隔重复 API（必须在 SPA fallback 之前）──────────
    @app.post("/api/data/notes/analyze")
    async def notes_analyze(req: NoteAnalyzeRequest):
        """分析笔记：解析 :::env 环境，自动分类，提取三元组"""
        if not _note_analyzer_ready:
            return err(msg="note_analyzer 模块未加载")
        full_path = _resolve_ws_path(req.path)
        if not os.path.exists(full_path):
            return err(msg="笔记文件不存在: " + req.path)
        try:
            from pathlib import Path as _P
            elements = await _run_data(NoteAnalyzer.extract_from_file, _P(full_path), req.course_id, req.lesson_num)
            NoteAnalyzer.auto_classify_elements(elements)
            NoteAnalyzer.auto_extract_triples(elements)
            by_type = {}
            by_category = {}
            for e in elements:
                by_type[e.elem_type] = by_type.get(e.elem_type, 0) + 1
                for c in (e.categories or []):
                    by_category[c] = by_category.get(c, 0) + 1
            return ok(data={
                "elements": [e.to_dict() for e in elements],
                "stats": {"total": len(elements), "by_type": by_type, "by_category": by_category}
            })
        except Exception as e:
            return err(msg=f"分析失败: {e}")

    @app.post("/api/data/notes/fill-blank")
    async def notes_fill_blank(req: FillBlankRequest):
        """生成填空题"""
        if not _note_analyzer_ready:
            return err(msg="note_analyzer 模块未加载")
        try:
            masked, answers, full = NoteAnalyzer.generate_fill_in_blank(req.content)
            # 将 ___ 替换为 ◆BLANK_N◆ 占位符，避免破坏 LaTeX/Markdown 语法
            import re
            _counter = [0]
            def _ph(m):
                r = f"◆BLANK{_counter[0]}◆"
                _counter[0] += 1
                return r
            masked = re.sub(r'___', _ph, masked)
            # 修复被 ◆BLANK 破坏的 LaTeX 公式：移除孤立的 $ 定界符
            # 例如 $E=◆BLANK_0◆$ → E=◆BLANK_0◆（占位符作为普通文本）
            masked = re.sub(r'\$\$([^$]*◆BLANK[^$]*)\$\$', r'\1', masked)
            masked = re.sub(r'\$([^$]*◆BLANK[^$]*)\$', r'\1', masked)
            masked = re.sub(r'\\\[([^\\]*◆BLANK[^\\]*)\\\]', r'\1', masked)
            masked = re.sub(r'\\\(([^\\]*◆BLANK[^\\]*)\\\)', r'\1', masked)
            return ok(data={"masked_text": masked, "answers": answers, "full_text": full})
        except Exception as e:
            return err(msg=f"生成填空失败: {e}")

    @app.post("/api/data/notes/pairs")
    async def notes_pairs(req: NotePairsRequest):
        """获取笔记中定理-证明 / 问题-解答配对"""
        if not _note_analyzer_ready:
            return err(msg="note_analyzer 模块未加载")
        full_path = _resolve_ws_path(req.path)
        if not os.path.exists(full_path):
            return err(msg="笔记文件不存在: " + req.path)
        try:
            from pathlib import Path as _P
            elements = await _run_data(NoteAnalyzer.extract_from_file, _P(full_path), req.course_id, req.lesson_num)
            if req.pair_type == "problem_solution":
                pairs = NoteAnalyzer.find_problem_solution_pair(elements)
            else:
                pairs = NoteAnalyzer.find_theorem_proof_pair(elements)
            # 序列化（包含 raw_start/raw_end 供后续保存定位）
            pair_list = []
            for primary, secondary in pairs:
                pair_list.append({
                    "primary": primary.to_dict(include_raw_pos=True),
                    "secondary": secondary.to_dict(include_raw_pos=True) if secondary else None,
                })
            return ok(data={"pairs": pair_list, "total": len(pair_list)})
        except Exception as e:
            return err(msg=f"获取配对失败: {e}")

    @app.post("/api/data/notes/pairs/save")
    async def notes_pairs_save(req: NotePairSaveRequest):
        """重读写入定理/证明 或 问题/解答到笔记文件"""
        if not _note_analyzer_ready:
            return err(msg="note_analyzer 模块未加载")
        full_path = _resolve_ws_path(req.path)
        if not os.path.exists(full_path):
            return err(msg="笔记文件不存在: " + req.path)
        try:
            from pathlib import Path as _P
            path = _P(full_path)
            note_content = path.read_text(encoding="utf-8")

            # 决定环境类型标签
            if req.pair_type == "problem_solution":
                primary_tag = "problem"
                secondary_tag = "solution"
            else:
                # 定理类：用原始 elem_type（theorem/corollary/lemma/...）
                primary_tag = req.primary_elem_type or "theorem"
                secondary_tag = "proof"

            primary_env = f"::: {primary_tag}\n{req.primary_content}\n:::"
            has_secondary = bool(req.secondary_content and req.secondary_content.strip())

            start_replace = req.primary_raw_start
            if req.secondary_raw_start >= 0 and req.secondary_raw_end >= 0:
                # 已有证明/解答段落
                if has_secondary:
                    # 替换从 primary 开始到 secondary 结束的整段
                    end_replace = req.secondary_raw_end
                    combined_env = primary_env + "\n\n" + f"::: {secondary_tag}\n{req.secondary_content}\n:::"
                    new_content = note_content[:start_replace] + combined_env + note_content[end_replace:]
                else:
                    # secondary 为空：只替换 primary 部分，保留原 secondary 不变
                    end_replace = req.primary_raw_end
                    new_content = note_content[:start_replace] + primary_env + note_content[end_replace:]
            else:
                # 无旧证明/解答
                if has_secondary:
                    # 替换 primary 自身，然后在其后插入 secondary
                    end_replace = req.primary_raw_end
                    temp_content = note_content[:start_replace] + primary_env + note_content[end_replace:]
                    insert_pos = start_replace + len(primary_env)
                    new_content = temp_content[:insert_pos] + "\n\n" + f"::: {secondary_tag}\n{req.secondary_content}\n:::" + temp_content[insert_pos:]
                else:
                    # 只替换 primary
                    end_replace = req.primary_raw_end
                    new_content = note_content[:start_replace] + primary_env + note_content[end_replace:]

            path.write_text(new_content, encoding="utf-8")
            return ok(data={"saved": True, "length": len(new_content)})
        except Exception as e:
            return err(msg=f"保存失败: {e}")

    # ─── 课程级聚合分析（持久化 + 增量更新） ─────────────────────────
    def _course_analysis_path(workspace_dir: str, course_id: str) -> str:
        """课程聚合分析的持久化路径"""
        safe_id = "".join(c if c.isalnum() or c in "-_" else "_" for c in course_id)[:80] or "unknown"
        return os.path.join(workspace_dir, "data", f"course_analysis_{safe_id}.json")

    def _resolve_course_note_path(workspace_dir: str, course: dict, lesson: dict) -> str:
        """复刻前端 getNotePath 逻辑生成笔记路径"""
        title = course.get("course_title") or course.get("title") or "unknown"
        safe_title = "".join(c if c not in '\\/:*?"<>|' else "_" for c in title)[:40]
        lnum = lesson.get("lesson_number")
        if lnum is None:
            lnum = lesson.get("number")
        if lnum is None:
            lnum = 1
        raw_title = lesson.get("lesson_title") or lesson.get("title") or ""
        import re as _re
        ltitle = _re.sub(r'^L\d+[_\s]*', '', raw_title, flags=_re.IGNORECASE)
        safe_ltitle = "".join(c if c not in '\\/:*?"<>|' else "_" for c in ltitle)[:30]
        num_str = str(lnum).zfill(2)
        fn = f"L{num_str}_{safe_ltitle}.Rmd" if safe_ltitle else f"L{num_str}.Rmd"
        return os.path.join(workspace_dir, "Notes", safe_title, fn)

    def _scan_course_note_paths(workspace_dir: str, course: dict) -> list:
        """扫描课程目录下所有 L*.Rmd 笔记，返回 [{path, lesson_number, lesson_title}]"""
        title = course.get("course_title") or course.get("title") or "unknown"
        safe_title = "".join(c if c not in '\\/:*?"<>|' else "_" for c in title)[:40]
        notes_dir = os.path.join(workspace_dir, "Notes", safe_title)
        result = []
        if not os.path.isdir(notes_dir):
            return result
        import re as _re
        for name in sorted(os.listdir(notes_dir)):
            if not name.lower().endswith(".rmd"):
                continue
            m = _re.match(r'L(\d+)', name, _re.IGNORECASE)
            if not m:
                continue
            lnum = int(m.group(1))
            ltitle = _re.sub(r'^L\d+[_\s]*', '', os.path.splitext(name)[0], flags=_re.IGNORECASE)
            result.append({
                "path": os.path.join(notes_dir, name),
                "rel_path": f"Notes/{safe_title}/{name}",
                "lesson_number": lnum,
                "lesson_title": ltitle
            })
        return result

    @app.post("/api/data/courses/{course_id}/analyze-all")
    async def course_analyze_all(course_id: str):
        """触发课程级聚合分析（增量更新，返回最新已持久化的结果）"""
        if not _note_analyzer_ready:
            return err(msg="note_analyzer 模块未加载")

        # 复用 _read_courses_data（多源合并，同 /api/data/courses 逻辑）
        data = _read_courses_data(workspace_dir)
        courses = data.get("courses", [])
        course = None
        for c in courses:
            cid = c.get("note_id") or c.get("id") or c.get("_id") or c.get("course_title") or ""
            if str(cid) == str(course_id):
                course = c
                break
        if not course:
            return err(msg="课程未找到: " + course_id)
        course_title = course.get("course_title") or course.get("title") or ""
        note_files = _scan_course_note_paths(workspace_dir, course)
        try:
            # 读取已持久化的旧数据（用于增量对比）
            cache_path = _course_analysis_path(workspace_dir, course_id)
            old_cache = {}
            if os.path.exists(cache_path):
                try:
                    old_cache = _json_loads(Path(cache_path).read_text(encoding="utf-8"))
                except Exception:
                    old_cache = {}

            # 增量分析：仅对新增/修改过的笔记重新分析
            import time as _time
            lesson_results = old_cache.get("lessons", {}) if isinstance(old_cache, dict) else {}
            new_lessons = {}
            stats = {"total_elements": 0, "by_type": {}, "by_category": {}, "by_lesson": {}}
            analyzed_count = 0
            skipped_count = 0
            for nf in note_files:
                full_path = nf["path"]
                try:
                    mtime = os.path.getmtime(full_path)
                except Exception:
                    mtime = 0
                cache_key = nf["rel_path"]
                old_entry = lesson_results.get(cache_key)
                # 若文件未修改且已有缓存，直接复用
                if old_entry and old_entry.get("mtime") == mtime and old_entry.get("elements"):
                    new_lessons[cache_key] = old_entry
                    skipped_count += 1
                else:
                    # 重新分析
                    try:
                        elements = await _run_data(
                            NoteAnalyzer.extract_from_file,
                            _Path(full_path),
                            course_id,
                            nf["lesson_number"]
                        )
                        NoteAnalyzer.auto_classify_elements(elements)
                        NoteAnalyzer.auto_extract_triples(elements)
                        new_lessons[cache_key] = {
                            "lesson_number": nf["lesson_number"],
                            "lesson_title": nf["lesson_title"],
                            "path": nf["rel_path"],
                            "mtime": mtime,
                            "analyzed_at": _time.time(),
                            "elements": [e.to_dict() for e in elements],
                            "count": len(elements)
                        }
                        analyzed_count += 1
                    except Exception as e:
                        new_lessons[cache_key] = {
                            "lesson_number": nf["lesson_number"],
                            "lesson_title": nf["lesson_title"],
                            "path": nf["rel_path"],
                            "mtime": mtime,
                            "error": str(e),
                            "elements": [],
                            "count": 0
                        }

            # 汇总统计
            for lkey, ldata in new_lessons.items():
                cnt = ldata.get("count", 0)
                stats["total_elements"] += cnt
                stats["by_lesson"][str(ldata.get("lesson_number", "?"))] = cnt
                for e in ldata.get("elements", []):
                    t = e.get("type", "unknown")
                    stats["by_type"][t] = stats["by_type"].get(t, 0) + 1
                    for c in e.get("categories", []):
                        stats["by_category"][c] = stats["by_category"].get(c, 0) + 1

            cache_data = {
                "course_id": course_id,
                "course_title": course_title,
                "updated_at": _time.time(),
                "lessons": new_lessons,
                "stats": stats,
                "lesson_count": len(note_files),
                "analyzed_count": analyzed_count,
                "skipped_count": skipped_count
            }

            # 持久化
            try:
                os.makedirs(os.path.dirname(cache_path), exist_ok=True)
                Path(cache_path).write_text(
                    json.dumps(cache_data, ensure_ascii=False, indent=2),
                    encoding="utf-8"
                )
            except Exception as e:
                logger.warning("持久化课程分析失败: %s", e)

            return ok(data=cache_data)
        except Exception as e:
            return err(msg=f"课程聚合分析失败: {e}")

    @app.get("/api/data/courses/{course_id}/analysis")
    async def course_analysis_get(course_id: str):
        """读取已持久化的课程聚合分析数据"""
        cache_path = _course_analysis_path(workspace_dir, course_id)
        if not os.path.exists(cache_path):
            return ok(data={"cached": False, "lessons": {}, "stats": {"total_elements": 0, "by_type": {}, "by_category": {}, "by_lesson": {}}})
        try:
            cache = _json_loads(Path(cache_path).read_text(encoding="utf-8"))
            cache["cached"] = True
            return ok(data=cache)
        except Exception as e:
            return err(msg=f"读取课程分析数据失败: {e}")

    @app.get("/api/data/courses/{course_id}/analysis-data")
    async def course_analysis_data(course_id: str):
        """课程级笔记分析：从持久化缓存中提取所有课时元素并展平"""
        cache_path = _course_analysis_path(workspace_dir, course_id)
        if not os.path.exists(cache_path):
            return err(msg="请先运行课程聚合分析")
        try:
            cache = _json_loads(Path(cache_path).read_text(encoding="utf-8"))
            lessons = cache.get("lessons", {})
            all_elements = []
            by_type = {}
            by_category = {}
            lesson_titles = {}
            for lkey, ldata in lessons.items():
                ltitle = ldata.get("lesson_title") or ""
                lesson_titles[lkey] = ltitle
                for e in ldata.get("elements", []):
                    all_elements.append(e)
                    t = e.get("type", "unknown")
                    by_type[t] = by_type.get(t, 0) + 1
                    for c in e.get("categories", []):
                        by_category[c] = by_category.get(c, 0) + 1
            return ok(data={
                "elements": all_elements,
                "lesson_titles": lesson_titles,
                "stats": {"total": len(all_elements), "by_type": by_type, "by_category": by_category}
            })
        except Exception as e:
            return err(msg=f"读取课程分析数据失败: {e}")

    @app.get("/api/data/courses/{course_id}/pairs")
    async def course_pairs(course_id: str, pair_type: str = "theorem_proof"):
        """课程级定理-证明/问题-解答配对：从缓存元素中提取"""
        if not _note_analyzer_ready:
            return err(msg="note_analyzer 模块未加载")
        cache_path = _course_analysis_path(workspace_dir, course_id)
        if not os.path.exists(cache_path):
            return err(msg="请先运行课程聚合分析")
        try:
            cache = _json_loads(Path(cache_path).read_text(encoding="utf-8"))
            lessons = cache.get("lessons", {})
            all_elements = []
            lesson_map = {}
            for lkey, ldata in lessons.items():
                lnum = ldata.get("lesson_number", 0)
                ltitle = ldata.get("lesson_title") or ""
                for ed in ldata.get("elements", []):
                    elem = NoteElement.from_dict(ed)
                    all_elements.append(elem)
                    lesson_map[elem.id] = {"lesson_number": lnum, "lesson_title": ltitle}
            if pair_type == "problem_solution":
                pairs = NoteAnalyzer.find_problem_solution_pair(all_elements)
            else:
                pairs = NoteAnalyzer.find_theorem_proof_pair(all_elements)
            pair_list = []
            for primary, secondary in pairs:
                pair_list.append({
                    "primary": primary.to_dict(),
                    "secondary": secondary.to_dict() if secondary else None,
                    "primary_lesson": lesson_map.get(primary.id, {}),
                    "secondary_lesson": lesson_map.get(secondary.id, {}) if secondary else None,
                })
            return ok(data={"pairs": pair_list, "total": len(pair_list)})
        except Exception as e:
            return err(msg=f"获取课程级配对失败: {e}")

    @app.get("/api/data/courses/{course_id}/graph")
    async def course_graph(course_id: str):
        """课程级知识图谱：从缓存元素中构建"""
        if not _note_analyzer_ready:
            return err(msg="note_analyzer 模块未加载")
        cache_path = _course_analysis_path(workspace_dir, course_id)
        if not os.path.exists(cache_path):
            return err(msg="请先运行课程聚合分析")
        try:
            cache = _json_loads(Path(cache_path).read_text(encoding="utf-8"))
            lessons = cache.get("lessons", {})
            all_elements = []
            for ldata in lessons.values():
                for ed in ldata.get("elements", []):
                    all_elements.append(NoteElement.from_dict(ed))
            kg = KnowledgeGraph()
            kg.build_from_elements(all_elements)
            cyto = kg.to_cytoscape_json()
            stats = kg.get_statistics()
            return ok(data={"nodes": cyto.get("nodes", []), "edges": cyto.get("edges", []), "stats": stats})
        except Exception as e:
            return err(msg=f"课程级图谱生成失败: {e}")

    @app.get("/api/data/courses/{course_id}/spaced-load")
    async def course_spaced_load(course_id: str):
        """课程级：从聚合缓存中加载所有元素到间隔复习系统"""
        if not _note_analyzer_ready:
            return err(msg="note_analyzer 模块未加载")
        if not _srm:
            return err(msg="间隔重复模块未初始化")
        cache_path = _course_analysis_path(workspace_dir, course_id)
        if not os.path.exists(cache_path):
            return err(msg="请先运行课程聚合分析")
        try:
            cache = _json_loads(Path(cache_path).read_text(encoding="utf-8"))
            lessons = cache.get("lessons", {})
            count = 0
            for ldata in lessons.values():
                for ed in ldata.get("elements", []):
                    elem = NoteElement.from_dict(ed)
                    _srm.add_element(elem)
                    count += 1
            return ok(data={"loaded": count, "course_id": course_id})
        except Exception as e:
            return err(msg=f"课程级加载复习数据失败: {e}")

    @app.get("/api/data/spaced-repetition/all")
    async def spaced_all():
        """获取所有复习卡片"""
        if not _srm:
            return err(msg="间隔重复模块未初始化")
        from datetime import datetime as _dt
        def _check_due(next_review_str):
            if not next_review_str:
                return True
            try:
                return _dt.now() >= _dt.fromisoformat(next_review_str)
            except Exception:
                return True
        reviews = []
        levels = [0] * 7
        due_count = 0
        for elem_id, item in _srm.review_data.get("reviews", {}).items():
            is_due = _check_due(item.get("next_review"))
            if is_due:
                due_count += 1
            lv = min(item.get("level", 0), 6)
            levels[lv] += 1
            reviews.append({
                "elem_id": elem_id,
                "elem": item.get("elem", {}),
                "level": item.get("level", 0),
                "next_review": item.get("next_review"),
                "is_due": is_due,
                "history_count": len(item.get("history", []))
            })
        return ok(data={"reviews": reviews, "stats": {"total": len(reviews), "due": due_count, "levels": levels}})

    @app.post("/api/data/spaced-repetition/review")
    async def spaced_review(req: SpacedReviewRequest):
        """提交复习评分（SM-2 quality 0-5）"""
        if not _srm:
            return err(msg="间隔重复模块未初始化")
        try:
            _srm.review(req.elem_id, req.quality)
            item = _srm.review_data["reviews"].get(req.elem_id, {})
            return ok(data={"new_level": item.get("level", 0), "next_review": item.get("next_review")})
        except Exception as e:
            return err(msg=f"评分失败: {e}")

    @app.post("/api/data/spaced-repetition/load")
    async def spaced_load(req: SpacedLoadRequest):
        """从笔记文件加载复习卡片"""
        if not _note_analyzer_ready:
            return err(msg="note_analyzer 模块未加载")
        full_path = _resolve_ws_path(req.path)
        if not os.path.exists(full_path):
            return err(msg="笔记文件不存在: " + req.path)
        try:
            from pathlib import Path as _P
            elements = await _run_data(NoteAnalyzer.extract_from_file, _P(full_path), req.course_id, req.lesson_num)
            NoteAnalyzer.auto_classify_elements(elements)
            NoteAnalyzer.auto_extract_triples(elements)
            count = 0
            for e in elements:
                _srm.add_element(e)
                count += 1
            return ok(data={"loaded": count})
        except Exception as e:
            return err(msg=f"加载失败: {e}")

    @app.post("/api/data/spaced-repetition/reset")
    async def spaced_reset():
        """重置所有复习数据"""
        if not _srm:
            return err(msg="间隔重复模块未初始化")
        _srm.review_data["reviews"] = {}
        _srm._save()
        return ok(data={"reset": True})

    @app.get("/api/data/spaced-repetition/graph")
    async def spaced_graph():
        """获取知识图谱（Cytoscape.js 格式）"""
        if not _note_analyzer_ready:
            return err(msg="note_analyzer 模块未加载")
        try:
            elements = _srm.get_all_elements()
            kg = KnowledgeGraph()
            kg.build_from_elements(elements)
            cyto = kg.to_cytoscape_json()
            stats = kg.get_statistics()
            return ok(data={"nodes": cyto.get("nodes", []), "edges": cyto.get("edges", []), "stats": stats})
        except Exception as e:
            return err(msg=f"图谱生成失败: {e}")

    @app.post("/api/data/spaced-repetition/annotate")
    async def spaced_annotate(req: SpacedAnnotateRequest):
        """保存卡片批注"""
        if not _srm:
            return err(msg="间隔重复模块未初始化")
        try:
            import datetime
            reviews = _srm.review_data.setdefault("reviews", {})
            if req.elem_id not in reviews:
                return err(msg="卡片不存在")
            ann = reviews[req.elem_id].setdefault("annotations", [])
            entry = {
                "text": req.text,
                "timestamp": datetime.datetime.now().isoformat(timespec="seconds")
            }
            ann.append(entry)
            _srm._save()
            return ok(data={"annotations": ann})
        except Exception as e:
            return err(msg=f"保存批注失败: {e}")

    @app.get("/api/data/spaced-repetition/annotations")
    async def spaced_annotations(elem_id: str):
        """获取卡片批注"""
        if not _srm:
            return err(msg="间隔重复模块未初始化")
        try:
            reviews = _srm.review_data.get("reviews", {})
            item = reviews.get(elem_id, {})
            ann = item.get("annotations", [])
            return ok(data={"annotations": ann})
        except Exception as e:
            return err(msg=f"获取批注失败: {e}")

    # ─── SPA fallback（必须在所有 API 路由和 mount 之后）──────────────
    @app.get("/{path:path}", response_class=HTMLResponse)
    async def spa_fallback(path: str):
        """SPA 路由 fallback — 非 /api/ 和非静态文件的路径返回 index.html"""
        if path.startswith("api/") or path.startswith("ws") or path.startswith("static/"):
            return HTMLResponse(content="", status_code=404)
        if "." in path.split("/")[-1]:
            return HTMLResponse(content="", status_code=404)
        index_file = static_dir / "index.html"
        if index_file.exists():
            return HTMLResponse(content=index_file.read_text(encoding="utf-8"))
        return HTMLResponse(content="", status_code=404)

    # ─── Agent 并行预初始化（方案 A）───────────────────────────
    # 服务器启动时后台线程构建 Agent，与 uvicorn 启动、静态资源加载并行
    # 首次 HTTP 请求到来时 Agent 大概率已就绪，避免阻塞
    def _preinit_web_agent():
        try:
            _get_web_agent(workspace_dir)
        except Exception as _e:
            logger.warning(f"Agent 预初始化失败（不影响启动，首次请求时重试）: {_e}")

    threading.Thread(target=_preinit_web_agent, daemon=True, name="AgentPreInit").start()
    logger.info("[预初始化] Agent 后台构建已启动")

    return app


def _convert_office_to_pdf(file_path: str) -> str:
    """将 Office 文件（docx/xlsx/pptx）转换为 PDF，返回 PDF 路径"""
    import tempfile
    import fitz  # PyMuPDF

    ext = os.path.splitext(file_path)[1].lower()
    pdf_path = os.path.join(tempfile.gettempdir(), "ts2_office_preview.pdf")

    if ext == ".docx":
        from docx import Document
        doc = Document(file_path)
        pdf_doc = fitz.open()
        page = pdf_doc.new_page(width=595, height=842)  # A4
        y = 50
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                y += 12
                continue
            # 自动换行
            remaining = text
            while remaining:
                # 计算每行能放多少字符
                chars_per_line = 70
                line = remaining[:chars_per_line]
                remaining = remaining[chars_per_line:]
                if y > 800:
                    page = pdf_doc.new_page(width=595, height=842)
                    y = 50
                page.insert_text((50, y), line, fontsize=11, fontname="helv")
                y += 16
        pdf_doc.save(pdf_path)
        pdf_doc.close()

    elif ext == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        pdf_doc = fitz.open()
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            page = pdf_doc.new_page(width=842, height=595)  # 横向 A4
            page.insert_text((50, 40), f"Sheet: {sheet_name}", fontsize=14, fontname="helv")
            y = 65
            for row in ws.iter_rows(max_row=50, values_only=True):
                if y > 560:
                    page = pdf_doc.new_page(width=842, height=595)
                    y = 40
                line = "  |  ".join(str(c) if c is not None else "" for c in row)
                if line.strip(" |"):
                    page.insert_text((30, y), line[:120], fontsize=9, fontname="helv")
                    y += 14
        wb.close()
        pdf_doc.save(pdf_path)
        pdf_doc.close()

    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(file_path)
        pdf_doc = fitz.open()
        for slide in prs.slides:
            page = pdf_doc.new_page(width=960, height=720)
            y = 40
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text and y < 680:
                            page.insert_text((40, y), text[:100], fontsize=12, fontname="helv")
                            y += 18
        pdf_doc.save(pdf_path)
        pdf_doc.close()

    else:
        raise ValueError(f"不支持的文件类型: {ext}")

    return pdf_path


def _has_module(name: str) -> bool:
    """检测 Python 模块是否可用"""
    try:
        __import__(name)
        return True
    except ImportError:
        return False


def run_server(workspace_dir: Optional[str] = None, host: str = "0.0.0.0",
               port: int = 6906, open_browser: bool = True,
               auto_port: bool = True):
    """启动 TS2 服务器

    Args:
        auto_port: 如果为 True，当指定端口被占用时自动递增查找可用端口，
                   支持同一台机器运行多个 TS2 实例
    """
    import uvicorn

    # 自动端口检测：如果默认端口被占用，递增查找可用端口
    actual_port = port
    if auto_port:
        try:
            actual_port = find_available_port(port, host=host)
        except OSError:
            logger.error(f"端口 {port} 已被占用且无法找到可用端口")
            raise
        if actual_port != port:
            logger.info(f"端口 {port} 已被占用，自动切换到 {actual_port}")

    app = create_app(workspace_dir=workspace_dir, host=host, port=actual_port)
    local_ip = get_local_ip()

    if open_browser:
        import webbrowser
        import threading

        def _open():
            time.sleep(1.5)
            webbrowser.open(f"http://127.0.0.1:{actual_port}")

        threading.Thread(target=_open, daemon=True).start()

    logger.info(f"Starting TS2 Server on http://{host}:{actual_port}")
    logger.info(f"LAN access from phone: http://{local_ip}:{actual_port}")

    # macOS/Linux 自动启用 uvloop 提升性能
    try:
        import uvloop
        asyncio.set_event_loop_policy(uvloop.EventLoopPolicy())
        logger.info("已启用 uvloop (高性能事件循环)")
    except ImportError:
        if sys.platform == "darwin":
            logger.info("建议安装 uvloop: pip install uvloop (可显著提升 macOS 响应速度)")

    # 检测 httptools（比默认 h11 更快的 HTTP 解析器）
    try:
        import httptools  # noqa: F401
        logger.info("已启用 httptools (高性能 HTTP 解析)")
    except ImportError:
        if sys.platform == "darwin":
            logger.info("建议安装 httptools: pip install httptools (可提升 HTTP 解析速度)")

    # 检测 orjson（比标准 json 更快的序列化）
    try:
        import orjson  # noqa: F401
        logger.info("已启用 orjson (高性能 JSON 序列化)")
    except ImportError:
        if sys.platform == "darwin":
            logger.info("建议安装 orjson: pip install orjson (可提升 JSON 响应速度)")

    uvicorn.run(
        app,
        host=host,
        port=actual_port,
        log_level="info",
        # 连接保活
        timeout_keep_alive=120,
        # WebSocket 保活
        ws_ping_interval=30,
        ws_ping_timeout=15,
        # 高性能选项（自动检测，有则用）
        http="httptools" if _has_module("httptools") else "h11",
        loop="uvloop" if _has_module("uvloop") else "asyncio",
        # macOS 性能优化：启用 asyncio 的调试模式关闭（减少开销）
        # 并设置合理的缓冲区大小
        h11_max_incomplete_event_size=None,
    )


def run_server_in_thread(workspace_dir: Optional[str] = None,
                         host: str = "0.0.0.0",
                         port: int = 6906,
                         open_browser: bool = True,
                         auto_port: bool = True):
    """在后台线程中启动 TS2 服务器，返回 (uvicorn.Server, threading.Thread, actual_port)

    用于 PyInstaller 打包场景，避免 subprocess 调用 sys.executable 的兼容性问题。
    通过 server.should_exit = True 优雅停止。
    """
    import uvicorn
    from uvicorn.config import Config
    from uvicorn.server import Server

    actual_port = port
    if auto_port:
        try:
            actual_port = find_available_port(port, host=host)
        except OSError:
            logger.error(f"端口 {port} 已被占用且无法找到可用端口")
            raise
        if actual_port != port:
            logger.info(f"端口 {port} 已被占用，自动切换到 {actual_port}")

    app = create_app(workspace_dir=workspace_dir, host=host, port=actual_port)
    local_ip = get_local_ip()

    if open_browser:
        import webbrowser
        def _open():
            time.sleep(1.5)
            webbrowser.open(f"http://127.0.0.1:{actual_port}")
        threading.Thread(target=_open, daemon=True).start()

    logger.info(f"Starting TS2 Server on http://{host}:{actual_port}")
    logger.info(f"LAN access from phone: http://{local_ip}:{actual_port}")

    config = Config(
        app=app,
        host=host,
        port=actual_port,
        log_level="info",
        timeout_keep_alive=120,
        ws_ping_interval=30,
        ws_ping_timeout=15,
        http="httptools" if _has_module("httptools") else "h11",
        loop="uvloop" if _has_module("uvloop") else "asyncio",
    )
    server = Server(config)
    thread = threading.Thread(target=server.run, daemon=True, name="TS2-Server")
    thread.start()
    return server, thread, actual_port


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)
    run_server()



